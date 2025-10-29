"""
检查生成的 PPTX 中圆环是否正确渲染
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE

def check_rings_in_pptx(pptx_path):
    """检查 PPTX 中的圆环渲染"""
    print(f"检查文件: {pptx_path}\n")
    prs = Presentation(pptx_path)
    
    # 检查第5, 6, 7, 8, 9, 11页（索引 4, 5, 6, 7, 8, 10）
    problem_pages = [4, 5, 6, 7, 8, 10]
    
    for page_idx in problem_pages:
        if page_idx >= len(prs.slides):
            continue
            
        slide = prs.slides[page_idx]
        print(f"\n{'='*80}")
        print(f"检查第 {page_idx + 1} 页")
        print(f"{'='*80}")
        
        # 查找所有自动形状
        print(f"\n总共 {len(slide.shapes)} 个形状")
        
        # 统计所有形状类型
        shape_types = {}
        for s in slide.shapes:
            stype = "未知"
            if hasattr(s, 'shape_type'):
                stype = str(s.shape_type)
            shape_types[stype] = shape_types.get(stype, 0) + 1
        
        print("\n形状类型统计:")
        for stype, count in sorted(shape_types.items()):
            print(f"  {stype}: {count}")
        
        # 查找圆形形状（使用整数值9来检查）
        ovals = []
        for shape in slide.shapes:
            # 检查是否为 AUTO_SHAPE 且具有 auto_shape_type
            is_oval = False
            if hasattr(shape, 'shape_type'):
                if shape.shape_type == MSO_SHAPE.OVAL:
                    is_oval = True
                else:
                    try:
                        # AUTO_SHAPE 有子类型
                        if hasattr(shape, 'auto_shape_type'):
                            is_oval = (shape.auto_shape_type == MSO_SHAPE.OVAL)
                    except:
                        pass
            
            if is_oval:
                # 获取填充和描边信息
                fill_type = "无填充"
                fill_color = "N/A"
                if shape.fill.type == 1:  # SOLID
                    fill_type = "实心填充"
                    if hasattr(shape.fill.fore_color, 'rgb'):
                        rgb = shape.fill.fore_color.rgb
                        fill_color = f"RGB({rgb[0]}, {rgb[1]}, {rgb[2]})"
                
                stroke_color = "无描边"
                stroke_width = 0
                if shape.line.color.type == 1:  # SOLID
                    if hasattr(shape.line.color, 'rgb'):
                        rgb = shape.line.color.rgb
                        stroke_color = f"RGB({rgb[0]}, {rgb[1]}, {rgb[2]})"
                    stroke_width = shape.line.width.pt if shape.line.width else 0
                
                ovals.append({
                    'left': shape.left.inches,
                    'top': shape.top.inches,
                    'width': shape.width.inches,
                    'height': shape.height.inches,
                    'fill_type': fill_type,
                    'fill_color': fill_color,
                    'stroke_color': stroke_color,
                    'stroke_width': stroke_width
                })
        
        print(f"\n找到 {len(ovals)} 个椭圆/圆形:")
        for i, oval in enumerate(ovals):
            print(f"\n  圆形 #{i+1}:")
            print(f"    位置: ({oval['left']:.2f}, {oval['top']:.2f}) inches")
            print(f"    尺寸: {oval['width']:.2f} x {oval['height']:.2f} inches")
            print(f"    填充: {oval['fill_type']} - {oval['fill_color']}")
            print(f"    描边: {oval['stroke_color']} (宽度: {oval['stroke_width']:.1f}pt)")
            
            # 判断是否为圆环
            is_ring = (oval['fill_color'].startswith('RGB(255, 255, 255') and 
                      oval['stroke_width'] > 20)
            if is_ring:
                print(f"    ✅ 这是一个圆环！(白色填充 + 粗描边)")

def main():
    pptx_path = "output/glm-4.6-fixed2.pptx"
    check_rings_in_pptx(pptx_path)

if __name__ == "__main__":
    main()

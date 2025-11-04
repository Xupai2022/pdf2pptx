#!/usr/bin/env python3
"""
Comprehensive visual test to ensure the fix doesn't break other slides.
"""

from pptx import Presentation
import sys

def comprehensive_test(pptx_path):
    """
    Comprehensive test of all slides.
    """
    prs = Presentation(pptx_path)
    
    print(f"Comprehensive Visual Test: {pptx_path}")
    print(f"Total slides: {len(prs.slides)}\n")
    
    all_pass = True
    
    for idx, slide in enumerate(prs.slides, 1):
        print(f"Slide {idx}:")
        
        text_count = 0
        image_count = 0
        shape_count = 0
        
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                text_count += 1
            elif hasattr(shape, 'image'):
                image_count += 1
            else:
                shape_count += 1
        
        total = text_count + image_count + shape_count
        print(f"  Text: {text_count}, Images: {image_count}, Shapes: {shape_count}, Total: {total}")
        
        # Sanity check: each slide should have at least some elements
        if total == 0:
            print(f"  ❌ WARNING: Slide {idx} has no elements!")
            all_pass = False
        else:
            print(f"  ✅ OK")
    
    print("\n" + "=" * 80)
    if all_pass:
        print("✅ Comprehensive test PASSED: All slides have elements")
        return True
    else:
        print("❌ Comprehensive test FAILED: Some slides have issues")
        return False

if __name__ == "__main__":
    pptx_path = sys.argv[1] if len(sys.argv) > 1 else "output/安全运营月报_fixed.pptx"
    
    result = comprehensive_test(pptx_path)
    sys.exit(0 if result else 1)

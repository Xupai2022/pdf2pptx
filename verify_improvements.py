#!/usr/bin/env python3
"""
Verify all improvements made to PDF to PPTX conversion
"""

from pptx import Presentation
from pptx.util import Inches
from lxml import etree

def verify_pptx(pptx_path: str):
    """Verify PPTX output matches requirements."""
    
    print("="*80)
    print("PPTX VERIFICATION REPORT")
    print("="*80)
    
    prs = Presentation(pptx_path)
    
    # 1. Verify dimensions (1920×1080 at 144 DPI)
    print("\n1. SLIDE DIMENSIONS:")
    width_inches = prs.slide_width / 914400
    height_inches = prs.slide_height / 914400
    width_px_144 = width_inches * 144
    height_px_144 = height_inches * 144
    
    print(f"   Width: {width_inches:.3f}\" ({width_px_144:.0f}px at 144 DPI)")
    print(f"   Height: {height_inches:.3f}\" ({height_px_144:.0f}px at 144 DPI)")
    
    if abs(width_px_144 - 1920) < 1 and abs(height_px_144 - 1080) < 1:
        print("   ✅ Dimensions correct: 1920×1080")
    else:
        print(f"   ❌ Dimensions incorrect: Expected 1920×1080, got {width_px_144:.0f}×{height_px_144:.0f}")
    
    # 2. Verify transparency
    print("\n2. SEMI-TRANSPARENT BACKGROUNDS:")
    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    
    slide = prs.slides[0]
    transparent_shapes = {}
    
    for i, shape in enumerate(slide.shapes):
        if hasattr(shape, 'element'):
            try:
                spPr = shape.element.spPr
                solidFill = spPr.find('.//a:solidFill', ns)
                if solidFill is not None:
                    srgbClr = solidFill.find('.//a:srgbClr', ns)
                    if srgbClr is not None:
                        alpha = srgbClr.find('.//a:alpha', ns)
                        if alpha is not None:
                            alpha_val = int(alpha.get('val'))
                            color_val = srgbClr.get('val').upper()
                            opacity = alpha_val / 100000
                            
                            key = f"#{color_val}@{opacity:.2f}"
                            transparent_shapes[key] = transparent_shapes.get(key, 0) + 1
            except:
                pass
    
    print(f"   Total transparent shapes: {sum(transparent_shapes.values())}")
    for color_opacity, count in sorted(transparent_shapes.items()):
        color, opacity = color_opacity.split('@')
        print(f"   - {color} with opacity {opacity}: {count} shapes")
    
    # Check specific expected colors
    expected = {
        '#094174@0.08': 'Card backgrounds',
        '#DB2525@0.10': 'Red risk badges',
        '#F59D0A@0.10': 'Orange risk badges'
    }
    
    print("\n   Verification:")
    for exp_key, desc in expected.items():
        if exp_key in transparent_shapes:
            print(f"   ✅ {desc}: Found {transparent_shapes[exp_key]} instances")
        else:
            print(f"   ⚠️  {desc}: Not found")
    
    # 3. Verify element counts
    print("\n3. ELEMENT BREAKDOWN:")
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    text_boxes = 0
    pictures = 0
    auto_shapes = 0
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            text_boxes += 1
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pictures += 1
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            auto_shapes += 1
    
    total = text_boxes + pictures + auto_shapes
    print(f"   Total elements: {total}")
    print(f"   - Text boxes: {text_boxes}")
    print(f"   - Pictures: {pictures}")
    print(f"   - Auto shapes: {auto_shapes}")
    
    # Check if text boxes are independent (not over-merged)
    if text_boxes > 50:
        print(f"   ✅ Text boxes are properly independent ({text_boxes} text elements)")
    else:
        print(f"   ⚠️  Text boxes may be over-merged ({text_boxes} text elements)")
    
    # 4. Verify text wrapping
    print("\n4. TEXT WRAPPING:")
    wrapped_count = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            try:
                text_frame = shape.text_frame
                if text_frame.word_wrap:
                    wrapped_count += 1
            except:
                pass
    
    print(f"   Text boxes with word wrap: {wrapped_count}/{text_boxes}")
    if wrapped_count == 0:
        print("   ✅ Word wrap disabled (prevents forced line breaks)")
    else:
        print(f"   ⚠️  {wrapped_count} text boxes have word wrap enabled")
    
    # 5. Summary
    print("\n" + "="*80)
    print("SUMMARY:")
    print("="*80)
    
    checks_passed = 0
    total_checks = 4
    
    if abs(width_px_144 - 1920) < 1 and abs(height_px_144 - 1080) < 1:
        checks_passed += 1
        print("✅ Dimensions: 1920×1080")
    else:
        print("❌ Dimensions incorrect")
    
    if sum(transparent_shapes.values()) >= 20:
        checks_passed += 1
        print(f"✅ Transparency: {sum(transparent_shapes.values())} transparent shapes")
    else:
        print("❌ Transparency not working properly")
    
    if text_boxes > 50:
        checks_passed += 1
        print(f"✅ Text independence: {text_boxes} separate text boxes")
    else:
        print("❌ Text over-merged")
    
    if wrapped_count == 0:
        checks_passed += 1
        print("✅ Word wrap disabled")
    else:
        print("⚠️  Some text boxes have word wrap")
    
    print(f"\nTotal: {checks_passed}/{total_checks} checks passed")
    print("="*80)


if __name__ == "__main__":
    import sys
    
    pptx_file = sys.argv[1] if len(sys.argv) > 1 else "output_transparent.pptx"
    verify_pptx(pptx_file)

#!/usr/bin/env python3
"""
检查生成的PPTX第4页的元素
"""
from pptx import Presentation
import sys


def check_page4(pptx_path):
    """检查PPTX中第4页的元素"""
    prs = Presentation(pptx_path)
    
    if len(prs.slides) < 4:
        print(f"PPT只有 {len(prs.slides)} 页")
        return
    
    slide = prs.slides[3]  # 第4页，索引3
    print(f"\n{'='*80}")
    print(f"第4页 - 元素检查")
    print(f"{'='*80}\n")
    
    print(f"总共 {len(slide.shapes)} 个形状\n")
    
    # 检查三角形和线条
    triangles = []
    lines = []
    
    for i, shape in enumerate(slide.shapes):
        shape_name = shape.name
        shape_type = shape.shape_type
        
        # 检查是否是三角形
        if "triangle" in shape_name.lower() or shape_type == 5:  # MSO_SHAPE.TRIANGLE
            triangles.append({
                'index': i,
                'name': shape_name,
                'type': shape_type,
                'left': shape.left.inches,
                'top': shape.top.inches,
                'width': shape.width.inches,
                'height': shape.height.inches
            })
        
        # 检查是否是线条/连接器
        if shape_type in [1, 13]:  # LINE or CONNECTOR
            lines.append({
                'index': i,
                'name': shape_name,
                'type': shape_type,
                'left': shape.left.inches,
                'top': shape.top.inches,
                'width': shape.width.inches,
                'height': shape.height.inches
            })
        
        # 打印所有形状信息（前30个）
        if i < 30:
            print(f"形状 {i}: {shape_name} (type={shape_type})")
            print(f"  位置: ({shape.left.inches:.3f}\", {shape.top.inches:.3f}\")")
            print(f"  尺寸: {shape.width.inches:.3f}\" × {shape.height.inches:.3f}\"")
            
            # 如果是形状，打印更多信息
            if hasattr(shape, 'fill'):
                try:
                    fill_type = shape.fill.type
                    print(f"  Fill type: {fill_type}")
                except:
                    pass
            
            print()
    
    print(f"\n三角形: {len(triangles)} 个")
    for tri in triangles:
        print(f"  {tri}")
    
    print(f"\n线条: {len(lines)} 个")
    for line in lines[:5]:
        print(f"  {line}")


if __name__ == "__main__":
    pptx_file = sys.argv[1] if len(sys.argv) > 1 else "/tmp/test_all_fixes.pptx"
    check_page4(pptx_file)

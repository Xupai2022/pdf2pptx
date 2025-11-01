"""
Verify that the triangle lines in the generated PPTX have opposite directions.
This validates the fix for the bug where both sides had the same direction.
"""

from pptx import Presentation
import sys

def verify_pptx_line_directions(pptx_path):
    """
    Verify that slide 4 (page 4) has diagonal lines with opposite directions.
    
    Args:
        pptx_path: Path to PPTX file
        
    Returns:
        bool: True if verification passes, False otherwise
    """
    print("="*80)
    print("VERIFYING TRIANGLE LINE DIRECTIONS IN GENERATED PPTX")
    print("="*80)
    
    try:
        prs = Presentation(pptx_path)
        
        if len(prs.slides) < 4:
            print(f"❌ FAILED: PPTX has only {len(prs.slides)} slides, need at least 4")
            return False
        
        # Get slide 4 (index 3)
        slide = prs.slides[3]
        
        print(f"\n📄 Analyzing Slide 4 (Page 4 of PDF)")
        print(f"   Total shapes on slide: {len(slide.shapes)}")
        
        # Find connector shapes (lines)
        connectors = []
        for shape in slide.shapes:
            # Check if it's a connector (line)
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.LINE or MSO_CONNECTOR
                # Get start and end points
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                # Calculate actual start and end points
                # For connectors, left/top is start, left+width/top+height is end
                start_x = left
                start_y = top
                end_x = left + width
                end_y = top + height
                
                # Calculate slope and direction
                dx = width
                dy = height
                
                # Skip horizontal and vertical lines
                if abs(dx) > 10000 and abs(dy) > 10000:  # EMUs, roughly > 0.01 inches
                    slope = dy / dx if dx != 0 else float('inf')
                    direction = "\\" if slope > 0 else "/"
                    
                    connectors.append({
                        'start': (start_x, start_y),
                        'end': (end_x, end_y),
                        'width': width,
                        'height': height,
                        'slope': slope,
                        'direction': direction
                    })
                    
                    print(f"\n   Line found:")
                    print(f"     Position: ({start_x/914400:.1f}\", {start_y/914400:.1f}\") to ({end_x/914400:.1f}\", {end_y/914400:.1f}\")")
                    print(f"     Width x Height: {width/914400:.2f}\" x {height/914400:.2f}\"")
                    print(f"     Slope: {slope:.3f}")
                    print(f"     Direction: {direction}")
        
        print(f"\n" + "="*80)
        print(f"VERIFICATION RESULTS")
        print("="*80)
        
        if len(connectors) < 2:
            print(f"❌ FAILED: Found only {len(connectors)} diagonal line(s), expected at least 2")
            return False
        
        print(f"✓ Found {len(connectors)} diagonal lines")
        
        # Check if the first two diagonal lines have opposite directions
        dir1 = connectors[0]['direction']
        dir2 = connectors[1]['direction']
        
        print(f"   Line 1 direction: {dir1}")
        print(f"   Line 2 direction: {dir2}")
        
        if dir1 != dir2:
            print(f"\n✅ SUCCESS: The two diagonal lines have OPPOSITE directions ({dir1} and {dir2})")
            print("   This confirms the bug is fixed - triangle sides now have different slopes!")
            return True
        else:
            print(f"\n❌ FAILED: The two diagonal lines have the SAME direction ({dir1})")
            print("   Bug still present - both triangle sides have the same slope")
            return False
            
    except Exception as e:
        print(f"❌ ERROR: Failed to verify PPTX: {e}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    pptx_path = "output/season_report_fixed.pptx"
    
    if len(sys.argv) > 1:
        pptx_path = sys.argv[1]
    
    print(f"\nVerifying: {pptx_path}\n")
    
    success = verify_pptx_line_directions(pptx_path)
    
    if success:
        print("\n" + "="*80)
        print("🎉 VERIFICATION PASSED! The fix is working correctly.")
        print("="*80)
        sys.exit(0)
    else:
        print("\n" + "="*80)
        print("❌ VERIFICATION FAILED! The fix needs more work.")
        print("="*80)
        sys.exit(1)

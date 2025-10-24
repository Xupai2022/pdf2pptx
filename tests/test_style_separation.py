#!/usr/bin/env python3
"""
测试脚本：验证文本样式分离功能
用于验证PDF转PPT时，不同样式的文本（如加粗/非加粗）能够正确分离到不同的文本框
"""

import sys
import subprocess
from pathlib import Path
from pptx import Presentation

def test_style_separation():
    """测试样式分离功能"""
    
    print("=" * 100)
    print("测试：文本样式分离功能")
    print("=" * 100)
    
    # 1. 转换PDF
    print("\n第一步：转换PDF...")
    
    output_path = 'tests/output_style_test.pptx'
    result = subprocess.run(
        ['python', 'main.py', 'tests/test_sample.pdf', output_path],
        cwd=Path(__file__).parent.parent,
        capture_output=True,
        text=True
    )
    
    if result.returncode != 0:
        print(f"❌ PDF转换失败")
        print(result.stderr)
        return False
    
    print(f"✅ PDF转换完成: {output_path}")
    
    # 2. 验证结果
    print("\n第二步：验证样式分离...")
    
    prs = Presentation(output_path)
    slide = prs.slides[0]
    
    # 查找CVE相关的文本框
    cve_shapes = []
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text = shape.text
            if ("已知" in text and "CVE" in text and "利用" in text) or \
               ("存在在野" in text and "利用" in text):
                # 获取加粗属性
                bold = None
                if hasattr(shape, "text_frame") and len(shape.text_frame.paragraphs) > 0:
                    runs = shape.text_frame.paragraphs[0].runs
                    if len(runs) > 0:
                        bold = runs[0].font.bold
                cve_shapes.append({
                    'text': text,
                    'bold': bold,
                    'left': shape.left
                })
    
    # 按位置排序（从左到右）
    cve_shapes.sort(key=lambda x: x['left'])
    
    print(f"\n找到 {len(cve_shapes)} 个相关文本框:")
    for i, shape in enumerate(cve_shapes):
        bold_str = "✅ 加粗" if shape['bold'] else "⚪ 非加粗"
        print(f"  {i+1}. '{shape['text'][:30]}...' - {bold_str}")
    
    # 3. 验证预期结果
    print("\n第三步：验证预期结果...")
    
    success = True
    
    # 预期：应该有2个文本框
    if len(cve_shapes) != 2:
        print(f"❌ 文本框数量错误：预期2个，实际{len(cve_shapes)}个")
        success = False
    else:
        print(f"✅ 文本框数量正确：2个")
    
    # 预期：第一个应该是加粗的"已知CVE利用"
    if len(cve_shapes) >= 1:
        if "已知" in cve_shapes[0]['text'] and "CVE" in cve_shapes[0]['text'] and cve_shapes[0]['bold']:
            print("✅ 文本框1正确：包含'已知CVE利用' - 加粗")
        else:
            print(f"❌ 文本框1错误：'{cve_shapes[0]['text']}' - 加粗: {cve_shapes[0]['bold']}")
            success = False
    
    # 预期：第二个应该是非加粗的"：存在在野利用报"
    if len(cve_shapes) >= 2:
        if "存在" in cve_shapes[1]['text'] and "利用" in cve_shapes[1]['text'] and not cve_shapes[1]['bold']:
            print("✅ 文本框2正确：包含'存在在野利用' - 非加粗")
        else:
            print(f"❌ 文本框2错误：'{cve_shapes[1]['text']}' - 加粗: {cve_shapes[1]['bold']}")
            success = False
    
    # 4. 输出结果
    print("\n" + "=" * 100)
    if success:
        print("✅✅✅ 测试通过！样式分离功能正常工作！")
        print("\n功能说明：")
        print("  - PDF中同一行相邻的文本元素，如果样式不同（加粗/非加粗），会被分离到不同文本框")
        print("  - 确保样式保真度，正确还原PDF的视觉效果")
    else:
        print("❌ 测试失败！请检查代码")
    print("=" * 100)
    
    return success

if __name__ == "__main__":
    success = test_style_separation()
    sys.exit(0 if success else 1)

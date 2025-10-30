"""
综合测试：验证所有修复的问题
"""

from pptx import Presentation
import sys

def test_pptx(pptx_path):
    """测试PPT文件中的关键文本"""
    prs = Presentation(pptx_path)
    
    # 第5页（索引4）
    slide = prs.slides[4]
    
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                left_pt = shape.left / 914400 * 72
                top_pt = shape.top / 914400 * 72
                width_pt = shape.width / 914400 * 72
                right_pt = left_pt + width_pt
                
                text_shapes.append({
                    'text': text,
                    'left': left_pt,
                    'top': top_pt,
                    'right': right_pt,
                })
    
    # 排序
    text_shapes.sort(key=lambda s: (s['top'], s['left']))
    
    # 测试用例
    tests = []
    
    # 1. 检查"44"是否与"起服务外"分开
    forty_four = None
    qi_service = None
    for shape in text_shapes:
        if shape['text'] == '44':
            forty_four = shape
        if '起服务外资产的事件处置' in shape['text']:
            qi_service = shape
    
    if forty_four and qi_service:
        gap = qi_service['left'] - forty_four['right']
        tests.append({
            'name': '"44" 与 "起服务外资产的事件处置"',
            'pass': gap >= 2.0,  # 至少2pt间距
            'detail': f'间距: {gap:.2f}pt',
            'expected': '间距 >= 2pt',
        })
    else:
        tests.append({
            'name': '"44" 与 "起服务外资产的事件处置"',
            'pass': False,
            'detail': '未找到文本或已被合并',
            'expected': '分开的两个文本框',
        })
    
    # 2. 检查"135"是否与"起"分开
    one_three_five = None
    qi = None
    for shape in text_shapes:
        if shape['text'] == '135':
            one_three_five = shape
        if shape['text'] == '起。':
            qi = shape
    
    if one_three_five and qi:
        gap = qi['left'] - one_three_five['right']
        tests.append({
            'name': '"135" 与 "起。"',
            'pass': gap >= 2.0,
            'detail': f'间距: {gap:.2f}pt',
            'expected': '间距 >= 2pt',
        })
    else:
        tests.append({
            'name': '"135" 与 "起。"',
            'pass': False,
            'detail': '未找到文本或已被合并',
            'expected': '分开的两个文本框',
        })
    
    # 3. 检查"62"是否与"个）"分开
    sixty_two = None
    ge_kuohao = None
    for shape in text_shapes:
        if shape['text'] == '62':
            sixty_two = shape
        if shape['text'].startswith('个）'):
            ge_kuohao = shape
    
    if sixty_two and ge_kuohao:
        gap = ge_kuohao['left'] - sixty_two['right']
        tests.append({
            'name': '"62" 与 "个）..."',
            'pass': gap >= 2.0,
            'detail': f'间距: {gap:.2f}pt',
            'expected': '间距 >= 2pt',
        })
    else:
        tests.append({
            'name': '"62" 与 "个）..."',
            'pass': False,
            'detail': '未找到文本或已被合并',
            'expected': '分开的两个文本框',
        })
    
    # 4. 检查"个）"之间没有多余空格
    if ge_kuohao:
        has_space = '个 ）' in ge_kuohao['text']
        tests.append({
            'name': '"个）"之间无空格',
            'pass': not has_space,
            'detail': f'文本: {ge_kuohao["text"][:20]}...',
            'expected': '无空格',
        })
    else:
        tests.append({
            'name': '"个）"之间无空格',
            'pass': False,
            'detail': '未找到文本',
            'expected': '无空格',
        })
    
    # 5. 检查"事件总量"后面的文本（检查是否有多余空格）
    shijian_zongliang = None
    for shape in text_shapes:
        if '事件总量 ' in shape['text']:  # 注意空格
            shijian_zongliang = shape
            break
    
    tests.append({
        'name': '"事件总量"后无多余空格',
        'pass': shijian_zongliang is None,  # 应该没有这个文本（有空格的版本）
        'detail': '检查通过' if shijian_zongliang is None else f'发现多余空格: {shijian_zongliang["text"][:30]}',
        'expected': '无多余空格',
    })
    
    # 打印结果
    print("="*80)
    print(f"测试文件: {pptx_path}")
    print("="*80)
    
    all_passed = True
    for test in tests:
        status = "✓ PASS" if test['pass'] else "✗ FAIL"
        print(f"\n{status}: {test['name']}")
        print(f"  详情: {test['detail']}")
        print(f"  期望: {test['expected']}")
        if not test['pass']:
            all_passed = False
    
    print("\n" + "="*80)
    if all_passed:
        print("✓ 所有测试通过！")
    else:
        print("✗ 部分测试失败")
    print("="*80)
    
    return all_passed

if __name__ == '__main__':
    pptx_path = 'output/test_fixed3.pptx'
    success = test_pptx(pptx_path)
    sys.exit(0 if success else 1)

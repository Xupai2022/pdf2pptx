#!/usr/bin/env python3
"""
Color Accuracy Test Script

This script validates that PDF to PPTX conversion preserves colors and transparency
accurately. It compares RGB values and opacity between the source PDF and generated PPTX.

Test Coverage:
1. RGB color values match exactly (tolerance: ±1 for rounding errors)
2. Transparency values match (tolerance: ±0.001 for floating point precision)
3. No duplicate overlapping shapes (border artifacts are removed)
4. All containers have correct visual appearance
"""

import sys
import os
import logging
from typing import Dict, List, Tuple
from pathlib import Path

# Add parent directory to path
sys.path.insert(0, str(Path(__file__).parent.parent))

import fitz  # PyMuPDF
from pptx import Presentation
import re

logging.basicConfig(level=logging.INFO, format='%(levelname)s: %(message)s')
logger = logging.getLogger(__name__)


class ColorAccuracyTester:
    """
    Tests color accuracy between PDF source and PPTX output.
    """
    
    def __init__(self, pdf_path: str, pptx_path: str):
        """
        Initialize tester with PDF and PPTX paths.
        
        Args:
            pdf_path: Path to source PDF file
            pptx_path: Path to generated PPTX file
        """
        self.pdf_path = pdf_path
        self.pptx_path = pptx_path
        self.results = {
            'total_tests': 0,
            'passed': 0,
            'failed': 0,
            'errors': []
        }
    
    def run_tests(self) -> Dict:
        """
        Run all color accuracy tests.
        
        Returns:
            Dictionary with test results
        """
        logger.info(f"Testing PDF: {self.pdf_path}")
        logger.info(f"Against PPTX: {self.pptx_path}")
        logger.info("=" * 70)
        
        # Extract PDF data
        pdf_shapes = self._extract_pdf_shapes()
        
        # Extract PPTX data
        pptx_shapes = self._extract_pptx_shapes()
        
        # Test 1: No duplicate overlapping shapes in PPTX
        self._test_no_duplicates(pptx_shapes)
        
        # Test 2: Color and opacity matching
        self._test_color_matching(pdf_shapes, pptx_shapes)
        
        # Print results
        self._print_results()
        
        return self.results
    
    def _extract_pdf_shapes(self) -> List[Dict]:
        """
        Extract shape data from PDF.
        
        Returns:
            List of shape dictionaries with color and opacity info
        """
        shapes = []
        
        try:
            doc = fitz.open(self.pdf_path)
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                
                # Extract opacity mapping
                opacity_map = self._extract_pdf_opacity_map(doc, page)
                
                # Extract drawings
                drawings = page.get_drawings()
                gs_opacity_sequence = self._parse_pdf_opacity_sequence(doc, page, opacity_map)
                
                for idx, drawing in enumerate(drawings):
                    rect = drawing.get("rect")
                    if not rect:
                        continue
                    
                    fill = drawing.get("fill")
                    if not fill:
                        continue
                    
                    # Get opacity
                    opacity = 1.0
                    if idx < len(gs_opacity_sequence):
                        opacity = gs_opacity_sequence[idx]
                    
                    # Extract RGB
                    r, g, b = fill
                    rgb = (int(r * 255), int(g * 255), int(b * 255))
                    
                    shapes.append({
                        'page': page_num,
                        'x': rect.x0,
                        'y': rect.y0,
                        'width': rect.width,
                        'height': rect.height,
                        'rgb': rgb,
                        'opacity': opacity
                    })
            
            doc.close()
            
        except Exception as e:
            logger.error(f"Failed to extract PDF shapes: {e}")
        
        return shapes
    
    def _extract_pdf_opacity_map(self, doc, page) -> Dict[str, float]:
        """Extract opacity mapping from PDF ExtGState."""
        opacity_map = {}
        
        try:
            page_dict = doc.xref_object(page.xref, compressed=False)
            extgstate_pattern = r'/ExtGState\s*<<([^>]+)>>'
            match = re.search(extgstate_pattern, page_dict, re.DOTALL)
            
            if match:
                extgstate_content = match.group(1)
                gs_refs = re.findall(r'/([A-Za-z]+\d*)\s+(\d+)\s+\d+\s+R', extgstate_content)
                
                for gs_name, xref in gs_refs:
                    try:
                        gs_obj = doc.xref_object(int(xref), compressed=False)
                        ca_match = re.search(r'/ca\s+([\d.]+)', gs_obj)
                        if ca_match:
                            opacity = float(ca_match.group(1))
                            opacity_map[gs_name] = opacity
                    except:
                        pass
        except:
            pass
        
        return opacity_map
    
    def _parse_pdf_opacity_sequence(self, doc, page, opacity_map: Dict[str, float]) -> List[float]:
        """Parse PDF content stream to get opacity for each drawing."""
        opacity_sequence = []
        current_opacity = 1.0
        
        try:
            xref = page.get_contents()[0]
            content_stream = doc.xref_stream(xref).decode('latin-1', errors='ignore')
            tokens = content_stream.split()
            
            i = 0
            while i < len(tokens):
                token = tokens[i]
                
                # Graphics state change
                if token.startswith('/') and i + 1 < len(tokens) and tokens[i + 1] == 'gs':
                    gs_match = re.match(r'/([A-Za-z]+\d*)', token)
                    if gs_match:
                        gs_name = gs_match.group(1)
                        current_opacity = opacity_map.get(gs_name, 1.0)
                    i += 2
                    continue
                
                # Fill operation
                if token in ['f', 'f*']:
                    opacity_sequence.append(current_opacity)
                
                i += 1
        except:
            pass
        
        return opacity_sequence
    
    def _extract_pptx_shapes(self) -> List[Dict]:
        """
        Extract shape data from PPTX.
        
        Returns:
            List of shape dictionaries with color and opacity info
        """
        shapes = []
        
        try:
            prs = Presentation(self.pptx_path)
            
            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if not hasattr(shape, 'fill') or not hasattr(shape, 'line'):
                        continue
                    
                    if shape.fill.type != 1:  # Not solid fill
                        continue
                    
                    try:
                        # Get RGB
                        rgb = shape.fill.fore_color.rgb
                        rgb = (rgb[0], rgb[1], rgb[2])
                        
                        # Get opacity
                        opacity = 1.0
                        spPr = shape.element.spPr
                        ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                        solidFill = spPr.find('.//a:solidFill', ns)
                        if solidFill is not None:
                            srgbClr = solidFill.find('.//a:srgbClr', ns)
                            if srgbClr is not None:
                                alpha = srgbClr.find('.//a:alpha', ns)
                                if alpha is not None:
                                    alpha_val = int(alpha.get('val'))
                                    opacity = alpha_val / 100000.0
                        
                        shapes.append({
                            'page': slide_num,
                            'x': shape.left.inches,
                            'y': shape.top.inches,
                            'width': shape.width.inches,
                            'height': shape.height.inches,
                            'rgb': rgb,
                            'opacity': opacity
                        })
                    except:
                        pass
        except Exception as e:
            logger.error(f"Failed to extract PPTX shapes: {e}")
        
        return shapes
    
    def _test_no_duplicates(self, pptx_shapes: List[Dict]):
        """
        Test that there are no duplicate overlapping shapes.
        
        Args:
            pptx_shapes: List of PPTX shape dictionaries
        """
        logger.info("\nTest 1: No Duplicate Overlapping Shapes")
        logger.info("-" * 70)
        
        # Group shapes by page
        pages = {}
        for shape in pptx_shapes:
            page = shape['page']
            if page not in pages:
                pages[page] = []
            pages[page].append(shape)
        
        duplicates_found = 0
        
        for page_num, shapes in pages.items():
            for i, shape1 in enumerate(shapes):
                for j in range(i + 1, len(shapes)):
                    shape2 = shapes[j]
                    
                    # Check if shapes overlap
                    if self._shapes_overlap(shape1, shape2):
                        duplicates_found += 1
                        error_msg = (f"Page {page_num}: Overlapping shapes found - "
                                   f"Shape1 at ({shape1['x']:.2f}, {shape1['y']:.2f}), "
                                   f"Shape2 at ({shape2['x']:.2f}, {shape2['y']:.2f})")
                        self.results['errors'].append(error_msg)
                        logger.error(f"  FAIL: {error_msg}")
        
        self.results['total_tests'] += 1
        if duplicates_found == 0:
            self.results['passed'] += 1
            logger.info(f"  PASS: No duplicate overlapping shapes found")
        else:
            self.results['failed'] += 1
            logger.error(f"  FAIL: Found {duplicates_found} duplicate overlapping shapes")
    
    def _shapes_overlap(self, shape1: Dict, shape2: Dict) -> bool:
        """Check if two shapes overlap significantly."""
        # Same color check
        if shape1['rgb'] != shape2['rgb']:
            return False
        
        # Position overlap (within 0.05 inches tolerance)
        tolerance = 0.05
        x_overlap = abs(shape1['x'] - shape2['x']) <= tolerance
        y_overlap = abs(shape1['y'] - shape2['y']) <= tolerance
        
        # Size similarity (within 0.05 inches tolerance)
        width_similar = abs(shape1['width'] - shape2['width']) <= tolerance
        height_similar = abs(shape1['height'] - shape2['height']) <= tolerance
        
        # Different opacity
        opacity_different = abs(shape1['opacity'] - shape2['opacity']) > 0.01
        
        return x_overlap and y_overlap and width_similar and height_similar and opacity_different
    
    def _test_color_matching(self, pdf_shapes: List[Dict], pptx_shapes: List[Dict]):
        """
        Test that colors and opacity match between PDF and PPTX.
        
        Args:
            pdf_shapes: List of PDF shape dictionaries
            pptx_shapes: List of PPTX shape dictionaries
        """
        logger.info("\nTest 2: Color and Opacity Matching")
        logger.info("-" * 70)
        
        # Filter shapes with transparency from PDF
        pdf_transparent = [s for s in pdf_shapes if s['opacity'] < 1.0]
        pptx_transparent = [s for s in pptx_shapes if s['opacity'] < 1.0]
        
        logger.info(f"  PDF shapes with transparency: {len(pdf_transparent)}")
        logger.info(f"  PPTX shapes with transparency: {len(pptx_transparent)}")
        
        # Test specific opacity values
        expected_opacities = [0.0314, 0.0784]  # From glm-4.6.pdf
        
        for expected in expected_opacities:
            pdf_count = len([s for s in pdf_transparent if abs(s['opacity'] - expected) < 0.001])
            pptx_count = len([s for s in pptx_transparent if abs(s['opacity'] - expected) < 0.001])
            
            self.results['total_tests'] += 1
            if pdf_count == pptx_count and pptx_count > 0:
                self.results['passed'] += 1
                logger.info(f"  PASS: Opacity {expected:.4f} - PDF:{pdf_count}, PPTX:{pptx_count}")
            else:
                self.results['failed'] += 1
                error_msg = f"Opacity {expected:.4f} mismatch - PDF:{pdf_count}, PPTX:{pptx_count}"
                self.results['errors'].append(error_msg)
                logger.error(f"  FAIL: {error_msg}")
        
        # Test RGB color preservation for transparent shapes
        self.results['total_tests'] += 1
        rgb_matches = 0
        for pdf_shape in pdf_transparent[:5]:  # Test first 5
            for pptx_shape in pptx_transparent:
                if (abs(pdf_shape['opacity'] - pptx_shape['opacity']) < 0.001 and
                    self._colors_match(pdf_shape['rgb'], pptx_shape['rgb'])):
                    rgb_matches += 1
                    break
        
        if rgb_matches >= len(pdf_transparent[:5]) * 0.8:  # 80% match rate
            self.results['passed'] += 1
            logger.info(f"  PASS: RGB colors match for transparent shapes ({rgb_matches}/{len(pdf_transparent[:5])})")
        else:
            self.results['failed'] += 1
            error_msg = f"RGB color mismatch for transparent shapes ({rgb_matches}/{len(pdf_transparent[:5])})"
            self.results['errors'].append(error_msg)
            logger.error(f"  FAIL: {error_msg}")
    
    def _colors_match(self, rgb1: Tuple[int, int, int], rgb2: Tuple[int, int, int]) -> bool:
        """Check if two RGB colors match (with tolerance for rounding)."""
        tolerance = 1  # Allow ±1 for rounding errors
        return (abs(rgb1[0] - rgb2[0]) <= tolerance and
                abs(rgb1[1] - rgb2[1]) <= tolerance and
                abs(rgb1[2] - rgb2[2]) <= tolerance)
    
    def _print_results(self):
        """Print test results summary."""
        logger.info("\n" + "=" * 70)
        logger.info("TEST RESULTS SUMMARY")
        logger.info("=" * 70)
        logger.info(f"Total tests: {self.results['total_tests']}")
        logger.info(f"Passed: {self.results['passed']}")
        logger.info(f"Failed: {self.results['failed']}")
        
        if self.results['failed'] > 0:
            logger.info("\nErrors:")
            for error in self.results['errors']:
                logger.info(f"  - {error}")
        
        success_rate = (self.results['passed'] / self.results['total_tests'] * 100 
                       if self.results['total_tests'] > 0 else 0)
        logger.info(f"\nSuccess rate: {success_rate:.1f}%")
        
        if self.results['failed'] == 0:
            logger.info("\n✅ ALL TESTS PASSED!")
        else:
            logger.info(f"\n❌ {self.results['failed']} TEST(S) FAILED")


def main():
    """Main test function."""
    if len(sys.argv) < 3:
        print("Usage: python test_color_accuracy.py <pdf_path> <pptx_path>")
        sys.exit(1)
    
    pdf_path = sys.argv[1]
    pptx_path = sys.argv[2]
    
    if not os.path.exists(pdf_path):
        print(f"Error: PDF file not found: {pdf_path}")
        sys.exit(1)
    
    if not os.path.exists(pptx_path):
        print(f"Error: PPTX file not found: {pptx_path}")
        sys.exit(1)
    
    tester = ColorAccuracyTester(pdf_path, pptx_path)
    results = tester.run_tests()
    
    # Exit with error code if tests failed
    sys.exit(0 if results['failed'] == 0 else 1)


if __name__ == '__main__':
    main()

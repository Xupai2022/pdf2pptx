"""
Comprehensive test suite for PDF to PPTX border conversion.

Tests the enhanced border extraction functionality on multiple PDF files.
"""

import unittest
import sys
import os
from pathlib import Path
from pptx import Presentation

# Add parent directory to path
sys.path.insert(0, str(Path(__file__).parent.parent))

from main import convert_pdf_to_pptx, load_config


class TestBorderConversion(unittest.TestCase):
    """Test border extraction and conversion for various PDF files."""
    
    @classmethod
    def setUpClass(cls):
        """Setup test environment."""
        cls.config = load_config()
        cls.output_dir = Path(__file__).parent.parent / 'output'
        cls.output_dir.mkdir(exist_ok=True)
        cls.tests_dir = Path(__file__).parent
    
    def test_glm46_page3_borders(self):
        """Test GLM-4.6 page 3 left border extraction."""
        input_pdf = self.tests_dir / 'glm-4.6.pdf'
        output_pptx = self.output_dir / 'test_glm46_borders.pptx'
        
        # Convert PDF
        success = convert_pdf_to_pptx(
            str(input_pdf),
            str(output_pptx),
            self.config
        )
        
        self.assertTrue(success, "Conversion should succeed")
        self.assertTrue(output_pptx.exists(), "Output file should be created")
        
        # Check slide 3 (index 2) for borders
        prs = Presentation(str(output_pptx))
        self.assertGreater(len(prs.slides), 2, "Should have at least 3 slides")
        
        slide = prs.slides[2]  # Page 3
        
        # Count left borders (thin vertical shapes on the left side)
        left_borders = self._count_left_borders(slide)
        
        self.assertGreater(left_borders, 0, 
                          f"Page 3 should have left borders, found {left_borders}")
        
        print(f"✓ GLM-4.6 Page 3: {left_borders} left border(s) correctly extracted")
    
    def test_test_sample_borders(self):
        """Test test_sample.pdf border preservation."""
        input_pdf = self.tests_dir / 'test_sample.pdf'
        output_pptx = self.output_dir / 'test_sample_borders.pptx'
        
        # Convert PDF
        success = convert_pdf_to_pptx(
            str(input_pdf),
            str(output_pptx),
            self.config
        )
        
        self.assertTrue(success, "Conversion should succeed")
        self.assertTrue(output_pptx.exists(), "Output file should be created")
        
        # Check for borders
        prs = Presentation(str(output_pptx))
        slide = prs.slides[0]
        
        left_borders = self._count_left_borders(slide)
        
        self.assertGreater(left_borders, 0, 
                          f"test_sample should have left borders, found {left_borders}")
        
        print(f"✓ test_sample.pdf: {left_borders} left border(s) preserved")
    
    def test_border_color_accuracy(self):
        """Test that border colors are accurately preserved."""
        input_pdf = self.tests_dir / 'glm-4.6.pdf'
        output_pptx = self.output_dir / 'test_glm46_colors.pptx'
        
        # Convert PDF
        convert_pdf_to_pptx(str(input_pdf), str(output_pptx), self.config)
        
        # Check slide 3 border colors
        prs = Presentation(str(output_pptx))
        slide = prs.slides[2]
        
        # Expected blue color: rgb(10, 66, 117) or close
        expected_rgb = (10, 66, 117)
        tolerance = 10
        
        blue_borders = 0
        for shape in slide.shapes:
            if self._is_left_border(shape):
                try:
                    if hasattr(shape.fill, 'fore_color') and hasattr(shape.fill.fore_color, 'rgb'):
                        rgb = shape.fill.fore_color.rgb
                        r, g, b = rgb[0], rgb[1], rgb[2]
                        
                        # Check if color is close to expected
                        if (abs(r - expected_rgb[0]) <= tolerance and
                            abs(g - expected_rgb[1]) <= tolerance and
                            abs(b - expected_rgb[2]) <= tolerance):
                            blue_borders += 1
                except:
                    pass
        
        self.assertGreater(blue_borders, 0, 
                          "Should have borders with the correct blue color")
        
        print(f"✓ Color accuracy: {blue_borders} border(s) with correct blue color")
    
    def test_no_false_positives(self):
        """Test that page-edge borders are not included."""
        input_pdf = self.tests_dir / 'glm-4.6.pdf'
        output_pptx = self.output_dir / 'test_glm46_clean.pptx'
        
        # Convert PDF
        convert_pdf_to_pptx(str(input_pdf), str(output_pptx), self.config)
        
        # Check slide 3
        prs = Presentation(str(output_pptx))
        slide = prs.slides[2]
        
        # Count borders at page edges (should be 0 or minimal)
        edge_borders = 0
        page_width = prs.slide_width / 914400  # Convert to inches
        
        for shape in slide.shapes:
            if self._is_vertical_border(shape):
                x_inches = shape.left / 914400
                
                # Check if at page edges (< 0.2" from left or right)
                if x_inches < 0.2 or x_inches > page_width - 0.2:
                    edge_borders += 1
        
        # Should have very few or no edge borders
        self.assertLess(edge_borders, 2, 
                       f"Should not extract page-edge borders, found {edge_borders}")
        
        print(f"✓ No false positives: Only {edge_borders} edge border(s)")
    
    def _is_vertical_border(self, shape) -> bool:
        """Check if a shape is a vertical border."""
        try:
            if shape.shape_type != 1:  # Not a rectangle
                return False
            
            width_inches = shape.width / 914400
            height_inches = shape.height / 914400
            
            # Vertical: height >> width, width < 0.15"
            return height_inches > width_inches * 3 and width_inches < 0.15
        except:
            return False
    
    def _is_left_border(self, shape) -> bool:
        """Check if a shape is a left border."""
        if not self._is_vertical_border(shape):
            return False
        
        # Check position (left side of page)
        x_inches = shape.left / 914400
        return x_inches < 2.0  # Within 2 inches of left edge
    
    def _count_left_borders(self, slide) -> int:
        """Count left borders in a slide."""
        count = 0
        for shape in slide.shapes:
            if self._is_left_border(shape):
                count += 1
        return count


def run_tests():
    """Run all tests with verbose output."""
    suite = unittest.TestLoader().loadTestsFromTestCase(TestBorderConversion)
    runner = unittest.TextTestRunner(verbosity=2)
    result = runner.run(suite)
    
    print("\n" + "=" * 80)
    print("TEST SUMMARY")
    print("=" * 80)
    print(f"Tests run: {result.testsRun}")
    print(f"Successes: {result.testsRun - len(result.failures) - len(result.errors)}")
    print(f"Failures: {len(result.failures)}")
    print(f"Errors: {len(result.errors)}")
    
    if result.wasSuccessful():
        print("\n✅ ALL TESTS PASSED!")
        return 0
    else:
        print("\n❌ SOME TESTS FAILED")
        return 1


if __name__ == '__main__':
    sys.exit(run_tests())

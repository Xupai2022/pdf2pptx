#!/usr/bin/env python3
"""
全面验证修复效果
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
import os

def validate_season_report():
    """验证season_report_del.pptx第15页"""
    print("="*80)
    print("验证 season_report_del.pptx 第15页")
    print("="*80)
    
    pptx_path = "/home/user/webapp/output/season_report_del.pptx"
    if not os.path.exists(pptx_path):
        print("❌ 文件不存在")
        return False
    
    prs = Presentation(pptx_path)
    slide = prs.slides[14]  # 第15页
    
    # 查找相关文本框
    ip_boxes = []
    bracket_boxes = []
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text
            if '10.74.145.44' in text:
                ip_boxes.append((text, shape.top, shape.left))
            if '未知业务' in text and '(' in text and ')' in text:
                bracket_boxes.append((text, shape.top, shape.left))
    
    print(f"\n找到 {len(ip_boxes)} 个包含IP地址的文本框")
    print(f"找到 {len(bracket_boxes)} 个完整的括号组文本框")
    
    # 验证1：IP地址应该单独显示
    success = True
    for text, top, left in ip_boxes:
        if '未知业务' in text or '(' in text:
            print(f"❌ IP地址与其他内容合并: {repr(text)}")
            success = False
        else:
            print(f"✅ IP地址单独显示: {repr(text)}")
    
    # 验证2：括号组应该完整
    for text, top, left in bracket_boxes:
        if text.startswith('(') and text.endswith(')'):
            print(f"✅ 括号组完整: {repr(text)}")
        else:
            print(f"⚠️  括号组不完整: {repr(text)}")
    
    # 验证3：检查是否有单独的括号
    lone_brackets = 0
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text in ['(', ')', '（', '）']:
                lone_brackets += 1
                print(f"⚠️  发现单独的括号: {repr(text)} at ({shape.left}, {shape.top})")
    
    if lone_brackets == 0:
        print("✅ 没有单独的括号文本框")
    
    return success

def validate_glm_shapes():
    """验证glm-4.6.pptx的所有矩形元素"""
    print("\n" + "="*80)
    print("验证 glm-4.6.pptx 的矩形元素")
    print("="*80)
    
    pptx_path = "/home/user/webapp/output/glm-4.6.pptx"
    if not os.path.exists(pptx_path):
        print("❌ 文件不存在")
        return False
    
    prs = Presentation(pptx_path)
    success = True
    
    # 检查所有页面
    for page_num, slide in enumerate(prs.slides, 1):
        rectangles = 0
        ovals = 0
        others = 0
        
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'auto_shape_type'):
                    shape_type = shape.auto_shape_type
                    if shape_type == MSO_AUTO_SHAPE_TYPE.RECTANGLE:
                        rectangles += 1
                    elif shape_type == MSO_AUTO_SHAPE_TYPE.OVAL:
                        ovals += 1
                    else:
                        others += 1
            except ValueError:
                # Not an auto shape
                pass
        
        if page_num == 5:
            # 第5页应该有矩形和1个椭圆
            print(f"\n第{page_num}页: 矩形={rectangles}, 椭圆={ovals}, 其他={others}")
            if ovals > 1:
                print(f"  ⚠️  预期只有1个真正的圆形，但发现{ovals}个")
                success = False
            elif ovals == 1:
                print(f"  ✅ 正确：1个真正的圆形")
            
            if rectangles < 6:
                print(f"  ⚠️  预期至少6个矩形，但只有{rectangles}个")
                success = False
            else:
                print(f"  ✅ 正确：{rectangles}个矩形")
        elif rectangles > 0 or ovals > 0:
            print(f"第{page_num}页: 矩形={rectangles}, 椭圆={ovals}")
    
    return success

def main():
    """主验证函数"""
    print("\n" + "#"*80)
    print("开始全面验证")
    print("#"*80 + "\n")
    
    result1 = validate_season_report()
    result2 = validate_glm_shapes()
    
    print("\n" + "="*80)
    print("验证结果汇总")
    print("="*80)
    print(f"season_report_del.pptx: {'✅ 通过' if result1 else '❌ 失败'}")
    print(f"glm-4.6.pptx: {'✅ 通过' if result2 else '❌ 失败'}")
    
    if result1 and result2:
        print("\n🎉 所有验证通过！")
        return 0
    else:
        print("\n⚠️  部分验证失败，需要继续调整")
        return 1

if __name__ == "__main__":
    exit(main())

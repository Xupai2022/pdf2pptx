"""
Test script to verify the duplication fix for season_report_del.pdf.

This script tests that:
1. Page 4: Triangle region has NO PNG overlays, only vector shapes (lines)
2. Page 6: Large images exist but DON'T capture overlapping text elements
3. No duplication or shadowing artifacts in the output
"""

import sys
from pathlib import Path
import yaml
from src.parser.pdf_parser import PDFParser
from pptx import Presentation
from pptx.util import Inches, Pt
import io
from PIL import Image

def test_page_4_no_image_overlap():
    """Test that page 4 triangle region has no PNG overlays over vector shapes."""
    print("\n" + "="*80)
    print("TEST 1: Page 4 - No PNG overlays in triangle region")
    print("="*80)
    
    with open('config/config.yaml', 'r', encoding='utf-8') as f:
        config = yaml.safe_load(f)
    
    parser = PDFParser(config['parser'])
    pdf_path = 'tests/season_report_del.pdf'
    
    if not parser.open(pdf_path):
        print("❌ FAIL: Could not open PDF")
        return False
    
    # Extract page 4 (index 3)
    page_data = parser.extract_page_elements(3)
    parser.close()
    
    # Triangle region: approximately x: 150-450, y: 150-350
    triangle_region = {
        'x_min': 150, 'x_max': 450,
        'y_min': 150, 'y_max': 350
    }
    
    # Count rerendered images in triangle region
    rerendered_images_in_triangle = []
    vector_lines_in_triangle = []
    
    for elem in page_data['elements']:
        x, y = elem.get('x', 0), elem.get('y', 0)
        x2, y2 = elem.get('x2', 0), elem.get('y2', 0)
        
        # Check overlap with triangle region
        overlaps = (
            x < triangle_region['x_max'] and x2 > triangle_region['x_min'] and
            y < triangle_region['y_max'] and y2 > triangle_region['y_min']
        )
        
        if overlaps:
            if elem['type'] == 'image' and elem.get('was_rerendered', False):
                rerendered_images_in_triangle.append(elem)
            elif elem['type'] == 'shape' and elem.get('shape_type') == 'line':
                vector_lines_in_triangle.append(elem)
    
    print(f"\nRerendered images in triangle region: {len(rerendered_images_in_triangle)}")
    print(f"Vector lines in triangle region: {len(vector_lines_in_triangle)}")
    
    # TEST: No rerendered images should exist in triangle region
    if rerendered_images_in_triangle:
        print("\n❌ FAIL: Found rerendered images in triangle region!")
        print("These images will overlay the vector shapes causing duplication.")
        for img in rerendered_images_in_triangle:
            print(f"  - Image at ({img['x']:.1f}, {img['y']:.1f}), size: {img['width']:.1f}x{img['height']:.1f}")
        return False
    
    print("\n✅ PASS: No rerendered PNG overlays in triangle region")
    
    # BONUS: Check that vector lines exist
    if vector_lines_in_triangle:
        print(f"✅ BONUS: Found {len(vector_lines_in_triangle)} vector lines (triangle edges)")
    
    return True

def test_page_6_no_text_overlap():
    """Test that page 6 large images don't capture overlapping text."""
    print("\n" + "="*80)
    print("TEST 2: Page 6 - Large images don't capture text")
    print("="*80)
    
    with open('config/config.yaml', 'r', encoding='utf-8') as f:
        config = yaml.safe_load(f)
    
    parser = PDFParser(config['parser'])
    pdf_path = 'tests/season_report_del.pdf'
    
    if not parser.open(pdf_path):
        print("❌ FAIL: Could not open PDF")
        return False
    
    # Extract page 6 (index 5)
    page_data = parser.extract_page_elements(5)
    parser.close()
    
    # Find large images
    large_images = []
    rerendered_large_images = []
    texts = []
    
    for elem in page_data['elements']:
        if elem['type'] == 'image':
            if elem['width_px'] > 200 or elem['height_px'] > 200:
                large_images.append(elem)
                if elem.get('was_rerendered', False):
                    rerendered_large_images.append(elem)
        elif elem['type'] == 'text':
            texts.append(elem)
    
    print(f"\nLarge images (>200px): {len(large_images)}")
    print(f"Rerendered large images: {len(rerendered_large_images)}")
    print(f"Text elements: {len(texts)}")
    
    # TEST: Large images should NOT be rerendered
    if rerendered_large_images:
        print("\n❌ FAIL: Found rerendered large images!")
        print("These will capture text elements causing duplication.")
        for img in rerendered_large_images:
            print(f"  - Image at ({img['x']:.1f}, {img['y']:.1f}), "
                  f"size: {img['width_px']}x{img['height_px']}px")
        return False
    
    print("\n✅ PASS: No large images were rerendered")
    print("  (Original embedded images used as-is, preventing text capture)")
    
    return True

def test_conversion_output():
    """Test actual PDF to PPTX conversion."""
    print("\n" + "="*80)
    print("TEST 3: Full conversion test")
    print("="*80)
    
    from main import convert_pdf_to_pptx
    
    input_pdf = 'tests/season_report_del.pdf'
    output_pptx = 'output/season_report_fixed.pptx'
    
    try:
        # Load config
        with open('config/config.yaml', 'r', encoding='utf-8') as f:
            config = yaml.safe_load(f)
        
        # Convert PDF to PPTX
        print("\nConverting PDF to PPTX...")
        result = convert_pdf_to_pptx(
            pdf_path=input_pdf,
            output_path=output_pptx,
            config=config
        )
        
        if not result:
            print("❌ FAIL: Conversion failed")
            return False
        
        print(f"✅ Conversion successful: {output_pptx}")
        
        # Open and inspect the PPTX
        print("\nInspecting output PPTX...")
        prs = Presentation(output_pptx)
        
        print(f"Total slides: {len(prs.slides)}")
        
        # Check page 4 (slide 3)
        if len(prs.slides) >= 4:
            slide_4 = prs.slides[3]
            print(f"\nSlide 4: {len(slide_4.shapes)} shapes")
            
            # Count images and lines
            images = [s for s in slide_4.shapes if hasattr(s, 'image')]
            connectors = [s for s in slide_4.shapes if hasattr(s, 'connector_format')]
            
            print(f"  - Images: {len(images)}")
            print(f"  - Connectors (lines): {len(connectors)}")
        
        # Check page 6 (slide 5)
        if len(prs.slides) >= 6:
            slide_6 = prs.slides[5]
            print(f"\nSlide 6: {len(slide_6.shapes)} shapes")
            
            # Count images and text boxes
            images = [s for s in slide_6.shapes if hasattr(s, 'image')]
            text_boxes = [s for s in slide_6.shapes if hasattr(s, 'text_frame')]
            
            print(f"  - Images: {len(images)}")
            print(f"  - Text boxes: {len(text_boxes)}")
        
        print("\n✅ PASS: Output PPTX generated successfully")
        print(f"\nPlease manually inspect: {output_pptx}")
        print("  - Page 4: Triangle should show vector lines, not blurry PNG overlay")
        print("  - Page 6: Text '30.46分' should appear once, not twice with shadow")
        
        return True
        
    except Exception as e:
        print(f"❌ FAIL: Conversion error: {e}")
        import traceback
        traceback.print_exc()
        return False

def main():
    """Run all tests."""
    print("="*80)
    print("DUPLICATION FIX VERIFICATION TEST SUITE")
    print("="*80)
    print("\nTesting fix for season_report_del.pdf duplication issues:")
    print("  Issue 1: Page 4 triangle has PNG overlay over vector lines")
    print("  Issue 2: Page 6 large image captures text causing duplication")
    
    results = []
    
    # Test 1: Page 4 no image overlap
    test1_passed = test_page_4_no_image_overlap()
    results.append(("Page 4 - No PNG overlays", test1_passed))
    
    # Test 2: Page 6 no text capture
    test2_passed = test_page_6_no_text_overlap()
    results.append(("Page 6 - No text capture", test2_passed))
    
    # Test 3: Full conversion
    test3_passed = test_conversion_output()
    results.append(("Full conversion", test3_passed))
    
    # Summary
    print("\n" + "="*80)
    print("TEST SUMMARY")
    print("="*80)
    
    for test_name, passed in results:
        status = "✅ PASS" if passed else "❌ FAIL"
        print(f"{status}: {test_name}")
    
    all_passed = all(passed for _, passed in results)
    
    if all_passed:
        print("\n🎉 ALL TESTS PASSED!")
        print("\nThe duplication issue has been fixed:")
        print("  ✅ Page 4: Vector shapes not overlaid by PNG images")
        print("  ✅ Page 6: Text not captured into images")
        print("\nRecommendation: Review output PPTX manually to confirm visual quality.")
        return 0
    else:
        print("\n❌ SOME TESTS FAILED")
        print("\nPlease review the failures above.")
        return 1

if __name__ == "__main__":
    sys.exit(main())

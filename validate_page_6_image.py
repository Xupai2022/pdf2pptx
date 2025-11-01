#!/usr/bin/env python3
"""
Validate that page 6 has high-quality image without text duplication.
"""

from pptx import Presentation
from PIL import Image
import io

def validate_page_6_image(pptx_path):
    """Validate that page 6 has high-quality image."""
    print("="*80)
    print("Validating Page 6 - Image Quality and No Text Duplication")
    print("="*80)
    
    prs = Presentation(pptx_path)
    slide = prs.slides[5]  # Page 6 (0-indexed)
    
    # Count images
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    images = []
    text_boxes = []
    text_with_3046 = []
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            images.append(shape)
            
            # Get image details
            try:
                img_blob = shape.image.blob
                pil_img = Image.open(io.BytesIO(img_blob))
                print(f"\n📷 Image found:")
                print(f"  Position: ({shape.left}, {shape.top})")
                print(f"  Size: {shape.width} x {shape.height} EMUs")
                print(f"  PIL Size: {pil_img.size}")
                print(f"  PIL Mode: {pil_img.mode}")
                
                # Check if high quality (pixel size should be large for rerendered image)
                if pil_img.width >= 800 or pil_img.height >= 600:
                    print(f"  ✓ High quality (rerendered with zoom=4.0)")
                else:
                    print(f"  ⚠ Original quality (not rerendered)")
            except Exception as e:
                print(f"  Error reading image: {e}")
        
        elif hasattr(shape, "text_frame"):
            text = shape.text
            text_boxes.append((shape, text))
            
            # Check for "30.46" or its parts ("30", "46", "分")
            if "30.46" in text or "30." in text or "46" == text.strip() or "分" == text.strip():
                text_with_3046.append({
                    'shape': shape,
                    'text': text,
                    'bbox': (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)
                })
    
    print(f"\n📝 Text analysis:")
    print(f"  Total text boxes: {len(text_boxes)}")
    print(f"  Text boxes containing '30.46': {len(text_with_3046)}")
    
    if text_with_3046:
        for t in text_with_3046:
            print(f"\n  Text: '{t['text']}'")
            print(f"    Position: ({t['bbox'][0]}, {t['bbox'][1]})")
            print(f"    Size: {t['shape'].width} x {t['shape'].height} EMUs")
    
    # Check for duplication
    # Count instances of each part (30, 46, 分 should each appear once)
    count_30 = sum(1 for t in text_with_3046 if "30" in t['text'])
    count_46 = sum(1 for t in text_with_3046 if "46" == t['text'].strip())
    count_fen = sum(1 for t in text_with_3046 if "分" == t['text'].strip())
    
    print(f"\n  Count '30': {count_30}")
    print(f"  Count '46': {count_46}")
    print(f"  Count '分': {count_fen}")
    
    # Check if there's duplication (each part should appear exactly once)
    has_duplication = (count_30 > 1 or count_46 > 1 or count_fen > 1)
    has_all_parts = (count_30 >= 1 and count_46 >= 1 and count_fen >= 1)
    
    if has_duplication:
        print(f"\n❌ FAILURE: Text duplication detected!")
        return False
    elif not has_all_parts:
        print(f"\n⚠️  WARNING: Not all text parts found (may have been filtered out)")
        return False
    else:
        print(f"\n✅ SUCCESS: All text parts found exactly once - NO DUPLICATION!")
    
    # Check if there's at least one large image
    has_large_image = any(img.width > 2000000 for img in images)  # > ~2 inches
    if has_large_image:
        print(f"✅ SUCCESS: Found large high-quality image on page 6")
        return True
    else:
        print(f"⚠️  WARNING: No large high-quality image found")
        return False


if __name__ == "__main__":
    result = validate_page_6_image("output/season_report_final.pptx")
    exit(0 if result else 1)

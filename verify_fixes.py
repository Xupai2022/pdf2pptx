#!/usr/bin/env python3
"""
验证修复效果
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def verify_page15_rotation(pptx_path):
    """验证第15页文字旋转"""
    prs = Presentation(pptx_path)
    slide = prs.slides[14]  # Page 15 (0-indexed)
    
    logger.info("=" * 80)
    logger.info("验证第15页：X轴标签旋转")
    logger.info("=" * 80)
    
    rotated_texts = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text:
            # 检查是否包含X轴标签文本
            text = shape.text
            if "10.74" in text or "xos-" in text:
                rotation = shape.rotation if hasattr(shape, 'rotation') else 0
                rotated_texts.append({
                    'text': text[:30],
                    'rotation': rotation
                })
                logger.info(f"文本: '{text[:30]}...'")
                logger.info(f"  旋转角度: {rotation}°")
    
    # 检查是否有-45度或45度旋转
    has_45_rotation = any(abs(abs(t['rotation']) - 45) < 1 for t in rotated_texts)
    
    if has_45_rotation:
        logger.info("✅ 发现45度旋转的文本")
        return True
    else:
        logger.warning("⚠️ 未发现45度旋转的文本")
        return False


def verify_page15_star(pptx_path):
    """验证第15页星星形状"""
    prs = Presentation(pptx_path)
    slide = prs.slides[14]  # Page 15 (0-indexed)
    
    logger.info("\n" + "=" * 80)
    logger.info("验证第15页：星星形状")
    logger.info("=" * 80)
    
    star_count = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if hasattr(shape, 'auto_shape_type'):
                # 检查是否是星星形状
                shape_type_str = str(shape.auto_shape_type)
                if 'STAR' in shape_type_str:
                    star_count += 1
                    logger.info(f"找到星星形状: {shape_type_str}")
                    logger.info(f"  位置: ({shape.left/914400:.2f}in, {shape.top/914400:.2f}in)")
                    logger.info(f"  尺寸: {shape.width/914400:.2f}in x {shape.height/914400:.2f}in")
    
    if star_count > 0:
        logger.info(f"✅ 找到 {star_count} 个星星形状")
        return True
    else:
        logger.warning("⚠️ 未找到星星形状")
        return False


def verify_page6_no_zero_rectangles(pptx_path):
    """验证第6页没有零尺寸矩形"""
    prs = Presentation(pptx_path)
    slide = prs.slides[5]  # Page 6 (0-indexed)
    
    logger.info("\n" + "=" * 80)
    logger.info("验证第6页：无零尺寸矩形")
    logger.info("=" * 80)
    
    zero_count = 0
    total_shapes = 0
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            total_shapes += 1
            if shape.width == 0 or shape.height == 0:
                zero_count += 1
                logger.warning(f"  发现零尺寸矩形: {shape.name}")
                logger.warning(f"    尺寸: {shape.width} x {shape.height}")
    
    logger.info(f"总形状数: {total_shapes}")
    logger.info(f"零尺寸矩形数: {zero_count}")
    
    if zero_count == 0:
        logger.info("✅ 没有零尺寸矩形")
        return True
    else:
        logger.warning(f"⚠️ 发现 {zero_count} 个零尺寸矩形")
        return False


def compare_element_counts(baseline_path, fixed_path):
    """比较基线和修复版本的元素数量"""
    logger.info("\n" + "=" * 80)
    logger.info("比较元素数量（基线 vs 修复）")
    logger.info("=" * 80)
    
    baseline_prs = Presentation(baseline_path)
    fixed_prs = Presentation(fixed_path)
    
    all_match = True
    
    for page_num in range(min(len(baseline_prs.slides), len(fixed_prs.slides))):
        baseline_slide = baseline_prs.slides[page_num]
        fixed_slide = fixed_prs.slides[page_num]
        
        baseline_count = len(baseline_slide.shapes)
        fixed_count = len(fixed_slide.shapes)
        
        # 对于第6页，修复版应该少一些元素（去掉了零尺寸矩形）
        if page_num == 5:  # Page 6
            if fixed_count < baseline_count:
                logger.info(f"第{page_num+1}页: {baseline_count} → {fixed_count} (-{baseline_count-fixed_count}) ✅ 预期减少")
            else:
                logger.warning(f"第{page_num+1}页: {baseline_count} → {fixed_count} ⚠️ 未减少")
                all_match = False
        elif page_num == 14:  # Page 15
            # 第15页：过滤了零尺寸形状，所以元素会减少
            if fixed_count < baseline_count:
                logger.info(f"第{page_num+1}页: {baseline_count} → {fixed_count} (-{baseline_count-fixed_count}) ✅ 过滤零尺寸形状")
            elif fixed_count == baseline_count:
                logger.info(f"第{page_num+1}页: {baseline_count} → {fixed_count} ✅")
            else:
                logger.warning(f"第{page_num+1}页: {baseline_count} → {fixed_count} ⚠️ 元素增加")
                all_match = False
        else:
            # 其他页面：修复版应该少一些元素（去掉了零尺寸形状）
            # 这是预期的改进，不是bug
            if fixed_count <= baseline_count:
                logger.debug(f"第{page_num+1}页: {baseline_count} → {fixed_count} ✅")
            else:
                logger.warning(f"第{page_num+1}页: {baseline_count} → {fixed_count} ⚠️ 元素增加了")
                all_match = False
    
    return all_match


def main():
    baseline_path = "output_baseline.pptx"
    fixed_path = "output_fixed.pptx"
    
    results = {
        'rotation': verify_page15_rotation(fixed_path),
        'star': verify_page15_star(fixed_path),
        'no_zero': verify_page6_no_zero_rectangles(fixed_path),
        'counts': compare_element_counts(baseline_path, fixed_path)
    }
    
    logger.info("\n" + "=" * 80)
    logger.info("验证结果总结")
    logger.info("=" * 80)
    
    all_passed = all(results.values())
    
    for test_name, passed in results.items():
        status = "✅ 通过" if passed else "❌ 失败"
        logger.info(f"{test_name}: {status}")
    
    if all_passed:
        logger.info("\n🎉 所有验证通过！")
    else:
        logger.warning("\n⚠️ 部分验证失败，请检查")
    
    return all_passed


if __name__ == '__main__':
    import sys
    success = main()
    sys.exit(0 if success else 1)

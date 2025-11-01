#!/usr/bin/env python3
"""
Validate that page 4 of the generated PPTX has two diagonal lines.
"""

from pptx import Presentation

def validate_page_4_lines(pptx_path):
    """Validate that page 4 has two diagonal lines."""
    print("="*80)
    print("Validating Page 4 - Triangle Diagonal Lines")
    print("="*80)
    
    prs = Presentation(pptx_path)
    slide = prs.slides[3]  # Page 4 (0-indexed)
    
    # Count connectors (lines)
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    connectors = []
    diagonal_connectors = []
    
    for shape in slide.shapes:
        # Check if it's a connector (line)
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            connectors.append(shape)
            
            # Check if diagonal (both width and height > 0)
            width = abs(shape.width)
            height = abs(shape.height)
            
            # Diagonal if both dimensions are significant
            if width > 100000 and height > 100000:  # EMUs (>= 0.1 inch)
                diagonal_connectors.append(shape)
                print(f"\n✓ Found diagonal line:")
                print(f"  Position: ({shape.left}, {shape.top})")
                print(f"  Size: {width} x {height} EMUs")
                if hasattr(shape.line, 'color') and hasattr(shape.line.color, 'rgb'):
                    print(f"  Color: {shape.line.color.rgb}")
                if hasattr(shape.line, 'width'):
                    print(f"  Width: {shape.line.width}")
    
    print(f"\n" + "="*80)
    print(f"Summary:")
    print(f"  Total connectors (lines): {len(connectors)}")
    print(f"  Diagonal lines: {len(diagonal_connectors)}")
    
    if len(diagonal_connectors) >= 2:
        print(f"\n✅ SUCCESS: Found {len(diagonal_connectors)} diagonal lines (expected >= 2)")
        print(f"   Page 4 has the triangle side lines properly rendered!")
        return True
    else:
        print(f"\n❌ FAILURE: Found only {len(diagonal_connectors)} diagonal lines (expected >= 2)")
        print(f"   The triangle side lines are still missing!")
        return False


if __name__ == "__main__":
    result = validate_page_4_lines("output/season_report_fixed.pptx")
    exit(0 if result else 1)

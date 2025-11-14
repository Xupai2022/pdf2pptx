"""
测试第16页转换成PPTX
"""
import sys
sys.path.append('src')

from src.parser.pdf_parser import PDFParser
from src.generator.pptx_generator import PPTXGenerator

def test_page16():
    pdf_path = "tests/season_report_del.pdf"
    output_path = "output/season_report_page16_test.pptx"
    
    print("解析PDF...")
    parser = PDFParser({'dpi': 300, 'extract_images': True})
    parser.open(pdf_path)
    
    # 只提取第16页
    page_data = parser.extract_page_elements(15)  # 0-indexed
    
    parser.close()
    
    print(f"\n第16页元素统计:")
    tables = [e for e in page_data['elements'] if e['type'] == 'table']
    shapes = [e for e in page_data['elements'] if e['type'] == 'shape']
    texts = [e for e in page_data['elements'] if e['type'] == 'text']
    
    print(f"  表格: {len(tables)}")
    print(f"  形状: {len(shapes)}")
    print(f"  文本: {len(texts)}")
    
    if tables:
        for i, table in enumerate(tables):
            print(f"\n  表格{i+1}详情:")
            print(f"    bbox: {(table['x'], table['y'], table['x2'], table['y2'])}")
            print(f"    行数: {table['rows']}")
            print(f"    列数: {table['cols']}")
            
            grid = table.get('grid', [])
            print(f"    实际网格: {len(grid)} 行")
            
            print(f"    行高: {table.get('row_heights', [])}")
            print(f"    列宽: {table.get('col_widths', [])}")
            
            print(f"\n    网格内容:")
            for row_idx, row in enumerate(grid):
                texts_in_row = []
                for col_idx, cell in enumerate(row):
                    text = cell.get('text', '').strip()
                    height = cell.get('height', 0)
                    width = cell.get('width', 0)
                    if text:
                        texts_in_row.append(f'{text[:20]}({width:.0f}x{height:.0f})')
                    else:
                        texts_in_row.append(f'空({width:.0f}x{height:.0f})')
                print(f"      行{row_idx+1}: {' | '.join(texts_in_row)}")
    
    # 生成PPTX
    print(f"\n生成PPTX...")
    from src.rebuilder.slide_model import SlideModel
    from src.rebuilder.coordinate_mapper import CoordinateMapper
    from src.analyzer.layout_analyzer_v2 import LayoutAnalyzerV2
    
    # 分析布局
    analyzer = LayoutAnalyzerV2({})
    layout_data = analyzer.analyze_page(page_data)
    
    # 转换为slide model
    mapper = CoordinateMapper({})
    slide_model = mapper.create_slide_model(layout_data)
    
    generator = PPTXGenerator({})
    generator.generate_from_models([slide_model])
    generator.save(output_path)
    print(f"✓ PPTX已生成: {output_path}")
    
    # 读取生成的PPTX检查表格
    print(f"\n检查生成的PPTX...")
    from pptx import Presentation
    prs = Presentation(output_path)
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n  幻灯片{slide_idx + 1}:")
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                print(f"    表格: {len(table.rows)}行 x {len(table.columns)}列")
                print(f"    位置: ({shape.left}, {shape.top}), 尺寸: ({shape.width}, {shape.height})")
                
                # 显示行高
                print(f"    行高信息:")
                for row_idx, row in enumerate(table.rows):
                    print(f"      行{row_idx+1}: 高度={row.height} ({row.height/914400:.1f}pt)")
                
                # 显示前几行内容
                print(f"    表格内容(前{min(len(table.rows), 7)}行):")
                for row_idx in range(min(len(table.rows), 7)):
                    row = table.rows[row_idx]
                    texts = []
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            texts.append(text[:15])
                        else:
                            texts.append('(空)')
                    print(f"      行{row_idx+1}: {' | '.join(texts)}")

if __name__ == "__main__":
    test_page16()

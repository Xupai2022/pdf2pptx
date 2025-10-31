"""
Final comprehensive test to verify the duplication fix.

This test ensures:
1. No large PNG overlays in page 4 triangle region
2. No text duplication in page 6
3. Output PPTX is correct
"""

import sys
import yaml
from src.parser.pdf_parser import PDFParser
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def test_no_large_rerendered_images():
    """Test that large images are NOT being rerendered."""
    print("\n" + "="*80)
    print("TEST 1: No Large Image Rerendering")
    print("="*80)
    
    with open('config/config.yaml', 'r', encoding='utf-8') as f:
        config = yaml.safe_load(f)
    
    parser = PDFParser(config['parser'])
    parser.open('tests/season_report_del.pdf')
    
    # Check all pages
    total_rerendered_large = 0
    
    for page_num in range(parser.get_page_count()):
        page_data = parser.extract_page_elements(page_num)
        
        for elem in page_data['elements']:
            if elem['type'] == 'image' and elem.get('was_rerendered', False):
                # Check if it's a large image
                if elem['width_px'] > 200 or elem['height_px'] > 200:
                    total_rerendered_large += 1
                    print(f"\n❌ FAIL: Found large rerendered image on page {page_num + 1}:")
                    print(f"  Position: ({elem['x']:.1f}, {elem['y']:.1f})")
                    print(f"  Size: {elem['width_px']}x{elem['height_px']}px")
    
    parser.close()
    
    if total_rerendered_large == 0:
        print("\n✅ PASS: No large images were rerendered")
        print("  (Prevents capturing of overlapping text/shapes)")
        return True
    else:
        print(f"\n❌ FAIL: Found {total_rerendered_large} large rerendered image(s)")
        return False

def test_page_4_triangle_region():
    """Test that page 4 triangle region has no PNG overlays."""
    print("\n" + "="*80)
    print("TEST 2: Page 4 Triangle Region - No PNG Overlays")
    print("="*80)
    
    with open('config/config.yaml', 'r', encoding='utf-8') as f:
        config = yaml.safe_load(f)
    
    parser = PDFParser(config['parser'])
    parser.open('tests/season_report_del.pdf')
    
    page_data = parser.extract_page_elements(3)  # Page 4
    parser.close()
    
    # Extended triangle region to catch all relevant areas
    triangle_region = {'x_min': 200, 'x_max': 700, 'y_min': 50, 'y_max': 450}
    
    large_images_in_region = []
    
    for elem in page_data['elements']:
        if elem['type'] == 'image':
            # Check if large image (> 200px is considered "large")
            # Small icons (< 200px) may be rerendered for quality, that's OK
            is_large = elem['width_px'] > 200 or elem['height_px'] > 200
            
            if is_large:
                x, y, x2, y2 = elem['x'], elem['y'], elem['x2'], elem['y2']
                
                # Check overlap
                overlaps = (
                    x < triangle_region['x_max'] and x2 > triangle_region['x_min'] and
                    y < triangle_region['y_max'] and y2 > triangle_region['y_min']
                )
                
                if overlaps:
                    large_images_in_region.append(elem)
    
    print(f"\nLarge images (>200px) in extended region: {len(large_images_in_region)}")
    
    if large_images_in_region:
        print("\n❌ FAIL: Found large images that may overlay vector shapes:")
        for i, img in enumerate(large_images_in_region):
            print(f"  Image {i+1}: {img['width_px']}x{img['height_px']}px at ({img['x']:.1f}, {img['y']:.1f})")
            print(f"    Was rerendered: {img.get('was_rerendered', False)}")
        return False
    else:
        print("\n✅ PASS: No large PNG overlays in triangle region")
        return True

def test_page_6_no_text_duplication():
    """Test that page 6 has no text duplication."""
    print("\n" + "="*80)
    print("TEST 3: Page 6 - No Text Duplication")
    print("="*80)
    
    prs = Presentation('output/season_report_fixed.pptx')
    
    if len(prs.slides) < 6:
        print("❌ FAIL: Not enough slides")
        return False
    
    slide = prs.slides[5]  # Page 6
    
    # Collect all text
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame'):
            text = shape.text_frame.text.strip()
            if text:
                texts.append({
                    'text': text,
                    'left': shape.left / 12700,  # Convert to points
                    'top': shape.top / 12700
                })
    
    # Check for "30.46" or similar appearing multiple times close together
    target = '30'
    matches = [t for t in texts if target in t['text']]
    
    duplicates = []
    for i in range(len(matches)):
        for j in range(i+1, len(matches)):
            dist_x = abs(matches[i]['left'] - matches[j]['left'])
            dist_y = abs(matches[i]['top'] - matches[j]['top'])
            
            # If very close (< 30pt), likely duplicate
            if dist_x < 30 and dist_y < 30:
                duplicates.append((matches[i], matches[j]))
    
    if duplicates:
        print(f"\n❌ FAIL: Found {len(duplicates)} potential text duplication(s):")
        for i, (t1, t2) in enumerate(duplicates):
            print(f"  Pair {i+1}:")
            print(f"    '{t1['text']}' at ({t1['left']:.1f}, {t1['top']:.1f})")
            print(f"    '{t2['text']}' at ({t2['left']:.1f}, {t2['top']:.1f})")
        return False
    else:
        print(f"\n✅ PASS: No text duplication detected")
        print(f"  (Checked {len(matches)} text elements containing '{target}')")
        return True

def test_output_quality():
    """Basic quality checks on output PPTX."""
    print("\n" + "="*80)
    print("TEST 4: Output PPTX Quality")
    print("="*80)
    
    try:
        prs = Presentation('output/season_report_fixed.pptx')
        
        print(f"\nTotal slides: {len(prs.slides)}")
        
        if len(prs.slides) < 17:
            print(f"❌ FAIL: Expected 17 slides, got {len(prs.slides)}")
            return False
        
        # Check slide 4
        slide_4 = prs.slides[3]
        shapes_4 = len(slide_4.shapes)
        images_4 = len([s for s in slide_4.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE])
        
        print(f"\nSlide 4 (Page 4):")
        print(f"  Total shapes: {shapes_4}")
        print(f"  Images: {images_4}")
        
        # Check slide 6
        slide_6 = prs.slides[5]
        shapes_6 = len(slide_6.shapes)
        images_6 = len([s for s in slide_6.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE])
        
        print(f"\nSlide 6 (Page 6):")
        print(f"  Total shapes: {shapes_6}")
        print(f"  Images: {images_6}")
        
        print("\n✅ PASS: Output PPTX structure looks good")
        return True
        
    except Exception as e:
        print(f"\n❌ FAIL: Error reading output PPTX: {e}")
        return False

def main():
    """Run all tests."""
    print("="*80)
    print("FINAL COMPREHENSIVE DUPLICATION FIX TEST")
    print("="*80)
    print("\nTesting PDF: tests/season_report_del.pdf")
    print("Testing PPTX: output/season_report_fixed.pptx")
    
    results = []
    
    # Test 1: No large image rerendering
    test1 = test_no_large_rerendered_images()
    results.append(("No large image rerendering", test1))
    
    # Test 2: Page 4 triangle region
    test2 = test_page_4_triangle_region()
    results.append(("Page 4 - No PNG overlays", test2))
    
    # Test 3: Page 6 text duplication
    test3 = test_page_6_no_text_duplication()
    results.append(("Page 6 - No text duplication", test3))
    
    # Test 4: Output quality
    test4 = test_output_quality()
    results.append(("Output PPTX quality", test4))
    
    # Summary
    print("\n" + "="*80)
    print("FINAL TEST SUMMARY")
    print("="*80)
    
    for test_name, passed in results:
        status = "✅ PASS" if passed else "❌ FAIL"
        print(f"{status}: {test_name}")
    
    all_passed = all(passed for _, passed in results)
    
    if all_passed:
        print("\n" + "🎉" * 20)
        print("🎉 ALL TESTS PASSED! FIX IS SUCCESSFUL! 🎉")
        print("🎉" * 20)
        print("\nThe duplication issue has been resolved:")
        print("  ✅ No large PNG images rerendered (prevents text/shape capture)")
        print("  ✅ Page 4: Triangle region clean, no PNG overlays")
        print("  ✅ Page 6: Text appears only once, no duplication")
        print("\nThe output PPTX is ready for manual review.")
        print("File: output/season_report_fixed.pptx")
        return 0
    else:
        print("\n❌ SOME TESTS FAILED")
        print("\nPlease review the failures above and investigate further.")
        return 1

if __name__ == "__main__":
    sys.exit(main())

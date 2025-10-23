#!/usr/bin/env python3
"""
深度分析所有问题：字体大小、边框、透明度等
"""

import yaml
from pathlib import Path
from pptx import Presentation
from lxml import etree

def analyze_html_specs():
    """分析 HTML 规格"""
    print("="*80)
    print("HTML 规格分析 (tests/slide11_reference.html)")
    print("="*80)
    
    specs = {
        "页面尺寸": "1920px × 1080px",
        "字体大小": {
            "h1": "48px (主标题)",
            "h2": "36px (副标题)", 
            "h3": "28px (卡片标题)",
            "p": "25px (正文)",
            ".text-4xl": "36px (大数字)",
            ".text-2xl": "24px (卡片标题)",
            ".text-lg": "18px (说明文字)",
            ".text-sm": "14px (小字)",
        },
        "top-bar": {
            "高度": "10px",
            "颜色": "rgb(10, 66, 117) = #0A4275",
            "透明度": "无 (实色)"
        },
        "stat-card": {
            "背景": "rgba(10, 66, 117, 0.08) - 8% 不透明度",
            "左边框": "4px solid rgb(10, 66, 117) - 无透明度",
            "圆角": "8px",
            "内边距": "15px 20px"
        },
        "data-card": {
            "背景": "rgba(10, 66, 117, 0.03) - 3% 不透明度",
            "左边框": "4px solid rgb(10, 66, 117) - 无透明度",
            "圆角": "8px",
            "内边距": "15px 20px"
        },
        "risk-level badges": {
            ".risk-high": "rgba(220, 38, 38, 0.1) - 10% 不透明度",
            ".risk-medium": "rgba(245, 158, 11, 0.1) - 10% 不透明度",
            ".risk-low": "rgba(59, 130, 246, 0.1) - 10% 不透明度"
        }
    }
    
    for key, value in specs.items():
        if isinstance(value, dict):
            print(f"\n{key}:")
            for k, v in value.items():
                print(f"  {k}: {v}")
        else:
            print(f"{key}: {value}")
    
    return specs


def analyze_pdf_extraction():
    """分析 PDF 提取结果"""
    print("\n" + "="*80)
    print("PDF 提取分析")
    print("="*80)
    
    from src.parser.pdf_parser import PDFParser
    
    with open('config/config.yaml', 'r') as f:
        config = yaml.safe_load(f)
    
    parser = PDFParser(config['parser'])
    parser.open('tests/test_sample.pdf')
    pages = parser.extract_all_pages()
    parser.close()
    
    page = pages[0]
    
    # Font sizes
    font_sizes = {}
    for elem in page['elements']:
        if elem.get('type') == 'text':
            size = round(elem.get('font_size', 0), 1)
            if size > 0:
                font_sizes[size] = font_sizes.get(size, 0) + 1
    
    print("\n提取的字体大小分布:")
    for size in sorted(font_sizes.keys(), reverse=True):
        print(f"  {size}pt: {font_sizes[size]} 个元素")
    
    # Shapes
    shapes = [e for e in page['elements'] if e.get('type') == 'shape']
    print(f"\n提取的形状: {len(shapes)} 个")
    
    # Analyze specific shapes
    top_bar = None
    borders = []
    backgrounds = []
    
    for shape in shapes:
        h = shape.get('y2', 0) - shape.get('y', 0)
        w = shape.get('x2', 0) - shape.get('x', 0)
        
        # Top bar (full width, ~10px high)
        if h < 15 and w > 1400:
            top_bar = shape
        # Borders (narrow, tall)
        elif w < 10 and h > 50:
            borders.append(shape)
        # Backgrounds (large rectangles)
        elif w > 200 and h > 50:
            backgrounds.append(shape)
    
    print(f"\n识别的元素:")
    if top_bar:
        h = top_bar['y2'] - top_bar['y']
        print(f"  Top bar: 高度 {h:.1f}pt, 颜色 {top_bar.get('fill_color')}")
    print(f"  边框: {len(borders)} 个")
    print(f"  背景: {len(backgrounds)} 个")
    
    return {
        'font_sizes': font_sizes,
        'top_bar': top_bar,
        'borders': borders,
        'backgrounds': backgrounds
    }


def analyze_pptx_output(pptx_path='output_transparent.pptx'):
    """分析 PPTX 输出"""
    print("\n" + "="*80)
    print("PPTX 输出分析")
    print("="*80)
    
    if not Path(pptx_path).exists():
        print(f"  ⚠️ 文件不存在: {pptx_path}")
        return
    
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    
    # Dimensions
    width = prs.slide_width / 914400
    height = prs.slide_height / 914400
    print(f"\n幻灯片尺寸: {width:.3f}\" × {height:.3f}\" ({width*144:.0f}px × {height*144:.0f}px @ 144 DPI)")
    
    # Font sizes
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    font_sizes = {}
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            try:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            size = round(run.font.size.pt, 1)
                            font_sizes[size] = font_sizes.get(size, 0) + 1
            except:
                pass
    
    print("\n字体大小分布:")
    for size in sorted(font_sizes.keys(), reverse=True):
        print(f"  {size}pt: {font_sizes[size]} 个")
    
    # Transparency
    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    transparent_shapes = 0
    for shape in slide.shapes:
        if hasattr(shape, 'element'):
            try:
                spPr = shape.element.spPr
                solidFill = spPr.find('.//a:solidFill', ns)
                if solidFill is not None:
                    srgbClr = solidFill.find('.//a:srgbClr', ns)
                    if srgbClr is not None:
                        alpha = srgbClr.find('.//a:alpha', ns)
                        if alpha is not None:
                            transparent_shapes += 1
            except:
                pass
    
    print(f"\n透明形状: {transparent_shapes} 个")


def identify_issues():
    """识别所有问题"""
    print("\n" + "="*80)
    print("问题识别")
    print("="*80)
    
    issues = [
        {
            "问题": "字体大小不匹配",
            "详情": "PDF 提取的字体比 HTML 小约 25-30%",
            "原因": "可能是 PDF points 到 screen pixels 的转换问题",
            "解决方案": "需要在 style_mapper 中添加字体大小缩放因子"
        },
        {
            "问题": "边框识别不准确",
            "详情": "4px 实色边框可能被识别为背景或其他元素",
            "原因": "布局分析器可能将细条形状归类错误",
            "解决方案": "改进形状角色检测，专门识别边框（宽度<10px，高度>50px）"
        },
        {
            "问题": "透明度应用不当",
            "详情": "边框应该是实色，但可能被错误地应用了透明度",
            "原因": "transparency_map 将 #094174 映射到 0.08，但边框应该是 1.0",
            "解决方案": "根据形状角色区分：背景用透明，边框用实色"
        },
        {
            "问题": "top-bar 尺寸",
            "详情": "应该精确是 10px 高度",
            "原因": "PDF 提取可能不是精确的 10pt",
            "解决方案": "需要验证并可能手动调整"
        }
    ]
    
    for i, issue in enumerate(issues, 1):
        print(f"\n问题 {i}: {issue['问题']}")
        print(f"  详情: {issue['详情']}")
        print(f"  原因: {issue['原因']}")
        print(f"  解决方案: {issue['解决方案']}")


if __name__ == "__main__":
    analyze_html_specs()
    pdf_data = analyze_pdf_extraction()
    analyze_pptx_output()
    identify_issues()
    
    print("\n" + "="*80)
    print("总结")
    print("="*80)
    print("""
需要修复的关键问题:

1. 字体大小缩放
   - 添加 font_size_scale 配置参数 (建议 1.33 倍)
   - HTML 48px → PDF 36pt → PPTX 48pt

2. 边框 vs 背景区分
   - 边框：窄条 (width < 10px)，实色 #094174，无透明度
   - 背景：大块，#094174，0.08 透明度

3. 形状角色识别改进
   - border: width < 10px, height > 50px
   - decoration (top-bar): width > 1400px, height < 15px
   - card_background: width > 200px, height > 50px

4. 透明度逻辑
   - 只对 card_background 和 risk badges 应用透明度
   - border 和 decoration 保持实色
""")

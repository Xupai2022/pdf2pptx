"""
Verify that PNG images in the generated PPTX have proper transparency
"""

from pptx import Presentation
from PIL import Image
import io
import sys

def verify_pptx_images(pptx_path, pages_to_check):
    """
    Verify PNG images in specific slides of a PPTX
    """
    prs = Presentation(pptx_path)
    
    print(f"=== Verifying PPTX: {pptx_path} ===\n")
    print(f"Total slides: {len(prs.slides)}")
    
    for page_idx in pages_to_check:
        if page_idx >= len(prs.slides):
            print(f"Slide {page_idx + 1} does not exist (total slides: {len(prs.slides)})")
            continue
        
        slide = prs.slides[page_idx]
        
        print(f"\n{'='*60}")
        print(f"Slide {page_idx + 1}")
        print(f"{'='*60}")
        
        image_count = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                image_count += 1
                
                # Get image data
                image_blob = shape.image.blob
                pil_image = Image.open(io.BytesIO(image_blob))
                
                width, height = pil_image.size
                mode = pil_image.mode
                
                # Get position
                left = shape.left.inches
                top = shape.top.inches
                w = shape.width.inches
                h = shape.height.inches
                
                print(f"\n  Image {image_count}:")
                print(f"    Size: {width}x{height}px")
                print(f"    Mode: {mode}")
                print(f"    Position: ({left:.2f}, {top:.2f}) inches")
                print(f"    Dimensions: {w:.2f}x{h:.2f} inches")
                
                # Check for transparency
                has_alpha = mode in ['RGBA', 'LA', 'PA']
                print(f"    Has alpha channel: {has_alpha}")
                
                if has_alpha:
                    # Sample alpha channel values
                    sample_points = [
                        (0, 0), (width//2, 0), (width-1, 0),
                        (0, height//2), (width//2, height//2), (width-1, height//2),
                        (0, height-1), (width//2, height-1), (width-1, height-1)
                    ]
                    
                    alpha_values = []
                    for x, y in sample_points:
                        try:
                            pixel = pil_image.getpixel((x, y))
                            if len(pixel) >= 4:
                                alpha_values.append(pixel[3])
                        except:
                            continue
                    
                    if alpha_values:
                        avg_alpha = sum(alpha_values) / len(alpha_values)
                        min_alpha = min(alpha_values)
                        max_alpha = max(alpha_values)
                        
                        print(f"    Alpha statistics:")
                        print(f"      Average: {avg_alpha:.1f}/255")
                        print(f"      Range: {min_alpha}-{max_alpha}/255")
                        
                        # Check if truly transparent
                        has_transparency = min_alpha < 255
                        print(f"    ✅ Has transparency: {has_transparency}")
                        
                        if not has_transparency:
                            print(f"    ⚠️  WARNING: Has alpha channel but no transparent pixels!")
                    
                    # Sample RGB values too
                    colors = []
                    for x, y in sample_points:
                        try:
                            pixel = pil_image.getpixel((x, y))
                            colors.append(pixel[:3])
                        except:
                            continue
                    
                    # Check for black pixels
                    black_count = sum(1 for c in colors if all(ch < 30 for ch in c))
                    
                    if black_count > 0:
                        print(f"    Black pixels: {black_count}/{len(colors)}")
                        if black_count >= len(colors) * 0.5:
                            print(f"    ⚠️  WARNING: Still has many black pixels!")
                        else:
                            print(f"    ✅ Black pixels are within normal range")
                else:
                    # No alpha channel - check for black background
                    sample_points = [
                        (0, 0), (width//2, 0), (width-1, 0),
                        (0, height//2), (width//2, height//2), (width-1, height//2),
                        (0, height-1), (width//2, height-1), (width-1, height-1)
                    ]
                    
                    colors = []
                    for x, y in sample_points:
                        try:
                            pixel = pil_image.getpixel((x, y))
                            if not isinstance(pixel, tuple):
                                pixel = (pixel, pixel, pixel)
                            colors.append(pixel)
                        except:
                            continue
                    
                    black_count = sum(1 for c in colors if all(ch < 30 for ch in c[:3]))
                    
                    print(f"    Black pixels: {black_count}/{len(colors)}")
                    
                    if black_count >= len(colors) * 0.5:
                        print(f"    ❌ ERROR: RGB image with black background!")
                    elif black_count > 0:
                        print(f"    ⚠️  Has some black pixels (may be intentional)")
                    else:
                        print(f"    ✅ No black background issue")
        
        print(f"\nTotal images in slide: {image_count}")

if __name__ == "__main__":
    pptx_path = "output/test_fixed.pptx"
    
    # Check pages mentioned: 3, 4, 5, 10, 11, 13 (convert to 0-indexed)
    pages_to_check = [2, 3, 4, 9, 10, 12]  # 0-indexed
    
    verify_pptx_images(pptx_path, pages_to_check)

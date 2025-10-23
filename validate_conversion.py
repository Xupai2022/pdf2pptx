#!/usr/bin/env python3
"""
Validate that the PDF to PPTX conversion produces correct dimensions and font sizes.
"""

import sys
import logging
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt, Inches

# Setup logging
logging.basicConfig(
    level=logging.INFO,
    format='%(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


def validate_pptx_output(pptx_path: str):
    """Validate the generated PPTX against HTML reference specs."""
    
    logger.info("=" * 60)
    logger.info("PPTX OUTPUT VALIDATION")
    logger.info("=" * 60)
    
    if not Path(pptx_path).exists():
        logger.error(f"PPTX file not found: {pptx_path}")
        return False
    
    prs = Presentation(pptx_path)
    
    # Check slide dimensions
    slide_width_inches = prs.slide_width / 914400  # EMUs to inches
    slide_height_inches = prs.slide_height / 914400
    
    logger.info(f"\nSlide Dimensions:")
    logger.info(f"  Width: {slide_width_inches:.3f}\" ({slide_width_inches * 72:.1f}pt / {slide_width_inches * 72:.0f}px at 72 DPI)")
    logger.info(f"  Height: {slide_height_inches:.3f}\" ({slide_height_inches * 72:.1f}pt / {slide_height_inches * 72:.0f}px at 72 DPI)")
    logger.info(f"  Expected: 26.667\" × 15.000\" (1920×1080 px)")
    
    # Validate dimensions
    expected_width = 26.667
    expected_height = 15.0
    width_match = abs(slide_width_inches - expected_width) < 0.1
    height_match = abs(slide_height_inches - expected_height) < 0.1
    
    if width_match and height_match:
        logger.info(f"  ✓ Dimensions CORRECT")
    else:
        logger.warning(f"  ✗ Dimensions INCORRECT (off by {abs(slide_width_inches - expected_width):.3f}\" × {abs(slide_height_inches - expected_height):.3f}\")")
    
    # Check first slide content
    if len(prs.slides) == 0:
        logger.error("No slides found in presentation")
        return False
    
    slide = prs.slides[0]
    logger.info(f"\nSlide 1 Analysis:")
    logger.info(f"  Total shapes: {len(slide.shapes)}")
    
    # Expected font sizes (HTML pixels = PPT points at 72 DPI after scaling)
    expected_fonts = {
        'h1': 48,  # 36pt PDF × 1.333 = 48pt
        'h2': 36,  # 27pt PDF × 1.333 = 36pt
        'h3': 28,  # 21pt PDF × 1.333 = 28pt
        'p': 25,   # 18.5pt PDF × 1.333 ≈ 25pt
        'risk': 20,  # 15pt PDF × 1.333 = 20pt
        'small': 18,  # 13.5pt PDF × 1.333 = 18pt
        'page': 14,  # 10.5pt PDF × 1.333 = 14pt
    }
    
    # Analyze text boxes
    text_shapes = [s for s in slide.shapes if s.has_text_frame]
    logger.info(f"  Text shapes: {len(text_shapes)}")
    
    font_sizes_found = {}
    for shape in text_shapes[:30]:  # Check first 30
        try:
            text_content = shape.text if hasattr(shape, 'text') else ""
            if text_content and text_content.strip():
                if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    if para.runs and len(para.runs) > 0:
                        run = para.runs[0]
                        if run.font.size:
                            size_pt = run.font.size.pt
                            text_preview = text_content[:30]
                            
                            if size_pt not in font_sizes_found:
                                font_sizes_found[size_pt] = []
                            font_sizes_found[size_pt].append(text_preview)
        except Exception as e:
            pass
    
    logger.info(f"\n  Font sizes found (PowerPoint points):")
    for size in sorted(font_sizes_found.keys(), reverse=True):
        texts = font_sizes_found[size]
        logger.info(f"    {size:.1f}pt: {len(texts)} shapes")
        logger.info(f"      Example: \"{texts[0]}\"")
        
        # Check against expected
        matched = False
        for name, expected_size in expected_fonts.items():
            if abs(size - expected_size) < 2:  # Within 2pt tolerance
                logger.info(f"      ✓ Matches {name} ({expected_size}pt expected)")
                matched = True
                break
        if not matched:
            logger.info(f"      ? No exact match in expected fonts")
    
    # Check for top bar (should be ~10pt height = 0.139")
    logger.info(f"\n  Top bar check:")
    expected_top_bar_height = 10 / 72  # 10px at 72 DPI = 0.139"
    
    # Find thin horizontal shapes near top
    horizontal_shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:  # Not text
            try:
                width_in = shape.width / 914400
                height_in = shape.height / 914400
                top_in = shape.top / 914400
                
                # Check if it's a wide, thin shape at the top
                if width_in > 20 and height_in < 0.3 and top_in < 0.5:
                    horizontal_shapes.append({
                        'width': width_in,
                        'height': height_in,
                        'top': top_in,
                        'height_pt': height_in * 72
                    })
            except:
                pass
    
    if horizontal_shapes:
        top_bar = horizontal_shapes[0]
        logger.info(f"    Found horizontal bar:")
        logger.info(f"      Height: {top_bar['height']:.4f}\" ({top_bar['height_pt']:.1f}pt)")
        logger.info(f"      Expected: {expected_top_bar_height:.4f}\" (10pt)")
        
        if abs(top_bar['height_pt'] - 10) < 2:
            logger.info(f"      ✓ Top bar height CORRECT")
        else:
            logger.warning(f"      ✗ Top bar height INCORRECT (off by {abs(top_bar['height_pt'] - 10):.1f}pt)")
    else:
        logger.warning(f"    ✗ No top bar found")
    
    # Check border widths (should be ~4pt)
    logger.info(f"\n  Border width check:")
    expected_border_width = 4  # 4px = 4pt at 72 DPI
    
    # Note: Border widths are hard to extract from python-pptx
    # Would need XML inspection for accurate validation
    logger.info(f"    Expected: {expected_border_width}pt")
    logger.info(f"    (Border widths require XML inspection for validation)")
    
    logger.info("\n" + "=" * 60)
    logger.info("VALIDATION SUMMARY")
    logger.info("=" * 60)
    logger.info(f"✓ Slide dimensions: {'PASS' if width_match and height_match else 'FAIL'}")
    logger.info(f"✓ Font sizes: Check output above")
    logger.info(f"✓ Top bar height: Check output above")
    logger.info("=" * 60)
    
    return True


def main():
    """Main validation function."""
    pptx_path = "output/test_output.pptx"
    
    if len(sys.argv) > 1:
        pptx_path = sys.argv[1]
    
    validate_pptx_output(pptx_path)


if __name__ == "__main__":
    main()

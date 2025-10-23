#!/usr/bin/env python3
"""直接检查PPTX文件的详细内容"""

from pptx import Presentation
from pathlib import Path

def inspect_pptx(path):
    prs = Presentation(path)
    
    print("=" * 60)
    print("PPTX DETAILED INSPECTION")
    print("=" * 60)
    
    print(f"\nSlide dimensions: {prs.slide_width / 914400:.3f}\" × {prs.slide_height / 914400:.3f}\"")
    print(f"Number of slides: {len(prs.slides)}")
    
    if len(prs.slides) > 0:
        slide = prs.slides[0]
        print(f"\nSlide 1 - Total shapes: {len(slide.shapes)}")
        
        print("\n" + "=" * 60)
        print("SHAPES ANALYSIS")
        print("=" * 60)
        
        count = 0
        for i, shape in enumerate(slide.shapes):  # All shapes
            if count >= 15:
                break
            count += 1
            print(f"\nShape {i+1}:")
            print(f"  Type: {shape.shape_type}")
            print(f"  Has text: {shape.has_text_frame}")
            
            if shape.has_text_frame:
                text = shape.text[:50]
                print(f"  Text: \"{text}\"")
                
                # Check paragraphs
                if shape.text_frame.paragraphs:
                    para = shape.text_frame.paragraphs[0]
                    print(f"  Paragraphs: {len(shape.text_frame.paragraphs)}")
                    
                    # Check runs
                    if para.runs:
                        run = para.runs[0]
                        print(f"  Runs: {len(para.runs)}")
                        print(f"  Font name: {run.font.name}")
                        print(f"  Font size: {run.font.size}")
                        if run.font.size:
                            print(f"  Font size (pt): {run.font.size.pt:.1f}")
                        print(f"  Bold: {run.font.bold}")
                        print(f"  Color: {run.font.color.rgb if run.font.color.rgb else 'None'}")
            else:
                # Not text, might be shape/image
                try:
                    w = shape.width / 914400
                    h = shape.height / 914400
                    t = shape.top / 914400
                    l = shape.left / 914400
                    print(f"  Position: ({l:.3f}\", {t:.3f}\")")
                    print(f"  Size: {w:.3f}\" × {h:.3f}\" ({w*72:.1f}pt × {h*72:.1f}pt)")
                except:
                    print(f"  (Could not get dimensions)")

if __name__ == "__main__":
    inspect_pptx("output/test_output.pptx")

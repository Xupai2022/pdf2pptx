#!/usr/bin/env python3
"""
Test the text decoration shape filtering fix.
"""

from pptx import Presentation
import sys

def test_pptx_no_unwanted_rectangles(pptx_path):
    """
    Test that the generated PPTX doesn't have unwanted text decoration rectangles.
    
    Specifically checks for:
    - Page 3: Should not have black (#13161A) rectangles
    - Page 4: Should not have black (#13161A) or green (#12A678) rectangles
    """
    prs = Presentation(pptx_path)
    
    print(f"Testing: {pptx_path}")
    print(f"Total slides: {len(prs.slides)}\n")
    
    # Test page 3 (slide index 2)
    print("=" * 80)
    print("Testing Slide 3 (Page 3)")
    print("=" * 80)
    
    slide_3 = prs.slides[2]
    black_rects_p3 = 0
    
    for shape in slide_3.shapes:
        try:
            if hasattr(shape, 'fill') and shape.fill.type == 1:  # SOLID
                color = shape.fill.fore_color.rgb
                color_hex = f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"
                
                # Check for black decoration color
                if color_hex in ['13161A', '000000']:
                    # Check if it's a shape without text
                    has_text = False
                    if hasattr(shape, 'text_frame'):
                        has_text = len(shape.text_frame.text.strip()) > 0
                    
                    if not has_text:
                        black_rects_p3 += 1
                        print(f"❌ Found unwanted black rectangle at "
                              f"({shape.left/914400:.1f}\", {shape.top/914400:.1f}\"), "
                              f"color=#{color_hex}")
        except:
            pass
    
    if black_rects_p3 == 0:
        print("✅ Page 3: No unwanted black rectangles found!")
    else:
        print(f"❌ Page 3: Found {black_rects_p3} unwanted black rectangle(s)")
    
    # Test page 4 (slide index 3)
    print("\n" + "=" * 80)
    print("Testing Slide 4 (Page 4)")
    print("=" * 80)
    
    slide_4 = prs.slides[3]
    black_rects_p4 = 0
    green_rects_p4 = 0
    
    for shape in slide_4.shapes:
        try:
            if hasattr(shape, 'fill') and shape.fill.type == 1:  # SOLID
                color = shape.fill.fore_color.rgb
                color_hex = f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"
                
                # Check for black or green decoration colors
                is_black = color_hex in ['13161A', '000000']
                is_green = color_hex in ['12A678']
                
                if is_black or is_green:
                    # Check if it's a shape without text
                    has_text = False
                    if hasattr(shape, 'text_frame'):
                        has_text = len(shape.text_frame.text.strip()) > 0
                    
                    if not has_text:
                        if is_black:
                            black_rects_p4 += 1
                            print(f"❌ Found unwanted black rectangle at "
                                  f"({shape.left/914400:.1f}\", {shape.top/914400:.1f}\"), "
                                  f"color=#{color_hex}")
                        elif is_green:
                            green_rects_p4 += 1
                            print(f"❌ Found unwanted green rectangle at "
                                  f"({shape.left/914400:.1f}\", {shape.top/914400:.1f}\"), "
                                  f"color=#{color_hex}")
        except:
            pass
    
    if black_rects_p4 == 0 and green_rects_p4 == 0:
        print("✅ Page 4: No unwanted black or green rectangles found!")
    else:
        print(f"❌ Page 4: Found {black_rects_p4} black and {green_rects_p4} green unwanted rectangle(s)")
    
    # Summary
    print("\n" + "=" * 80)
    print("TEST SUMMARY")
    print("=" * 80)
    
    total_issues = black_rects_p3 + black_rects_p4 + green_rects_p4
    
    if total_issues == 0:
        print("✅ ALL TESTS PASSED! No unwanted rectangles found.")
        return True
    else:
        print(f"❌ TESTS FAILED! Found {total_issues} unwanted rectangle(s) total.")
        return False

def test_text_preservation(pptx_path):
    """
    Test that text is preserved correctly (not accidentally filtered).
    """
    prs = Presentation(pptx_path)
    
    print("\n" + "=" * 80)
    print("Testing Text Preservation")
    print("=" * 80)
    
    # Check that slide 4 still has text
    slide_4 = prs.slides[3]
    text_shapes = 0
    
    for shape in slide_4.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
            text_shapes += 1
    
    print(f"Slide 4: Found {text_shapes} text shapes")
    
    if text_shapes >= 10:  # Should have many text elements
        print("✅ Text preservation: PASSED")
        return True
    else:
        print("❌ Text preservation: FAILED (too few text elements)")
        return False

if __name__ == "__main__":
    pptx_path = sys.argv[1] if len(sys.argv) > 1 else "output/安全运营月报_fixed.pptx"
    
    test1 = test_pptx_no_unwanted_rectangles(pptx_path)
    test2 = test_text_preservation(pptx_path)
    
    if test1 and test2:
        print("\n🎉 ALL TESTS PASSED!")
        sys.exit(0)
    else:
        print("\n❌ SOME TESTS FAILED!")
        sys.exit(1)

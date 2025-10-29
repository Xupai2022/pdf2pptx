"""
全面测试圆环修复功能
"""
import sys
import logging
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE

# Setup logging
logging.basicConfig(level=logging.INFO, format='%(message)s')
logger = logging.getLogger(__name__)

def test_ring_rendering(pptx_path):
    """测试圆环渲染"""
    logger.info(f"\n{'='*80}")
    logger.info("测试 1: 圆环形状渲染验证")
    logger.info(f"{'='*80}\n")
    
    prs = Presentation(pptx_path)
    
    # 检查问题页面（第5, 6, 7, 8, 9, 11页）
    problem_pages = {
        4: "第5页 - 应该有圆环",
        5: "第6页 - 百分比圆环(可能复杂)",
        6: "第7页 - 可能有图形",
        7: "第8页 - 应该有圆环",
        8: "第9页 - 应该有圆环",
        10: "第11页 - 应该有圆环"
    }
    
    total_rings = 0
    for page_idx, description in problem_pages.items():
        if page_idx >= len(prs.slides):
            continue
        
        slide = prs.slides[page_idx]
        rings = []
        
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'auto_shape_type') and shape.auto_shape_type == MSO_SHAPE.OVAL:
                    # 检查是否为圆环（白色填充+粗描边）
                    if shape.fill.type == 1:  # SOLID
                        rgb = shape.fill.fore_color.rgb
                        if rgb == (255, 255, 255) or (rgb[0] > 250 and rgb[1] > 250 and rgb[2] > 250):
                            stroke_width = shape.line.width.pt if shape.line.width else 0
                            if stroke_width > 20:  # 粗描边
                                rings.append(shape)
            except:
                pass
        
        status = "✅ 通过" if len(rings) > 0 else "⚠️ 未发现圆环"
        logger.info(f"{description}: {status} (找到 {len(rings)} 个圆环)")
        total_rings += len(rings)
    
    logger.info(f"\n总共找到 {total_rings} 个圆环")
    return total_rings >= 4  # 期望至少4个圆环（第5, 8, 9, 11页）

def test_normal_shapes_not_affected(pptx_path):
    """测试普通形状没有受影响"""
    logger.info(f"\n{'='*80}")
    logger.info("测试 2: 普通形状不受影响验证")
    logger.info(f"{'='*80}\n")
    
    prs = Presentation(pptx_path)
    
    # 检查所有页面
    rectangles = 0
    text_boxes = 0
    images = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'shape_type'):
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or shape.shape_type == 17:
                    text_boxes += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.shape_type == 13:
                    images += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape.shape_type == 1:
                    try:
                        if hasattr(shape, 'auto_shape_type') and shape.auto_shape_type == MSO_SHAPE.RECTANGLE:
                            rectangles += 1
                    except:
                        pass
    
    logger.info(f"文本框: {text_boxes}")
    logger.info(f"图片: {images}")
    logger.info(f"矩形: {rectangles}")
    
    # 验证有合理数量的普通形状（文本框应该多于100，矩形应该多于50）
    has_normal_shapes = (text_boxes > 100 and rectangles > 50)
    status = "✅ 通过" if has_normal_shapes else "❌ 失败"
    logger.info(f"\n普通形状保留状态: {status}")
    logger.info(f"  文本框足够: {'✅' if text_boxes > 100 else '❌'} ({text_boxes} > 100)")
    logger.info(f"  矩形足够: {'✅' if rectangles > 50 else '❌'} ({rectangles} > 50)")
    
    return has_normal_shapes

def test_no_rounded_rectangles_added(pptx_path):
    """测试没有错误地给矩形添加圆角"""
    logger.info(f"\n{'='*80}")
    logger.info("测试 3: 验证没有错误添加圆角")
    logger.info(f"{'='*80}\n")
    
    # 这个测试检查普通矩形文本框没有被错误地添加圆角
    # 由于 python-pptx 不直接暴露圆角属性，我们通过检查形状类型来间接验证
    prs = Presentation(pptx_path)
    
    unexpected_shapes = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            # 检查是否有意外的圆角矩形（ROUNDED_RECTANGLE）
            try:
                if hasattr(shape, 'auto_shape_type'):
                    if shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
                        unexpected_shapes += 1
            except:
                pass
    
    status = "✅ 通过" if unexpected_shapes == 0 else f"⚠️ 发现 {unexpected_shapes} 个圆角矩形"
    logger.info(f"圆角矩形数量: {unexpected_shapes}")
    logger.info(f"状态: {status}")
    
    return unexpected_shapes == 0

def main():
    """运行所有测试"""
    pptx_path = "output/glm-4.6-fixed2.pptx"
    
    if not Path(pptx_path).exists():
        logger.error(f"测试文件不存在: {pptx_path}")
        logger.error("请先运行: python main.py tests/glm-4.6.pdf output/glm-4.6-fixed2.pptx")
        return 1
    
    logger.info(f"\n{'#'*80}")
    logger.info("圆环修复功能全面测试")
    logger.info(f"{'#'*80}\n")
    logger.info(f"测试文件: {pptx_path}\n")
    
    # 运行测试
    test1 = test_ring_rendering(pptx_path)
    test2 = test_normal_shapes_not_affected(pptx_path)
    test3 = test_no_rounded_rectangles_added(pptx_path)
    
    # 总结
    logger.info(f"\n{'='*80}")
    logger.info("测试总结")
    logger.info(f"{'='*80}\n")
    
    tests = [
        ("圆环形状渲染", test1),
        ("普通形状不受影响", test2),
        ("没有错误添加圆角", test3)
    ]
    
    passed = sum(1 for _, result in tests if result)
    total = len(tests)
    
    for test_name, result in tests:
        status = "✅ 通过" if result else "❌ 失败"
        logger.info(f"{test_name}: {status}")
    
    logger.info(f"\n总计: {passed}/{total} 测试通过")
    
    if passed == total:
        logger.info("\n🎉 所有测试通过！")
        return 0
    else:
        logger.info(f"\n⚠️ {total - passed} 个测试失败")
        return 1

if __name__ == "__main__":
    sys.exit(main())

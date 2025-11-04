"""
Compare original PDF PNG with PPTX PNG to verify transparency fix
"""

import fitz  # PyMuPDF
from PIL import Image
from pptx import Presentation
import io
import numpy as np

def analyze_pdf_png(pdf_path, page_num, img_index):
    """Extract and analyze a specific PNG from PDF"""
    doc = fitz.open(pdf_path)
    page = doc[page_num]
    image_list = page.get_images(full=True)
    
    if img_index >= len(image_list):
        print(f"Image index {img_index} not found on page {page_num}")
        return None
    
    xref = image_list[img_index][0]
    base_image = doc.extract_image(xref)
    
    if not base_image:
        return None
    
    image_bytes = base_image["image"]
    pil_image = Image.open(io.BytesIO(image_bytes))
    
    doc.close()
    return pil_image

def analyze_pptx_png(pptx_path, slide_num, img_index):
    """Extract and analyze a specific PNG from PPTX"""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_num]
    
    images = [shape for shape in slide.shapes if shape.shape_type == 13]
    
    if img_index >= len(images):
        print(f"Image index {img_index} not found on slide {slide_num}")
        return None
    
    shape = images[img_index]
    image_blob = shape.image.blob
    pil_image = Image.open(io.BytesIO(image_blob))
    
    return pil_image

def compare_images(pdf_img, pptx_img, label):
    """Compare two images and report differences"""
    print(f"\n{'='*60}")
    print(f"{label}")
    print(f"{'='*60}")
    
    if pdf_img is None or pptx_img is None:
        print("Cannot compare - one or both images missing")
        return
    
    print(f"\nPDF Image:")
    print(f"  Size: {pdf_img.size}")
    print(f"  Mode: {pdf_img.mode}")
    
    print(f"\nPPTX Image:")
    print(f"  Size: {pptx_img.size}")
    print(f"  Mode: {pptx_img.mode}")
    
    # Sample corner pixels
    sample_points = [(0, 0), (pdf_img.width-1, 0), (0, pdf_img.height-1), (pdf_img.width-1, pdf_img.height-1)]
    
    print(f"\nPDF corner pixels:")
    for i, (x, y) in enumerate(sample_points):
        try:
            pixel = pdf_img.getpixel((x, y))
            if not isinstance(pixel, tuple):
                pixel = (pixel,)
            print(f"  Corner {i+1}: {pixel}")
        except:
            pass
    
    # For PPTX, we need to account for potential size difference
    pptx_sample_points = []
    if pptx_img.size == pdf_img.size:
        pptx_sample_points = sample_points
    else:
        # Sample corners relative to PPTX size
        pptx_sample_points = [(0, 0), (pptx_img.width-1, 0), (0, pptx_img.height-1), (pptx_img.width-1, pptx_img.height-1)]
    
    print(f"\nPPTX corner pixels:")
    for i, (x, y) in enumerate(pptx_sample_points):
        try:
            pixel = pptx_img.getpixel((x, y))
            if not isinstance(pixel, tuple):
                pixel = (pixel,)
            print(f"  Corner {i+1}: {pixel}")
        except:
            pass
    
    # Check transparency improvement
    pdf_has_alpha = pdf_img.mode in ['RGBA', 'LA', 'PA']
    pptx_has_alpha = pptx_img.mode in ['RGBA', 'LA', 'PA']
    
    print(f"\nAlpha channel:")
    print(f"  PDF: {pdf_has_alpha}")
    print(f"  PPTX: {pptx_has_alpha}")
    
    if not pdf_has_alpha and pptx_has_alpha:
        print(f"  ✅ FIXED: PPTX now has alpha channel!")
    elif pdf_has_alpha and not pptx_has_alpha:
        print(f"  ❌ REGRESSION: PPTX lost alpha channel!")
    
    # Count black pixels in PDF
    pdf_arr = np.array(pdf_img.convert('RGB'))
    pdf_black_pixels = np.sum(np.all(pdf_arr < 30, axis=-1))
    pdf_total_pixels = pdf_arr.shape[0] * pdf_arr.shape[1]
    pdf_black_pct = (pdf_black_pixels / pdf_total_pixels) * 100
    
    print(f"\nBlack pixel analysis:")
    print(f"  PDF: {pdf_black_pixels}/{pdf_total_pixels} ({pdf_black_pct:.2f}%)")
    
    # Count black pixels in PPTX (only RGB channels)
    pptx_arr = np.array(pptx_img.convert('RGB'))
    pptx_black_pixels = np.sum(np.all(pptx_arr < 30, axis=-1))
    pptx_total_pixels = pptx_arr.shape[0] * pptx_arr.shape[1]
    pptx_black_pct = (pptx_black_pixels / pptx_total_pixels) * 100
    
    print(f"  PPTX: {pptx_black_pixels}/{pptx_total_pixels} ({pptx_black_pct:.2f}%)")
    
    if pdf_black_pct > 20 and pptx_black_pct < pdf_black_pct:
        print(f"  ✅ IMPROVED: Reduced black pixels by {pdf_black_pct - pptx_black_pct:.2f}%")
    elif pdf_black_pct > 20 and pptx_black_pct >= pdf_black_pct:
        print(f"  ⚠️  WARNING: Black pixels not reduced (still {pptx_black_pct:.2f}%)")

# Test cases: Compare specific problematic images
print("="*60)
print("PNG TRANSPARENCY FIX VERIFICATION")
print("="*60)

# Page 3 (0-indexed: 2), Image 1 (0-indexed: 0) - Large PNG with black background
compare_images(
    analyze_pdf_png("tests/安全运营月报.pdf", 2, 0),
    analyze_pptx_png("output/test_fixed.pptx", 2, 0),
    "Page 3, Image 1 (246x252px icon)"
)

# Page 3 (0-indexed: 2), Image 3 (0-indexed: 2) - Large PNG chart
compare_images(
    analyze_pdf_png("tests/安全运营月报.pdf", 2, 2),
    analyze_pptx_png("output/test_fixed.pptx", 2, 2),
    "Page 3, Image 3 (652x505px chart)"
)

# Page 4 (0-indexed: 3), Image 2 (0-indexed: 1) - Icon
compare_images(
    analyze_pdf_png("tests/安全运营月报.pdf", 3, 1),
    analyze_pptx_png("output/test_fixed.pptx", 3, 5),
    "Page 4, Image 2 (108x108px icon)"
)

# Page 13 (0-indexed: 12), Image 3 (0-indexed: 2) - Large PNG
compare_images(
    analyze_pdf_png("tests/安全运营月报.pdf", 12, 2),
    analyze_pptx_png("output/test_fixed.pptx", 12, 2),
    "Page 13, Image 3 (448x287px)"
)

print("\n" + "="*60)
print("VERIFICATION COMPLETE")
print("="*60)

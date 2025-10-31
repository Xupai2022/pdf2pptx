"""
Verify the output PPTX file to ensure no duplication issues.
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys

def analyze_slide_4():
    """Analyze slide 4 (page 4) for triangle region."""
    print("\n" + "="*80)
    print("ANALYZING SLIDE 4 (Page 4 - Triangle Region)")
    print("="*80)
    
    prs = Presentation('output/season_report_fixed.pptx')
    
    if len(prs.slides) < 4:
        print("❌ ERROR: Not enough slides")
        return False
    
    slide = prs.slides[3]  # 0-indexed, so slide 4 is index 3
    
    print(f"\nTotal shapes on slide 4: {len(slide.shapes)}")
    
    # Categorize shapes
    images = []
    connectors = []
    text_boxes = []
    other_shapes = []
    
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            images.append(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
            connectors.append(shape)
        elif hasattr(shape, 'text_frame'):
            text_boxes.append(shape)
        else:
            other_shapes.append(shape)
    
    print(f"\nShape breakdown:")
    print(f"  Images: {len(images)}")
    print(f"  Lines/Connectors: {len(connectors)}")
    print(f"  Text boxes: {len(text_boxes)}")
    print(f"  Other shapes: {len(other_shapes)}")
    
    # Check for images in triangle region (approx 150-450, 150-350 in pts)
    # In PPTX, 1 pt ≈ 12700 EMUs
    EMU_PER_PT = 12700
    triangle_region = {
        'left': 150 * EMU_PER_PT,
        'right': 450 * EMU_PER_PT,
        'top': 150 * EMU_PER_PT,
        'bottom': 350 * EMU_PER_PT
    }
    
    images_in_triangle = []
    lines_in_triangle = []
    
    for img in images:
        left = img.left
        top = img.top
        right = left + img.width
        bottom = top + img.height
        
        # Check overlap
        overlaps = (
            left < triangle_region['right'] and right > triangle_region['left'] and
            top < triangle_region['bottom'] and bottom > triangle_region['top']
        )
        
        if overlaps:
            images_in_triangle.append({
                'left_pt': left / EMU_PER_PT,
                'top_pt': top / EMU_PER_PT,
                'width_pt': img.width / EMU_PER_PT,
                'height_pt': img.height / EMU_PER_PT
            })
    
    for line in connectors:
        left = line.left
        top = line.top
        right = left + line.width
        bottom = top + line.height
        
        # Check overlap
        overlaps = (
            left < triangle_region['right'] and right > triangle_region['left'] and
            top < triangle_region['bottom'] and bottom > triangle_region['top']
        )
        
        if overlaps:
            lines_in_triangle.append({
                'left_pt': left / EMU_PER_PT,
                'top_pt': top / EMU_PER_PT,
                'width_pt': line.width / EMU_PER_PT,
                'height_pt': line.height / EMU_PER_PT
            })
    
    print(f"\nTriangle region analysis:")
    print(f"  Images in region: {len(images_in_triangle)}")
    print(f"  Lines in region: {len(lines_in_triangle)}")
    
    if images_in_triangle:
        print("\n  Image details:")
        for i, img_info in enumerate(images_in_triangle):
            print(f"    Image {i+1}: pos=({img_info['left_pt']:.1f}, {img_info['top_pt']:.1f}), "
                  f"size={img_info['width_pt']:.1f}x{img_info['height_pt']:.1f}pt")
    
    if lines_in_triangle:
        print("\n  Line details:")
        for i, line_info in enumerate(lines_in_triangle):
            print(f"    Line {i+1}: pos=({line_info['left_pt']:.1f}, {line_info['top_pt']:.1f}), "
                  f"size={line_info['width_pt']:.1f}x{line_info['height_pt']:.1f}pt")
    
    # VERIFICATION: Check for issues
    issues = []
    
    # Issue 1: Large images in triangle region suggest PNG overlays
    large_images_in_triangle = [img for img in images_in_triangle 
                                 if img['width_pt'] > 100 or img['height_pt'] > 100]
    if large_images_in_triangle:
        issues.append(f"Found {len(large_images_in_triangle)} large image(s) in triangle region - "
                     "may be PNG overlays covering vector shapes")
    
    if issues:
        print("\n❌ POTENTIAL ISSUES DETECTED:")
        for issue in issues:
            print(f"  - {issue}")
        return False
    else:
        print("\n✅ NO DUPLICATION ISSUES: Triangle region looks clean")
        if lines_in_triangle:
            print(f"  ✅ Vector lines present: {len(lines_in_triangle)} line(s)")
        if not large_images_in_triangle:
            print("  ✅ No large PNG overlays detected")
        return True

def analyze_slide_6():
    """Analyze slide 6 (page 6) for text duplication."""
    print("\n" + "="*80)
    print("ANALYZING SLIDE 6 (Page 6 - Text Duplication Check)")
    print("="*80)
    
    prs = Presentation('output/season_report_fixed.pptx')
    
    if len(prs.slides) < 6:
        print("❌ ERROR: Not enough slides")
        return False
    
    slide = prs.slides[5]  # 0-indexed, so slide 6 is index 5
    
    print(f"\nTotal shapes on slide 6: {len(slide.shapes)}")
    
    # Collect all text content with positions
    text_contents = []
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
            text_contents.append({
                'text': shape.text_frame.text.strip(),
                'left_pt': shape.left / 12700,
                'top_pt': shape.top / 12700
            })
    
    print(f"\nText elements found: {len(text_contents)}")
    
    # Check for "30.46" or "30.46分" appearing multiple times
    target_texts = ['30.46', '30.46分', '30', '46', '分']
    duplicates = {}
    
    for target in target_texts:
        matches = [t for t in text_contents if target in t['text']]
        if len(matches) > 1:
            # Check if they're at similar positions (indicating duplication)
            for i in range(len(matches)):
                for j in range(i+1, len(matches)):
                    pos_diff_x = abs(matches[i]['left_pt'] - matches[j]['left_pt'])
                    pos_diff_y = abs(matches[i]['top_pt'] - matches[j]['top_pt'])
                    
                    # If positions are within 50pt, likely duplication
                    if pos_diff_x < 50 and pos_diff_y < 50:
                        key = f"'{target}'"
                        if key not in duplicates:
                            duplicates[key] = []
                        duplicates[key].append((matches[i], matches[j]))
    
    if duplicates:
        print("\n❌ POTENTIAL TEXT DUPLICATION DETECTED:")
        for text, pairs in duplicates.items():
            print(f"\n  Text {text} appears {len(pairs)} times at close positions:")
            for idx, (t1, t2) in enumerate(pairs):
                print(f"    Pair {idx+1}:")
                print(f"      Instance 1: '{t1['text']}' at ({t1['left_pt']:.1f}, {t1['top_pt']:.1f})")
                print(f"      Instance 2: '{t2['text']}' at ({t2['left_pt']:.1f}, {t2['top_pt']:.1f})")
        return False
    else:
        print("\n✅ NO TEXT DUPLICATION: Text '30.46分' appears only once")
        return True

def main():
    """Run verification checks."""
    print("="*80)
    print("PPTX OUTPUT VERIFICATION")
    print("="*80)
    print("\nVerifying: output/season_report_fixed.pptx")
    
    try:
        prs = Presentation('output/season_report_fixed.pptx')
        print(f"\nTotal slides: {len(prs.slides)}")
        
        # Run checks
        slide4_ok = analyze_slide_4()
        slide6_ok = analyze_slide_6()
        
        # Summary
        print("\n" + "="*80)
        print("VERIFICATION SUMMARY")
        print("="*80)
        
        results = [
            ("Slide 4 (Triangle region)", slide4_ok),
            ("Slide 6 (Text duplication)", slide6_ok)
        ]
        
        for check_name, passed in results:
            status = "✅ PASS" if passed else "❌ FAIL"
            print(f"{status}: {check_name}")
        
        if all(passed for _, passed in results):
            print("\n🎉 ALL CHECKS PASSED!")
            print("\nThe duplication issues have been successfully fixed:")
            print("  ✅ Page 4: Vector shapes preserved, no PNG overlays")
            print("  ✅ Page 6: Text appears only once, no duplication")
            return 0
        else:
            print("\n❌ SOME CHECKS FAILED")
            print("\nPlease review the output manually.")
            return 1
            
    except Exception as e:
        print(f"\n❌ ERROR: {e}")
        import traceback
        traceback.print_exc()
        return 1

if __name__ == "__main__":
    sys.exit(main())

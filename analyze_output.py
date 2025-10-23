#!/usr/bin/env python3
"""
Analyze the generated PPTX and compare with HTML reference
"""

import sys
from pathlib import Path
from pptx import Presentation
from collections import defaultdict

def analyze_pptx_details(pptx_path):
    """Analyze PPTX in detail."""
    print("\n" + "="*80)
    print("  Detailed PPTX Analysis")
    print("="*80 + "\n")
    
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    
    print(f"Slide dimensions: {prs.slide_width.inches:.1f}\" × {prs.slide_height.inches:.1f}\"")
    print(f"Total shapes: {len(slide.shapes)}")
    
    # Group text by font size
    text_by_size = defaultdict(list)
    images = []
    shapes = []
    
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            # Get font size from first run
            if shape.text_frame.paragraphs:
                para = shape.text_frame.paragraphs[0]
                if para.runs:
                    font_size = para.runs[0].font.size
                    if font_size:
                        size_pt = font_size.pt
                        text_by_size[size_pt].append({
                            'text': shape.text.strip()[:60],
                            'left': shape.left.inches,
                            'top': shape.top.inches,
                            'width': shape.width.inches,
                            'height': shape.height.inches
                        })
        elif hasattr(shape, "image"):
            images.append({
                'left': shape.left.inches,
                'top': shape.top.inches,
                'width': shape.width.inches,
                'height': shape.height.inches
            })
        else:
            # Regular shape
            shapes.append({
                'left': shape.left.inches,
                'top': shape.top.inches,
                'width': shape.width.inches,
                'height': shape.height.inches
            })
    
    print(f"\n{'─'*80}")
    print("Text Elements by Font Size:")
    print(f"{'─'*80}")
    
    for size in sorted(text_by_size.keys(), reverse=True):
        texts = text_by_size[size]
        print(f"\n  Font size {size}pt ({len(texts)} elements):")
        for i, item in enumerate(texts[:5]):
            print(f"    [{i+1}] \"{item['text']}\"")
            print(f"        Position: ({item['left']:.2f}\", {item['top']:.2f}\")")
            print(f"        Size: {item['width']:.2f}\" × {item['height']:.2f}\"")
        if len(texts) > 5:
            print(f"    ... and {len(texts)-5} more")
    
    print(f"\n{'─'*80}")
    print(f"Images: {len(images)}")
    print(f"{'─'*80}")
    for i, img in enumerate(images):
        print(f"  Image {i+1}:")
        print(f"    Position: ({img['left']:.2f}\", {img['top']:.2f}\")")
        print(f"    Size: {img['width']:.2f}\" × {img['height']:.2f}\"")
    
    print(f"\n{'─'*80}")
    print(f"Shapes/Backgrounds: {len(shapes)}")
    print(f"{'─'*80}")
    
    # Show shapes by size
    large_shapes = [s for s in shapes if s['width'] > 2 or s['height'] > 1]
    small_shapes = [s for s in shapes if s['width'] <= 2 and s['height'] <= 1]
    
    print(f"  Large shapes (backgrounds, cards): {len(large_shapes)}")
    for i, s in enumerate(large_shapes[:10]):
        print(f"    [{i+1}] {s['width']:.2f}\" × {s['height']:.2f}\" @ ({s['left']:.2f}\", {s['top']:.2f}\")")
    
    print(f"\n  Small shapes (borders, decorations): {len(small_shapes)}")


def compare_with_html():
    """Print expected structure from HTML."""
    print("\n" + "="*80)
    print("  Expected from HTML (for comparison)")
    print("="*80 + "\n")
    
    print("Expected Font Sizes:")
    print("  - 48px (≈36pt): Main title '风险资产深度剖析'")
    print("  - 36px (≈27pt): Subtitle '典型高风险资产案例展示'")
    print("  - 28px (≈21pt): Section headings")
    print("  - 25px (≈18.75pt): Regular text")
    print("  - 20px (≈15pt): Small labels")
    
    print("\nExpected Colors:")
    print("  - Blue #0a4275 / rgb(10, 66, 117): Headers, top bar")
    print("  - Red #dc2626: High risk items ('8个', '3个')")
    print("  - Gray #333, #666: Regular text")
    
    print("\nExpected Layout:")
    print("  1. Top bar: Full width, 10px height, blue")
    print("  2. Title section: Top left")
    print("  3. 3 stat cards with backgrounds")
    print("  4. 2 detail sections side by side")
    print("  5. Impact analysis section at bottom")
    print("  6. 6 icons (FontAwesome, should be images)")


def main():
    """Main function."""
    pptx_path = "output/test_output.pptx"
    
    if not Path(pptx_path).exists():
        print(f"❌ PPTX file not found: {pptx_path}")
        return 1
    
    compare_with_html()
    analyze_pptx_details(pptx_path)
    
    print("\n" + "="*80)
    print("  Analysis Complete")
    print("="*80)
    print("\n✅ Check if font sizes match (36pt, 27pt, 21pt, 18pt)")
    print("✅ Verify all 6 images (icons) are present")
    print("✅ Check if large shapes create card backgrounds")
    print("✅ Compare text positions with HTML structure")
    print()
    
    return 0


if __name__ == '__main__':
    sys.exit(main())

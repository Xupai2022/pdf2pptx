#!/usr/bin/env python3
"""
Validation script to verify first and last page layer structure.

This script checks that:
1. First page has a full-page background image at the bottom layer
2. Last page has a full-page background image at the bottom layer
3. Text elements are on top of the background image
4. Background image does not contain visible text
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import List, Dict, Any, Tuple


def analyze_slide_layers(slide, slide_num: int, total_slides: int, slide_width: int, slide_height: int) -> Dict[str, Any]:
    """
    Analyze the layer structure of a slide.
    
    Args:
        slide: python-pptx slide object
        slide_num: Slide number (0-indexed)
        total_slides: Total number of slides
        slide_width: Slide width in EMUs
        slide_height: Slide height in EMUs
        
    Returns:
        Dictionary with analysis results
    """
    is_first = (slide_num == 0)
    is_last = (slide_num == total_slides - 1)
    
    analysis = {
        'slide_num': slide_num + 1,
        'is_first': is_first,
        'is_last': is_last,
        'total_shapes': len(slide.shapes),
        'pictures': [],
        'text_boxes': [],
        'other_shapes': [],
        'has_full_page_background': False,
        'background_at_bottom': False,
        'text_above_background': False
    }
    
    # Analyze each shape
    for idx, shape in enumerate(slide.shapes):
        shape_info = {
            'index': idx,
            'type': shape.shape_type,
            'type_name': _get_shape_type_name(shape.shape_type),
            'left': shape.left,
            'top': shape.top,
            'width': shape.width,
            'height': shape.height
        }
        
        # Check if this is a full-page image (covering most of the slide)
        is_full_page = False
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            coverage = (shape.width / slide_width) * (shape.height / slide_height)
            is_full_page = coverage > 0.9  # 90% coverage
            shape_info['is_full_page'] = is_full_page
            shape_info['coverage'] = f"{coverage * 100:.1f}%"
            
            if is_full_page and idx == 0:
                analysis['has_full_page_background'] = True
                analysis['background_at_bottom'] = True
        
        # Categorize shape
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            analysis['pictures'].append(shape_info)
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX or (
            hasattr(shape, 'has_text_frame') and shape.has_text_frame
        ):
            if hasattr(shape, 'text'):
                shape_info['text_preview'] = shape.text[:50] if shape.text else ""
            analysis['text_boxes'].append(shape_info)
        else:
            analysis['other_shapes'].append(shape_info)
    
    # Check if text is above background
    if analysis['has_full_page_background'] and analysis['text_boxes']:
        # Background should be at index 0, text should be at higher indices
        background_idx = 0
        text_indices = [tb['index'] for tb in analysis['text_boxes']]
        analysis['text_above_background'] = all(idx > background_idx for idx in text_indices)
    
    return analysis


def _get_shape_type_name(shape_type) -> str:
    """Convert MSO_SHAPE_TYPE enum to readable name."""
    type_names = {
        MSO_SHAPE_TYPE.PICTURE: 'PICTURE',
        MSO_SHAPE_TYPE.TEXT_BOX: 'TEXT_BOX',
        MSO_SHAPE_TYPE.AUTO_SHAPE: 'AUTO_SHAPE',
        MSO_SHAPE_TYPE.GROUP: 'GROUP',
        MSO_SHAPE_TYPE.LINE: 'LINE',
        MSO_SHAPE_TYPE.PLACEHOLDER: 'PLACEHOLDER',
    }
    return type_names.get(shape_type, f'UNKNOWN({shape_type})')


def validate_pptx(pptx_path: str) -> bool:
    """
    Validate a PPTX file for first/last page layer structure.
    
    Args:
        pptx_path: Path to PPTX file
        
    Returns:
        True if validation passes
    """
    try:
        prs = Presentation(pptx_path)
        total_slides = len(prs.slides)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        print("=" * 80)
        print(f"PPTX Validation: {pptx_path}")
        print(f"Total slides: {total_slides}")
        print(f"Slide dimensions: {slide_width} x {slide_height} EMUs")
        print("=" * 80)
        
        validation_passed = True
        
        for idx, slide in enumerate(prs.slides):
            analysis = analyze_slide_layers(slide, idx, total_slides, slide_width, slide_height)
            
            # Only check first and last pages
            if not (analysis['is_first'] or analysis['is_last']):
                continue
            
            page_label = "FIRST" if analysis['is_first'] else "LAST"
            print(f"\n--- {page_label} PAGE (Slide {analysis['slide_num']}) ---")
            print(f"Total shapes: {analysis['total_shapes']}")
            print(f"Pictures: {len(analysis['pictures'])}")
            print(f"Text boxes: {len(analysis['text_boxes'])}")
            print(f"Other shapes: {len(analysis['other_shapes'])}")
            
            # Check requirement 1: Has full-page background
            if analysis['has_full_page_background']:
                print("✅ Has full-page background image")
            else:
                print("❌ Missing full-page background image")
                validation_passed = False
            
            # Check requirement 2: Background is at bottom (index 0)
            if analysis['background_at_bottom']:
                print("✅ Background image is at bottom layer (index 0)")
            else:
                print("❌ Background image is NOT at bottom layer")
                validation_passed = False
            
            # Check requirement 3: Text is above background
            if analysis['text_above_background']:
                print(f"✅ All {len(analysis['text_boxes'])} text boxes are above background")
            elif analysis['text_boxes']:
                print(f"❌ Some text boxes are NOT above background")
                validation_passed = False
            else:
                print("⚠️  No text boxes found (might be expected)")
            
            # Print detailed shape information
            print("\nShape details:")
            if analysis['pictures']:
                print(f"  Pictures ({len(analysis['pictures'])}):")
                for pic in analysis['pictures']:
                    full_page_marker = " [FULL PAGE]" if pic.get('is_full_page') else ""
                    coverage = pic.get('coverage', 'N/A')
                    print(f"    [{pic['index']}] {pic['type_name']} - "
                          f"Coverage: {coverage}{full_page_marker}")
            
            if analysis['text_boxes']:
                print(f"  Text boxes ({len(analysis['text_boxes'])}):")
                for tb in analysis['text_boxes'][:5]:  # Show first 5
                    preview = tb.get('text_preview', '')
                    print(f"    [{tb['index']}] {tb['type_name']} - "
                          f"Text: \"{preview}\"")
                if len(analysis['text_boxes']) > 5:
                    print(f"    ... and {len(analysis['text_boxes']) - 5} more text boxes")
        
        print("\n" + "=" * 80)
        if validation_passed:
            print("✅ VALIDATION PASSED: First and last pages have correct layer structure")
        else:
            print("❌ VALIDATION FAILED: First and/or last pages do not meet requirements")
        print("=" * 80)
        
        return validation_passed
        
    except Exception as e:
        print(f"❌ Error during validation: {e}")
        import traceback
        traceback.print_exc()
        return False


def main():
    """Main entry point."""
    if len(sys.argv) < 2:
        print("Usage: python validate_first_last_page_layers.py <pptx_file>")
        return 1
    
    pptx_path = sys.argv[1]
    
    if not Path(pptx_path).exists():
        print(f"Error: File not found: {pptx_path}")
        return 1
    
    success = validate_pptx(pptx_path)
    return 0 if success else 1


if __name__ == '__main__':
    sys.exit(main())

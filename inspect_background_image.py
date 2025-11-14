#!/usr/bin/env python3
"""
Inspect the background image to verify it doesn't contain visible text.

This script extracts the first page's background image and performs
visual analysis to confirm text has been masked out.
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io


def inspect_background_image(pptx_path: str, output_image_path: str = None):
    """
    Extract and inspect the background image from the first page.
    
    Args:
        pptx_path: Path to PPTX file
        output_image_path: Optional path to save the background image
    """
    try:
        prs = Presentation(pptx_path)
        
        # Get first slide
        slide = prs.slides[0]
        
        print(f"Inspecting first slide of: {pptx_path}")
        print(f"Total shapes: {len(slide.shapes)}")
        
        # Find the background image (should be first shape)
        background_img = None
        for idx, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                print(f"\nFound picture at index {idx}:")
                print(f"  Position: ({shape.left}, {shape.top})")
                print(f"  Size: {shape.width} x {shape.height}")
                
                # Get image data
                image = shape.image
                image_bytes = image.blob
                img = Image.open(io.BytesIO(image_bytes))
                
                print(f"  Image dimensions: {img.size}")
                print(f"  Image mode: {img.mode}")
                print(f"  Image format: {image.content_type}")
                
                # Check if this is the full-page background (index 0)
                if idx == 0:
                    background_img = img
                    print(f"  ✅ This is the BACKGROUND IMAGE")
                    
                    # Save if requested
                    if output_image_path:
                        img.save(output_image_path)
                        print(f"  💾 Saved to: {output_image_path}")
                    
                    # Analyze image to check if text was masked
                    analyze_text_masking(img)
                    
                    break
        
        if not background_img:
            print("\n❌ No background image found!")
            return False
        
        return True
        
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return False


def analyze_text_masking(img: Image.Image):
    """
    Analyze the image to verify text regions were masked.
    
    This performs basic statistical analysis to check if the image
    has uniform regions where text was masked out.
    """
    import numpy as np
    
    print("\n--- Text Masking Analysis ---")
    
    # Convert to numpy array
    img_array = np.array(img)
    
    # Get color statistics
    if img.mode == 'RGB':
        r, g, b = img_array[:,:,0], img_array[:,:,1], img_array[:,:,2]
        print(f"Red channel: min={r.min()}, max={r.max()}, mean={r.mean():.1f}")
        print(f"Green channel: min={g.min()}, max={g.max()}, mean={g.mean():.1f}")
        print(f"Blue channel: min={b.min()}, max={b.max()}, mean={b.mean():.1f}")
    elif img.mode == 'RGBA':
        r, g, b, a = img_array[:,:,0], img_array[:,:,1], img_array[:,:,2], img_array[:,:,3]
        print(f"Red channel: min={r.min()}, max={r.max()}, mean={r.mean():.1f}")
        print(f"Green channel: min={g.min()}, max={g.max()}, mean={g.mean():.1f}")
        print(f"Blue channel: min={b.min()}, max={b.max()}, mean={b.mean():.1f}")
        print(f"Alpha channel: min={a.min()}, max={a.max()}, mean={a.mean():.1f}")
    
    # Detect uniform regions (likely masked text areas)
    # A masked region should have pixels with very similar values
    if img.mode in ['RGB', 'RGBA']:
        # Calculate variance in small patches
        patch_size = 20
        h, w = img_array.shape[:2]
        
        uniform_patches = 0
        total_patches = 0
        
        for y in range(0, h - patch_size, patch_size):
            for x in range(0, w - patch_size, patch_size):
                patch = img_array[y:y+patch_size, x:x+patch_size]
                
                # Calculate standard deviation of the patch
                if img.mode == 'RGB':
                    patch_std = patch.std()
                else:  # RGBA
                    patch_std = patch[:,:,:3].std()  # Ignore alpha for this check
                
                total_patches += 1
                
                # If std is very low, it's likely a uniform (masked) region
                if patch_std < 10:  # Very uniform region
                    uniform_patches += 1
        
        uniform_ratio = uniform_patches / total_patches if total_patches > 0 else 0
        print(f"\nUniform patches: {uniform_patches}/{total_patches} ({uniform_ratio*100:.1f}%)")
        
        if uniform_ratio > 0.1:
            print("✅ Image contains significant uniform regions (likely masked text areas)")
        else:
            print("⚠️  Image has few uniform regions (text might not be fully masked)")


def main():
    """Main entry point."""
    if len(sys.argv) < 2:
        print("Usage: python inspect_background_image.py <pptx_file> [output_image.png]")
        return 1
    
    pptx_path = sys.argv[1]
    output_image_path = sys.argv[2] if len(sys.argv) > 2 else "background_image.png"
    
    if not Path(pptx_path).exists():
        print(f"Error: File not found: {pptx_path}")
        return 1
    
    success = inspect_background_image(pptx_path, output_image_path)
    return 0 if success else 1


if __name__ == '__main__':
    sys.exit(main())

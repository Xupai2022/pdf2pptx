#!/usr/bin/env python3
"""
Analyze pages 4 and 6 of season_report_del.pdf to understand the issues.
"""

import fitz  # PyMuPDF
from pptx import Presentation
from PIL import Image
import io
import json

def analyze_pdf_page(pdf_path, page_num):
    """Analyze a specific page of the PDF."""
    print(f"\n{'='*80}")
    print(f"PDF Analysis - Page {page_num + 1}")
    print(f"{'='*80}")
    
    doc = fitz.open(pdf_path)
    page = doc[page_num]
    
    # Get all drawings (shapes/lines)
    drawings = page.get_drawings()
    print(f"\n📐 Total drawings found: {len(drawings)}")
    
    # Analyze triangles and diagonal lines
    triangle_lines = []
    for i, drawing in enumerate(drawings):
        items = drawing.get("items", [])
        for item in items:
            if item[0] == "l":  # line
                if len(item) >= 3:  # has point data
                    p1 = item[1]  # start point
                    p2 = item[2]  # end point
                    x1, y1 = p1.x, p1.y
                    x2, y2 = p2.x, p2.y
                    
                    # Check if it's a diagonal line (not horizontal or vertical)
                    is_diagonal = abs(x2 - x1) > 5 and abs(y2 - y1) > 5
                    if is_diagonal:
                        angle = abs((y2 - y1) / (x2 - x1)) if (x2 - x1) != 0 else float('inf')
                        triangle_lines.append({
                            'index': i,
                            'coords': (x1, y1, x2, y2),
                            'length': ((x2-x1)**2 + (y2-y1)**2)**0.5,
                            'angle': angle,
                            'stroke_color': drawing.get('color'),
                            'stroke_width': drawing.get('width'),
                            'fill': drawing.get('fill')
                        })
    
    if triangle_lines:
        print(f"\n🔺 Diagonal lines (potential triangle sides): {len(triangle_lines)}")
        for line in triangle_lines:
            print(f"  Line: {line['coords']}")
            print(f"    Length: {line['length']:.2f}, Angle ratio: {line['angle']:.2f}")
            print(f"    Stroke: color={line['stroke_color']}, width={line['stroke_width']}, fill={line['fill']}")
    
    # Get all images
    image_list = page.get_image_info()
    print(f"\n🖼️  Total images found: {len(image_list)}")
    
    for img_idx, img_info in enumerate(image_list):
        xref = img_info.get("xref", 0)
        bbox = img_info.get("bbox", (0, 0, 0, 0))
        
        try:
            pix = fitz.Pixmap(doc, xref)
            
            print(f"\n  Image {img_idx + 1}:")
            print(f"    BBox: {bbox}")
            print(f"    Size: {pix.width}x{pix.height}")
            print(f"    Colorspace: {pix.colorspace.name if pix.colorspace else 'N/A'}")
            print(f"    Alpha: {pix.alpha}")
            
            # Save to analyze
            if pix.colorspace and pix.colorspace.name not in ["DeviceGray", "DeviceRGB"]:
                pix = fitz.Pixmap(fitz.csRGB, pix)
            
            img_bytes = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_bytes))
            
            # Analyze image content
            colors = img.getcolors(maxcolors=1000)
            if colors:
                print(f"    Unique colors: {len(colors)}")
                # Show dominant colors
                sorted_colors = sorted(colors, key=lambda x: x[0], reverse=True)[:5]
                print(f"    Top colors: {sorted_colors}")
        except Exception as e:
            print(f"\n  Image {img_idx + 1}: Error - {e}")
    
    # Get text blocks
    text_blocks = page.get_text("dict")["blocks"]
    text_only = [b for b in text_blocks if b.get("type") == 0]
    print(f"\n📝 Text blocks: {len(text_only)}")
    
    # Show text near top of page (might overlap with images)
    for block in text_only[:5]:
        bbox = block.get("bbox", [])
        text = "".join([span.get("text", "") for line in block.get("lines", []) for span in line.get("spans", [])])
        if text.strip():
            print(f"  Text: '{text.strip()[:50]}' at y={bbox[1]:.2f}")
    
    doc.close()


def analyze_pptx_page(pptx_path, slide_idx):
    """Analyze a specific slide of the PPTX."""
    print(f"\n{'='*80}")
    print(f"PPTX Analysis - Slide {slide_idx + 1}")
    print(f"{'='*80}")
    
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_idx]
    
    # Count different shape types
    lines = []
    images = []
    text_boxes = []
    other_shapes = []
    
    for shape in slide.shapes:
        if hasattr(shape, "line"):
            lines.append(shape)
        elif shape.shape_type == 13:  # Picture
            images.append(shape)
        elif hasattr(shape, "text_frame"):
            text_boxes.append(shape)
        else:
            other_shapes.append(shape)
    
    print(f"\n📊 Shape Summary:")
    print(f"  Lines: {len(lines)}")
    print(f"  Images: {len(images)}")
    print(f"  Text boxes: {len(text_boxes)}")
    print(f"  Other shapes: {len(other_shapes)}")
    
    # Analyze lines (looking for triangle sides)
    if lines:
        print(f"\n📐 Line details:")
        diagonal_lines = []
        for i, line in enumerate(lines):
            try:
                # Get line coordinates
                left = line.left
                top = line.top
                width = line.width
                height = line.height
                
                # Check if diagonal (both width and height > 0)
                is_diagonal = width > 100000 and height > 100000  # EMUs
                
                if is_diagonal:
                    diagonal_lines.append(line)
                    print(f"  Diagonal Line {len(diagonal_lines)}:")
                    print(f"    Position: ({left}, {top})")
                    print(f"    Size: {width} x {height} EMUs")
                    if hasattr(line.line, 'color'):
                        print(f"    Color: {line.line.color.rgb if hasattr(line.line.color, 'rgb') else 'N/A'}")
                    if hasattr(line.line, 'width'):
                        print(f"    Width: {line.line.width}")
            except Exception as e:
                print(f"  Line {i}: Error analyzing - {e}")
        
        print(f"\n  Total diagonal lines: {len(diagonal_lines)}")
    
    # Analyze images
    if images:
        print(f"\n🖼️  Image details:")
        for i, img in enumerate(images):
            print(f"  Image {i + 1}:")
            print(f"    Position: ({img.left}, {img.top})")
            print(f"    Size: {img.width} x {img.height} EMUs")
            
            # Try to get image data
            try:
                img_blob = img.image.blob
                pil_img = Image.open(io.BytesIO(img_blob))
                print(f"    PIL Size: {pil_img.size}")
                print(f"    PIL Mode: {pil_img.mode}")
                
                # Check quality indicators
                colors = pil_img.getcolors(maxcolors=10000)
                if colors:
                    print(f"    Unique colors: {len(colors)}")
            except Exception as e:
                print(f"    Error reading image: {e}")


def main():
    pdf_path = "tests/season_report_del.pdf"
    pptx_path = "output/season_report_test.pptx"
    
    print("="*80)
    print("Season Report Analysis - Pages 4 & 6")
    print("="*80)
    
    # Analyze Page 4 (index 3)
    print("\n" + "🔍 PAGE 4 ANALYSIS 🔍".center(80))
    analyze_pdf_page(pdf_path, 3)
    analyze_pptx_page(pptx_path, 3)
    
    # Analyze Page 6 (index 5)
    print("\n\n" + "🔍 PAGE 6 ANALYSIS 🔍".center(80))
    analyze_pdf_page(pdf_path, 5)
    analyze_pptx_page(pptx_path, 5)


if __name__ == "__main__":
    main()

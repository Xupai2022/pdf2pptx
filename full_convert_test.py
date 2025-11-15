"""
完整转换测试 - 转换整个season_report_del.pdf
"""
import sys
import logging
sys.path.append('.')

# 设置日志级别为INFO,不要DEBUG
logging.basicConfig(level=logging.INFO, format='%(levelname)s - %(message)s')

from main import convert_pdf_to_pptx, load_config

def main():
    pdf_path = "tests/season_report_del.pdf"
    output_path = "output/season_report_del_full.pptx"
    
    # 加载配置
    config = load_config()
    
    print(f"开始转换: {pdf_path} -> {output_path}")
    success = convert_pdf_to_pptx(pdf_path, output_path, config)
    
    if success:
        print(f"\n✅ 转换成功！")
        print(f"请打开 {output_path} 查看第16页表格")
        
        # 检查PPTX第16页
        print(f"\n检查第16页表格...")
        from pptx import Presentation
        prs = Presentation(output_path)
        
        if len(prs.slides) >= 16:
            slide = prs.slides[15]  # 0-indexed
            print(f"\n第16页内容:")
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    print(f"  表格: {len(table.rows)}行 x {len(table.columns)}列")
                    print(f"  行高详情:")
                    for row_idx, row in enumerate(table.rows):
                        height_pt = row.height / 914400 * 72
                        print(f"    行{row_idx+1}: {height_pt:.1f}pt ({row.height} EMUs)")
                    
                    print(f"\n  前7行内容:")
                    for row_idx in range(min(len(table.rows), 7)):
                        row = table.rows[row_idx]
                        texts = []
                        for cell in row.cells:
                            text = cell.text.strip()
                            if text:
                                texts.append(text[:20])
                            else:
                                texts.append('(空)')
                        print(f"    行{row_idx+1}: {' | '.join(texts)}")
        else:
            print(f"⚠️  PDF只有 {len(prs.slides)} 页")
    else:
        print(f"\n❌ 转换失败")
        return 1
    
    return 0

if __name__ == "__main__":
    sys.exit(main())

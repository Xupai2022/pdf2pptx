#!/usr/bin/env python3
"""直接检查边框shape"""

from pptx import Presentation

prs = Presentation("output/test_output.pptx")
slide = prs.slides[0]

# Check shape 5 (index 5)
print("Checking first potential border shape...")

for i, shape in enumerate(slide.shapes):
    try:
        if not hasattr(shape, 'text') or not shape.text or not shape.text.strip():
            if not hasattr(shape, 'image'):
                w_in = shape.width / 914400
                h_in = shape.height / 914400
                
                w_pt = w_in * 72
                h_pt = h_in * 72
                
                # Look for vertical borders (height > width, thin width)
                if h_pt > w_pt and w_pt < 50:
                    print(f"\nShape {i}: Vertical border candidate")
                    print(f"  Width: {w_pt:.2f}pt ({w_in:.4f}\")")
                    print(f"  Height: {h_pt:.2f}pt ({h_in:.4f}\")")
                    print(f"  Expected width: 4pt")
                    if abs(w_pt - 4) < 1:
                        print(f"  ✓ CORRECT!")
                    else:
                        print(f"  ✗ Off by {abs(w_pt - 4):.2f}pt")
    except:
        pass

print("\n" + "=" * 60)
print("All thin vertical shapes found above")

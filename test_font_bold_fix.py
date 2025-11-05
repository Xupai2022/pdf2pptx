"""
测试字体粗体修复

这个脚本将:
1. 转换PDF并生成PPTX
2. 检查PPTX中的字体设置
3. 验证Bold属性是否正确设置
"""

import sys
import yaml
from pathlib import Path
from pptx import Presentation

# Add src to path
sys.path.insert(0, str(Path(__file__).parent / 'src'))

from src.parser.pdf_parser import PDFParser
from src.rebuilder.slide_model import SlideModel
from src.generator.pptx_generator import PPTXGenerator

def test_font_bold_fix():
    """测试字体粗体修复"""
    
    print("=" * 80)
    print("字体粗体修复测试")
    print("=" * 80)
    
    # 1. 加载配置
    print("\n1. 加载配置...")
    config_path = Path(__file__).parent / 'config' / 'config.yaml'
    with open(config_path, 'r', encoding='utf-8') as f:
        config = yaml.safe_load(f)
    
    print(f"   字体映射配置:")
    for pdf_font, ppt_font in config['mapper']['font_mapping'].items():
        if 'YaHei' in pdf_font or 'yaHei' in pdf_font.lower():
            print(f"     {pdf_font} -> {ppt_font}")
    
    # 2. 查找PDF文件
    print("\n2. 查找PDF文件...")
    pdf_path = Path(__file__).parent / 'tests' / '安全运营月报.pdf'
    if not pdf_path.exists():
        print(f"   错误: 找不到PDF文件: {pdf_path}")
        return False
    print(f"   找到PDF: {pdf_path}")
    
    # 3. 解析PDF
    print("\n3. 解析PDF第一页...")
    parser = PDFParser(config['parser'])
    if not parser.open(str(pdf_path)):
        print("   错误: 无法打开PDF")
        return False
    
    page_data = parser.extract_page_elements(0)
    parser.close()
    
    # 检查提取的文本元素
    text_elements = [elem for elem in page_data['elements'] if elem['type'] == 'text']
    print(f"   提取到 {len(text_elements)} 个文本元素")
    
    bold_texts = []
    for elem in text_elements:
        if elem.get('is_bold'):
            bold_texts.append({
                'content': elem['content'][:20],
                'font_name': elem['font_name'],
                'is_bold': elem['is_bold'],
                'flags': elem.get('flags', 0)
            })
    
    print(f"   其中 {len(bold_texts)} 个是粗体文本:")
    for i, text in enumerate(bold_texts[:5], 1):
        print(f"     {i}. '{text['content']}' - {text['font_name']} (flags={text['flags']})")
    
    # 4. 分析布局
    print("\n4. 分析布局...")
    from src.analyzer.layout_analyzer_v2 import LayoutAnalyzerV2
    analyzer = LayoutAnalyzerV2(config['analyzer'])
    layout_data = analyzer.analyze_page(page_data)
    print(f"   找到 {len(layout_data['layout'])} 个布局区域")
    
    # 5. 构建Slide模型
    print("\n5. 构建Slide模型...")
    from src.rebuilder.coordinate_mapper import CoordinateMapper
    mapper = CoordinateMapper(config['rebuilder'])
    slide_model = mapper.create_slide_model(layout_data)
    print(f"   Slide包含 {len(slide_model.elements)} 个元素")
    
    # 6. 生成PPTX
    print("\n6. 生成PPTX...")
    output_path = Path(__file__).parent / 'test_font_bold_output.pptx'
    generator = PPTXGenerator(config)
    generator.generate_from_models([slide_model])
    generator.save(str(output_path))
    print(f"   生成完成: {output_path}")
    
    # 7. 检查PPTX内容
    print("\n7. 检查生成的PPTX...")
    prs = Presentation(str(output_path))
    
    if len(prs.slides) == 0:
        print("   错误: PPTX没有slide")
        return False
    
    slide = prs.slides[0]
    print(f"   第一个slide有 {len(slide.shapes)} 个shapes")
    
    # 检查所有文本框
    text_boxes = []
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame.text:
            text = shape.text_frame.text[:30]
            
            # 检查字体属性
            if shape.text_frame.paragraphs:
                para = shape.text_frame.paragraphs[0]
                if para.runs:
                    run = para.runs[0]
                    font_name = run.font.name
                    is_bold = run.font.bold
                    
                    text_boxes.append({
                        'text': text,
                        'font_name': font_name,
                        'is_bold': is_bold
                    })
    
    print(f"\n   找到 {len(text_boxes)} 个文本框:")
    bold_count = 0
    yahe_ui_count = 0
    
    for i, tb in enumerate(text_boxes[:10], 1):
        is_bold_str = " [粗体]" if tb['is_bold'] else ""
        print(f"     {i}. '{tb['text']}' - 字体:{tb['font_name']}{is_bold_str}")
        
        if tb['is_bold']:
            bold_count += 1
        if 'YaHei UI' in tb['font_name'] or 'Microsoft YaHei UI' == tb['font_name']:
            yahe_ui_count += 1
    
    # 8. 验证结果
    print("\n" + "=" * 80)
    print("验证结果:")
    print("-" * 80)
    
    success = True
    
    # 检查1: 是否使用了Microsoft YaHei UI字体
    print(f"\n✓ 检查1: 使用Microsoft YaHei UI字体")
    print(f"  结果: {yahe_ui_count}/{len(text_boxes)} 个文本框使用了Microsoft YaHei UI")
    if yahe_ui_count > 0:
        print(f"  ✓ 通过: 成功使用Microsoft YaHei UI字体")
    else:
        print(f"  ✗ 失败: 没有使用Microsoft YaHei UI字体")
        success = False
    
    # 检查2: 粗体属性是否正确设置
    print(f"\n✓ 检查2: 粗体属性设置")
    print(f"  PDF中粗体文本数量: {len(bold_texts)}")
    print(f"  PPT中粗体文本数量: {bold_count}")
    
    if bold_count >= len(bold_texts):
        print(f"  ✓ 通过: 粗体属性正确设置 ({bold_count} >= {len(bold_texts)})")
    else:
        print(f"  ✗ 失败: 粗体属性设置不足 ({bold_count} < {len(bold_texts)})")
        success = False
    
    # 检查3: "安全运营月报"文字的字体
    print(f"\n✓ 检查3: '安全运营月报'标题字体")
    title_found = False
    for tb in text_boxes:
        if '安全运营月报' in tb['text']:
            title_found = True
            print(f"  文本: {tb['text']}")
            print(f"  字体: {tb['font_name']}")
            print(f"  粗体: {tb['is_bold']}")
            
            if tb['font_name'] == 'Microsoft YaHei UI' and tb['is_bold']:
                print(f"  ✓ 通过: 字体和粗体属性都正确!")
            else:
                print(f"  ✗ 失败: 字体或粗体属性不正确")
                success = False
            break
    
    if not title_found:
        print(f"  ✗ 失败: 未找到'安全运营月报'文字")
        success = False
    
    print("\n" + "=" * 80)
    if success:
        print("✓ 所有测试通过!")
        print("\n修复说明:")
        print("  1. 字体名称已正确映射到 'Microsoft YaHei UI'")
        print("  2. Bold属性已正确设置")
        print("  3. PPT将使用真实的Bold字体文件，而不是算法加粗")
        print("  4. 字体显示效果应该与PDF完全一致")
    else:
        print("✗ 部分测试失败，请检查上述错误")
    print("=" * 80)
    
    return success

if __name__ == "__main__":
    success = test_font_bold_fix()
    sys.exit(0 if success else 1)

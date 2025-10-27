#!/usr/bin/env python3
"""
全面测试 PDF 转 PPTX 功能
测试边框、图案、特殊符号等各类元素的准确识别与转换
"""
import sys
import fitz
from pathlib import Path
from pptx import Presentation
sys.path.insert(0, str(Path(__file__).parent))

from main import convert_pdf_to_pptx, load_config, setup_logging

def test_pdf_elements(pdf_path, page_num=0):
    """测试 PDF 元素"""
    print(f"\n{'='*60}")
    print(f"测试文件: {pdf_path.name}")
    print(f"{'='*60}\n")
    
    doc = fitz.open(pdf_path)
    if page_num >= len(doc):
        print(f"⚠️  页面 {page_num + 1} 不存在")
        return None
    
    page = doc[page_num]
    paths = page.get_drawings()
    images = page.get_images(full=True)
    text_dict = page.get_text("dict")
    
    # 统计不同类型的元素
    stats = {
        'shapes_total': len(paths),
        'images': len(images),
        'text_blocks': sum(1 for b in text_dict.get("blocks", []) if b.get("type") == 0),
        'narrow_shapes': 0,
        'blue_shapes': 0,
        'complex_paths': 0,
        'special_chars': 0,
        'fontawesome_icons': 0
    }
    
    # 分析形状
    for path in paths:
        rect = path.get('rect', fitz.Rect())
        fill = path.get('fill')
        items = path.get('items', [])
        
        # 窄条形状（边框）
        if fill and (rect.width < 10 or rect.height < 10):
            stats['narrow_shapes'] += 1
        
        # 蓝色形状
        if fill:
            r, g, b = fill
            if r < 0.1 and b > 0.4:
                stats['blue_shapes'] += 1
        
        # 复杂路径（可能是图标）
        if len(items) > 10:
            stats['complex_paths'] += 1
    
    # 分析特殊字符
    for block in text_dict.get("blocks", []):
        if block.get("type") == 0:
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    text = span.get("text", "")
                    font = span.get("font", "")
                    
                    # FontAwesome 图标
                    if "FontAwesome" in font:
                        stats['fontawesome_icons'] += len(text)
                    
                    # 特殊 Unicode 字符
                    for char in text:
                        if ord(char) > 0xE000:
                            stats['special_chars'] += 1
    
    doc.close()
    
    # 打印统计
    print("PDF 元素统计:")
    print(f"  总形状: {stats['shapes_total']}")
    print(f"  图像: {stats['images']}")
    print(f"  文本块: {stats['text_blocks']}")
    print(f"  窄条形状（边框）: {stats['narrow_shapes']}")
    print(f"  蓝色形状: {stats['blue_shapes']}")
    print(f"  复杂路径（>10项）: {stats['complex_paths']}")
    print(f"  特殊字符: {stats['special_chars']}")
    print(f"  FontAwesome 图标: {stats['fontawesome_icons']}")
    
    return stats

def test_pptx_elements(pptx_path, page_num=0):
    """测试 PPTX 元素"""
    prs = Presentation(pptx_path)
    if page_num >= len(prs.slides):
        print(f"⚠️  页面 {page_num + 1} 不存在")
        return None
    
    slide = prs.slides[page_num]
    
    stats = {
        'total_shapes': len(slide.shapes),
        'text_shapes': 0,
        'image_shapes': 0,
        'other_shapes': 0,
        'narrow_shapes': 0,
        'blue_shapes': 0,
        'special_chars': 0,
        'fontawesome_chars': 0
    }
    
    for shape in slide.shapes:
        # 文本形状
        if hasattr(shape, 'text_frame') and shape.text_frame is not None:
            text = shape.text.strip()
            if text:
                stats['text_shapes'] += 1
                
                # 特殊字符
                for char in text:
                    code = ord(char)
                    if code > 0xE000:
                        stats['special_chars'] += 1
                        if 0xE000 <= code <= 0xF8FF:  # Private use area (FontAwesome)
                            stats['fontawesome_chars'] += 1
            else:
                stats['other_shapes'] += 1
        # 图像
        elif shape.shape_type == 13:
            stats['image_shapes'] += 1
        else:
            stats['other_shapes'] += 1
        
        # 窄条形状（边框）
        if hasattr(shape, 'width') and hasattr(shape, 'height'):
            width = shape.width / 914400
            height = shape.height / 914400
            
            if width < 0.15 or height < 0.15:
                stats['narrow_shapes'] += 1
            
            # 蓝色形状
            if hasattr(shape, 'fill'):
                fill = shape.fill
                if fill.type == 1:  # SOLID
                    try:
                        color = fill.fore_color
                        rgb_int = int(str(color.rgb), 16)
                        r = (rgb_int >> 16) & 0xFF
                        b = rgb_int & 0xFF
                        
                        if r < 50 and b > 100:
                            stats['blue_shapes'] += 1
                    except:
                        pass
    
    print("\nPPTX 元素统计:")
    print(f"  总形状: {stats['total_shapes']}")
    print(f"  文本形状: {stats['text_shapes']}")
    print(f"  图像: {stats['image_shapes']}")
    print(f"  其他形状: {stats['other_shapes']}")
    print(f"  窄条形状（边框）: {stats['narrow_shapes']}")
    print(f"  蓝色形状: {stats['blue_shapes']}")
    print(f"  特殊字符: {stats['special_chars']}")
    print(f"  FontAwesome 图标: {stats['fontawesome_chars']}")
    
    return stats

def compare_results(pdf_stats, pptx_stats):
    """比较转换结果"""
    print(f"\n{'='*60}")
    print("转换结果对比")
    print(f"{'='*60}\n")
    
    # 计算保留率
    checks = [
        ("窄条形状（边框）", pdf_stats['narrow_shapes'], pptx_stats['narrow_shapes']),
        ("蓝色形状", pdf_stats['blue_shapes'], pptx_stats['blue_shapes']),
        ("特殊字符", pdf_stats['special_chars'], pptx_stats['special_chars']),
        ("FontAwesome 图标", pdf_stats['fontawesome_icons'], pptx_stats['fontawesome_chars'])
    ]
    
    all_passed = True
    
    for name, pdf_count, pptx_count in checks:
        if pdf_count == 0:
            print(f"⚪ {name}: PDF中无此元素")
        elif pptx_count >= pdf_count * 0.9:  # 90% 保留率视为通过
            print(f"✅ {name}: {pdf_count} -> {pptx_count} ({pptx_count/pdf_count*100:.0f}%)")
        else:
            print(f"❌ {name}: {pdf_count} -> {pptx_count} ({pptx_count/pdf_count*100:.0f}% - 未达标)")
            all_passed = False
    
    return all_passed

def run_comprehensive_test():
    """运行全面测试"""
    setup_logging("INFO")
    config = load_config()
    
    test_cases = [
        ("tests/glm-4.6.pdf", 2, "output/glm46_comprehensive.pptx"),
        ("tests/test_sample.pdf", 0, "output/test_sample_comprehensive.pptx")
    ]
    
    results = []
    
    for pdf_path_str, page_num, output_path in test_cases:
        pdf_path = Path(pdf_path_str)
        if not pdf_path.exists():
            print(f"\n⚠️  文件不存在: {pdf_path}")
            continue
        
        # 测试 PDF
        pdf_stats = test_pdf_elements(pdf_path, page_num)
        if not pdf_stats:
            continue
        
        # 转换
        print(f"\n转换中...")
        success = convert_pdf_to_pptx(str(pdf_path), output_path, config)
        
        if not success:
            print(f"❌ 转换失败")
            results.append(False)
            continue
        
        # 测试 PPTX
        pptx_stats = test_pptx_elements(Path(output_path), page_num)
        if not pptx_stats:
            results.append(False)
            continue
        
        # 比较
        passed = compare_results(pdf_stats, pptx_stats)
        results.append(passed)
    
    # 总结
    print(f"\n{'='*60}")
    print("测试总结")
    print(f"{'='*60}\n")
    
    passed_count = sum(1 for r in results if r)
    total_count = len(results)
    
    print(f"测试通过: {passed_count}/{total_count}")
    
    if all(results):
        print("\n🎉 所有测试通过！")
        return 0
    else:
        print("\n⚠️  部分测试失败")
        return 1

if __name__ == "__main__":
    sys.exit(run_comprehensive_test())

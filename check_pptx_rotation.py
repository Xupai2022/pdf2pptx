#!/usr/bin/env python3
"""
检查生成的PPTX文件中的文字旋转
"""
from pptx import Presentation


def check_pptx_rotation(pptx_path):
    """检查PPTX中的文字旋转"""
    prs = Presentation(pptx_path)
    
    # 检查第11页 (索引10) 和第15页 (索引14)
    pages_to_check = [11, 15]
    
    for page_num in pages_to_check:
        if page_num - 1 >= len(prs.slides):
            print(f"页面 {page_num} 不存在")
            continue
        
        slide = prs.slides[page_num - 1]
        print(f"\n{'='*80}")
        print(f"第 {page_num} 页")
        print(f"{'='*80}\n")
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                
                # 关注特定文本
                if "10.64.5.37" in text or "10.74.145.44" in text or "未知业务" in text:
                    print(f"🔍 文本: '{text}'")
                    print(f"   位置: left={shape.left.inches:.3f}\", top={shape.top.inches:.3f}\"")
                    print(f"   尺寸: width={shape.width.inches:.3f}\", height={shape.height.inches:.3f}\"")
                    print(f"   旋转角度: {shape.rotation}°")
                    
                    # 分析旋转方向
                    # 归一化角度到-180到180范围
                    norm_rotation = shape.rotation
                    while norm_rotation > 180:
                        norm_rotation -= 360
                    while norm_rotation < -180:
                        norm_rotation += 360
                    
                    if shape.rotation == 45:
                        print(f"   ⚠️ PPT中45°是顺时针旋转 (从左上到右下 \\)")
                        print(f"   ⚠️ 但PDF中是 -45° (逆时针)，应该是从左下到右上 /")
                        print(f"   ⚠️ 问题：旋转方向反了！")
                    elif norm_rotation == -45 or shape.rotation == 315:
                        print(f"   ✅ PPT中-45° (或315°) 是逆时针旋转 (从左下到右上 /)")
                        print(f"   ✅ 这与PDF中的旋转方向一致")
                    elif shape.rotation == 0:
                        print(f"   ℹ️ 无旋转")
                    
                    print()


if __name__ == "__main__":
    import sys
    pptx_file = sys.argv[1] if len(sys.argv) > 1 else "/tmp/test_rotation.pptx"
    check_pptx_rotation(pptx_file)

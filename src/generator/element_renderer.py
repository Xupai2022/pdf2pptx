"""
Element Renderer - Renders individual slide elements to PowerPoint
"""

import logging
import io
from typing import Dict, Any, List
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from ..rebuilder.slide_model import SlideElement
from ..mapper.style_mapper import StyleMapper
from ..utils.xml_utils import set_run_font_xml, is_cjk_font

logger = logging.getLogger(__name__)


class ElementRenderer:
    """
    Renders slide elements to PowerPoint shapes.
    """
    
    def __init__(self, style_mapper: StyleMapper):
        """
        Initialize Element Renderer.
        
        Args:
            style_mapper: StyleMapper instance
        """
        self.style_mapper = style_mapper
    
    def render_text(self, slide, element: SlideElement) -> Any:
        """
        Render a text element to the slide.
        
        Args:
            slide: PowerPoint slide object
            element: SlideElement with text content
            
        Returns:
            Created shape object
        """
        position = element.position
        content = element.content
        style = element.style
        
        # Store original for comparison
        original_content = content
        
        # Clean content - remove ONLY control characters that cause XML issues
        # Map private use area characters and non-characters to their standard Unicode equivalents
        # These are often used in PDFs for symbol fonts but aren't XML-compatible
        if isinstance(content, str):
            import re
            
            # Map common private use area characters to standard Unicode
            # Based on common symbol font mappings (Wingdings, Webdings, Font Awesome, etc.)
            # Font Awesome uses U+F000-U+F2FF range for icon glyphs
            symbol_map = {
                # Font Awesome icons (U+F000-U+F2FF range)
                '\uf002': '🔍',  # search
                '\uf007': '👤',  # user
                '\uf013': '⚙',  # cog / settings
                '\uf015': '🏠',  # home
                '\uf019': '⬇',  # download
                '\uf01c': '📁',  # folder
                '\uf023': '🔒',  # lock
                '\uf024': '$',  # dollar
                '\uf044': '$',  # dollar alternate
                '\uf055': '➕',  # plus-circle
                '\uf058': '✓',  # check-circle
                '\uf059': '❓',  # question-circle  
                '\uf05a': 'ℹ',  # info-circle
                '\uf05e': '🚫',  # ban / prohibited
                '\uf06a': '●',  # circle / bullet
                '\uf06e': '👁',  # eye
                '\uf071': '⚠',  # exclamation-triangle / warning
                '\uf073': '📅',  # calendar
                '\uf084': '🔑',  # key
                '\uf0c9': '☰',  # bars / menu
                '\uf0f6': '➕',  # plus-square
                '\uf146': '➖',  # minus-square
                '\uf188': '🐛',  # bug
                '\uf3ed': '📹',  # video camera
                # Wingdings / Webdings
                '\uf0a7': '■',  # Black square
                '\uf0b7': '●',  # Bullet
                '\uf0fc': '☁',  # Cloud
                # Non-characters - remove (can't be rendered in PowerPoint XML)
                '\uffff': '',   # Non-character, typically used as placeholder
                '\ufffe': '',   # Non-character
            }
            
            # Apply symbol mappings
            for old_char, new_char in symbol_map.items():
                content = content.replace(old_char, new_char)
            
            # Remove NULL bytes and problematic control characters
            # Keep normal printable characters, whitespace, and CJK characters
            # Remove only C0 control characters (0x00-0x1F) except tab, newline, carriage return
            content = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F]', '', content)
            # Remove C1 control characters (0x80-0x9F)
            content = re.sub(r'[\x80-\x9F]', '', content)
            # Remove Unicode non-characters (U+FDD0 - U+FDEF, U+FFFE, U+FFFF) after mapping
            content = re.sub(r'[\uFDD0-\uFDEF\uFFFE\uFFFF]', '', content)
            
            # For remaining unmapped private use area characters (U+E000-U+F8FF),
            # replace with a generic placeholder symbol (●) so content isn't lost entirely
            # This handles custom icon fonts that we haven't mapped yet
            def replace_pua(match):
                # Replace each PUA character with a bullet point as a visible placeholder
                return '●' * len(match.group(0))
            
            content = re.sub(r'[\uE000-\uF8FF]+', replace_pua, content)
        
        if not content or not content.strip():
            if original_content and original_content.strip():
                logger.warning(f"Text cleaned away: original had {len(original_content)} chars: {repr(original_content[:50])}")
            return None
        
        # Create text box
        left = Inches(position['x'])
        top = Inches(position['y'])
        width = Inches(position['width'])
        height = Inches(position['height'])
        
        # ANTI-OVERLAP FIX: 只在非富文本模式下应用防重叠调整
        # 富文本模式下所有文本在同一个文本框，不需要防重叠调整
        # 而且需要保持完整宽度以防止文本换行
        is_rich_text = style.get('rich_text_runs') is not None

        if not is_rich_text:
            # 非富文本模式：应用防重叠调整
            anti_overlap_left_gap = 2.0 / 72.0  # 2pt shift right
            anti_overlap_width_reduction = 3.0 / 72.0  # 3pt reduction

            left += Inches(anti_overlap_left_gap)
            if width.inches > anti_overlap_width_reduction * 2:
                width -= Inches(anti_overlap_width_reduction)
        else:
            # 富文本模式：计算文本实际需要的宽度，确保不换行
            # 获取富文本runs信息
            runs = style.get('rich_text_runs', [])
            if runs:
                # 估算每个run的宽度
                estimated_width_pt = 0
                for run in runs:
                    text = run['text']
                    run_style = run['style']
                    font_size = run_style.get('font_size', 18)

                    # 估算宽度：
                    # - 中文字符：宽度 ≈ 字号 (1:1比例)
                    # - 英文/数字：宽度 ≈ 字号 × 0.6
                    # - 空格：宽度 ≈ 字号 × 0.3
                    for char in text:
                        if '\u4e00' <= char <= '\u9fff':  # 中文
                            estimated_width_pt += font_size
                        elif char == ' ':  # 空格
                            estimated_width_pt += font_size * 0.3
                        elif char == '\n':  # 换行符不占宽度
                            pass
                        else:  # 英文、数字、符号
                            estimated_width_pt += font_size * 0.6

                # 转换为英寸并增加20%缓冲
                estimated_width_inches = (estimated_width_pt / 72.0) * 1.2

                # 如果估算宽度大于当前宽度，扩展文本框
                if estimated_width_inches > width.inches:
                    logger.debug(f"Expanding textbox width: {width.inches:.3f}\" → {estimated_width_inches:.3f}\" for text: '{content[:30]}'")
                    width = Inches(estimated_width_inches)
        
        # Check if this is rotated text
        rotation = style.get('rotation', 0)
        is_rotated = abs(rotation) > 5
        
        if not is_rotated:
            # Only adjust dimensions for non-rotated text
            # For rotated text, the bbox from PDF is already correct and should not be modified
            
            # CRITICAL FIX: Trust PDF bbox coordinates to prevent text overlap
            # PDF provides accurate text bounding boxes. Artificially expanding them
            # causes overlaps, especially between Chinese text and numbers/letters.
            # 
            # Issue: Setting minimum widths for short texts caused "共发现弱口令" + "3" + "个"
            # to overlap because "3" bbox was expanded beyond its actual width.
            #
            # Solution: Use PDF bbox as-is, only add small padding for text rendering
            font_size = style.get('font_size', 18)
            
            # Add minimal padding (10%) only if width is extremely small to ensure text renders
            # But keep it proportional to avoid overlaps
            min_width_pt = font_size * 0.3  # Minimum 0.3x font size
            min_width_inches = min_width_pt / 72.0
            
            if width < Inches(min_width_inches):
                width = Inches(min_width_inches)
            
            # Ensure reasonable minimum height for text rendering
            # Use font size as basis to avoid arbitrary fixed minimums
            min_height_pt = font_size * 1.2  # 1.2x font size for line height
            min_height_inches = min_height_pt / 72.0
            
            if height < Inches(min_height_inches):
                height = Inches(min_height_inches)
        else:
            # For rotated text, keep bbox dimensions exactly as provided by PDF
            # The PDF bbox already accounts for rotation
            # Modifying width/height will shift the rotated text position incorrectly
            pass
        
        try:
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame

            # 检查是否有富文本runs（多个样式）
            rich_text_runs = style.get('rich_text_runs')

            if rich_text_runs:
                # 富文本模式：为每个run应用独立样式
                text_frame.clear()  # 清除默认段落

                paragraph = text_frame.paragraphs[0]
                for i, run_info in enumerate(rich_text_runs):
                    run_text = run_info['text']
                    run_style = run_info['style']

                    # 添加run
                    if i == 0:
                        # 第一个run使用现有paragraph
                        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                        run.text = run_text
                    else:
                        run = paragraph.add_run()
                        run.text = run_text

                    # 应用run的独立样式（字体、字号、颜色等）
                    font_name = run_style.get('font_name', 'Arial')
                    run.font.name = font_name
                    run.font.size = Pt(run_style.get('font_size', 18))
                    run.font.bold = run_style.get('bold', False)
                    run.font.italic = run_style.get('italic', False)

                    # 应用颜色
                    color = run_style.get('color', '#000000')
                    if color and color.startswith('#'):
                        try:
                            rgb = tuple(int(color[i:i+2], 16) for i in (1, 3, 5))
                            run.font.color.rgb = RGBColor(*rgb)
                        except:
                            pass

                    # 为中文字体设置东亚字体（确保中文正确显示）
                    if is_cjk_font(font_name):
                        try:
                            set_run_font_xml(run, font_name)
                        except:
                            pass

                logger.debug(f"Applied rich text with {len(rich_text_runs)} runs: '{content[:50]}'")
            else:
                # 普通文本模式：统一样式
                text_frame.text = content

                # Apply style
                self.style_mapper.apply_text_style(text_frame, style)

            # Text frame properties
            # 如果内容包含换行符（多行段落），启用word_wrap以支持换行显示
            # 否则禁用word_wrap以防止单行文本强制换行
            has_newlines = '\n' in content
            text_frame.word_wrap = has_newlines
            
            # Apply rotation if specified
            rotation = style.get('rotation', 0)
            if rotation != 0:
                # PDF和PPT的旋转系统:
                # - PDF中atan2(dy, dx)计算的角度: dy>0表示向下倾斜(\), dy<0表示向上倾斜(/)
                # - PPT中rotation: 正值为顺时针(\), 负值为逆时针(/)
                # - PDF的dir=(0.707, -0.707)表示/方向, atan2计算得-45°
                # - PPT需要-45°来显示/方向
                # 结论: 直接使用PDF角度，不需要取反
                textbox.rotation = rotation
                logger.debug(f"Applied rotation {rotation}° to text: {content[:30]}")
            
            return textbox
            
        except Exception as e:
            logger.error(f"Failed to render text element: {e}")
            logger.error(f"Content length: {len(content)}, first 50 chars: {repr(content[:50])}")
            # Try to identify problematic characters
            for i, c in enumerate(content):
                code = ord(c)
                if code < 32 and c not in '\t\n\r':
                    logger.error(f"  Found control char at pos {i}: U+{code:04X} ({repr(c)})")
                elif 0x80 <= code <= 0x9F:
                    logger.error(f"  Found C1 control at pos {i}: U+{code:04X} ({repr(c)})")
            return None
    
    def render_image(self, slide, element: SlideElement) -> Any:
        """
        Render an image element to the slide.
        
        Args:
            slide: PowerPoint slide object
            element: SlideElement with image content
            
        Returns:
            Created picture object
        """
        position = element.position
        image_data = element.content
        
        # Create image stream
        image_stream = io.BytesIO(image_data)
        
        # Position and size
        left = Inches(position['x'])
        top = Inches(position['y'])
        width = Inches(position['width'])
        height = Inches(position['height'])
        
        # Ensure reasonable dimensions (keep original size for small icons)
        if width < Inches(0.05):
            width = Inches(0.1)
        if height < Inches(0.05):
            height = Inches(0.1)
        
        try:
            picture = slide.shapes.add_picture(image_stream, left, top, width, height)
            return picture
            
        except Exception as e:
            logger.error(f"Failed to render image element: {e}")
            return None
    
    def render_shape(self, slide, element: SlideElement) -> Any:
        """
        Render a shape element to the slide.
        
        Args:
            slide: PowerPoint slide object
            element: SlideElement with shape content
            
        Returns:
            Created shape object
        """
        position = element.position
        shape_type = element.content
        style = element.style
        
        # Validate shape_type
        if not shape_type or not isinstance(shape_type, str):
            shape_type = 'rectangle'
        
        # Position and size
        left = Inches(position['x'])
        top = Inches(position['y'])
        width = Inches(position['width'])
        height = Inches(position['height'])
        
        # Filter out TRUE zero-dimension shapes (both width and height are zero)
        # BUT allow thin lines (one dimension is very small but not zero)
        # This is critical for rendering horizontal/vertical lines in tables and charts
        if width.inches == 0 and height.inches == 0:
            logger.warning(f"Skipping zero-dimension shape: w={width.inches:.4f}, h={height.inches:.4f}")
            return None
        
        # For shapes with one zero dimension, convert to thin line
        # E.g., horizontal line: width > 0, height = 0 -> height = 0.01
        # E.g., vertical line: width = 0, height > 0 -> width = 0.01
        min_line_thickness = Inches(0.01)  # 0.01 inches = 0.72 pt (very thin but visible)
        
        if width.inches == 0 and height.inches > 0:
            # Vertical line (zero width but has height)
            width = min_line_thickness
            logger.info(f"Adjusted zero-width vertical line to {min_line_thickness.inches:.4f} inches (height={height.inches:.4f})")
        elif height.inches == 0 and width.inches > 0:
            # Horizontal line (zero height but has width)
            height = min_line_thickness
            logger.info(f"Adjusted zero-height horizontal line to {min_line_thickness.inches:.4f} inches (width={width.inches:.4f})")
        
        # Ensure minimum dimensions (but allow thin borders)
        # Check if it's a border (thin vertical or horizontal shape)
        is_vertical_border = (position['height'] > position['width'] and position['width'] < 0.1)
        is_horizontal_border = (position['width'] > position['height'] and position['height'] < 0.1)
        is_border = is_vertical_border or is_horizontal_border
        
        if not is_border:
            # Apply minimum size only to non-border shapes
            if width < Inches(0.1):
                width = Inches(0.5)
            if height < Inches(0.1):
                height = Inches(0.5)
        
        try:
            # Special handling for lines - use connectors for proper line rendering
            if shape_type.lower() in ['line', 'triangle']:
                # For lines and triangles, we need to use connectors or stroke-only rectangles
                # Check if this has stroke but no fill
                has_stroke = style.get('stroke_color') is not None
                has_fill = style.get('fill_color') is not None
                
                if shape_type.lower() == 'line' and has_stroke and not has_fill:
                    # This is a true line - render as a connector
                    from pptx.enum.shapes import MSO_CONNECTOR
                    
                    # Calculate start and end points
                    # Use explicit x2, y2 if available (preserves line direction / vs \)
                    # Otherwise fall back to left+width, top+height
                    needs_vertical_flip = False
                    
                    if 'x2' in position and 'y2' in position:
                        # Use exact endpoints to preserve direction
                        x1, y1 = position['x'], position['y']
                        x2, y2 = position['x2'], position['y2']
                        
                        # Determine if line goes upward (/ direction)
                        # PowerPoint connectors normalize to top-left start, so we need to detect
                        # if the line should be flipped vertically to get / instead of \
                        if y2 < y1:
                            # Line goes upward - we need vertical flip
                            needs_vertical_flip = True
                            logger.debug(f"Line goes upward ({y1:.3f} -> {y2:.3f}), will apply vertical flip")
                        
                        begin_x = Inches(x1)
                        begin_y = Inches(y1)
                        end_x = Inches(x2)
                        end_y = Inches(y2)
                        logger.debug(f"Using explicit endpoints: ({x1:.3f},{y1:.3f}) -> ({x2:.3f},{y2:.3f})")
                    else:
                        # Fallback to bounding box (legacy behavior)
                        begin_x = left
                        begin_y = top
                        end_x = left + width
                        end_y = top + height
                        logger.debug(f"Using bounding box for line endpoints")
                    
                    # Add a straight connector (line)
                    # Note: PowerPoint will normalize this to left-top start point
                    connector = slide.shapes.add_connector(
                        MSO_CONNECTOR.STRAIGHT,
                        begin_x, begin_y, end_x, end_y
                    )
                    
                    # Apply vertical flip if needed to preserve / direction
                    if needs_vertical_flip:
                        try:
                            # Access the shape's XML element
                            sp_element = connector._element
                            # Look for spPr (shape properties) element
                            spPr = sp_element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
                            if spPr is None:
                                spPr = sp_element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}spPr')
                            
                            if spPr is not None:
                                # Look for xfrm (transform) element
                                xfrm = spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
                                if xfrm is not None:
                                    # Set flipV attribute to true (1)
                                    xfrm.set('flipV', '1')
                                    logger.debug(f"Applied vertical flip to connector via XML")
                                else:
                                    logger.warning("Could not find xfrm element for vertical flip")
                            else:
                                logger.warning("Could not find spPr element for vertical flip")
                        except Exception as e:
                            logger.warning(f"Failed to apply vertical flip: {e}")
                    
                    # Apply line style
                    if hasattr(connector, 'line'):
                        line = connector.line
                        # Set line color
                        if style.get('stroke_color'):
                            from pptx.util import Pt
                            from pptx.dml.color import RGBColor
                            hex_color = style['stroke_color']
                            if hex_color.startswith('#'):
                                rgb = tuple(int(hex_color[i:i+2], 16) for i in (1, 3, 5))
                                line.color.rgb = RGBColor(*rgb)
                        
                        # Set line width
                        stroke_width = style.get('stroke_width', 1)
                        line.width = Pt(stroke_width)
                    
                    logger.debug(f"Rendered line from ({begin_x}, {begin_y}) to ({end_x}, {end_y})")
                    return connector
            
            # Map shape type to MSO_SHAPE
            shape_map = {
                'rectangle': MSO_SHAPE.RECTANGLE,
                'rect': MSO_SHAPE.RECTANGLE,
                'circle': MSO_SHAPE.OVAL,
                'oval': MSO_SHAPE.OVAL,
                'ellipse': MSO_SHAPE.OVAL,
                'line': MSO_SHAPE.RECTANGLE,  # Fallback for lines with fill
                'triangle': MSO_SHAPE.RECTANGLE,  # Triangles as rectangles for now
                'star': MSO_SHAPE.STAR_5_POINT,  # 5-point star
                'path': MSO_SHAPE.RECTANGLE,
                'f': MSO_SHAPE.RECTANGLE,  # Fill path - default to rectangle
                's': MSO_SHAPE.RECTANGLE   # Stroke path - default to rectangle
            }
            
            # Get the MSO_SHAPE based on shape type
            # First check if it's explicitly marked as oval/circle
            if shape_type.lower() in ['oval', 'circle', 'ellipse']:
                # Always render oval/circle/ellipse as OVAL, regardless of aspect ratio
                mso_shape = MSO_SHAPE.OVAL
                logger.debug(f"Rendering {shape_type} as OVAL at ({position['x']:.1f}, {position['y']:.1f})")
            elif style.get('is_ring', False):
                # Ring shapes must be rendered as OVAL (circle) with stroke
                mso_shape = MSO_SHAPE.OVAL
                logger.debug(f"Rendering ring shape as OVAL at ({position['x']:.1f}, {position['y']:.1f})")
            else:
                # Use the shape_map for other types
                mso_shape = shape_map.get(shape_type.lower(), MSO_SHAPE.RECTANGLE)
            
            # Add shape
            shape = slide.shapes.add_shape(mso_shape, left, top, width, height)
            
            # Apply style
            self.style_mapper.apply_shape_style(shape, style)
            
            return shape
            
        except Exception as e:
            logger.error(f"Failed to render shape element: {e}")
            return None
    
    def render_table(self, slide, element: SlideElement) -> Any:
        """
        Render a table element to the slide using python-pptx table API.
        Supports cell merging for "设备类型" column to match reference format.
        
        Args:
            slide: PowerPoint slide object
            element: SlideElement with table content
            
        Returns:
            Created table shape object
        """
        position = element.position
        content = element.content  # This should be the table grid
        
        # Extract table grid from content
        if not isinstance(content, dict):
            logger.error(f"Table content must be a dict, got {type(content)}")
            return None
        

        
        grid = content.get('grid', [])
        if not grid:
            logger.warning("Table has no grid data")
            return None
        
        rows = len(grid)
        cols = max(len(row) for row in grid) if grid else 0
        
        if rows == 0 or cols == 0:
            logger.warning(f"Invalid table dimensions: {rows}x{cols}")
            return None
        
        logger.info(f"Rendering table: {rows}x{cols} at ({position['x']:.2f}, {position['y']:.2f})")
        
        # Position and size for the table
        left = Inches(position['x'])
        top = Inches(position['y'])
        width = Inches(position['width'])
        height = Inches(position['height'])
        
        try:
            # Add table shape to slide
            table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
            table = table_shape.table
            
            # CRITICAL: Apply actual column widths from PDF
            # python-pptx defaults to equal-width columns, but we need to match PDF exactly
            col_widths = content.get('col_widths', [])
            if col_widths and len(col_widths) == cols:
                logger.info(f"Applying actual column widths from PDF: {[f'{w:.1f}pt' for w in col_widths]}")
                
                # Convert points to EMUs (1 pt = 12700 EMUs)
                for col_idx, col_width_pt in enumerate(col_widths):
                    col_width_emus = int(col_width_pt * 12700)
                    table.columns[col_idx].width = col_width_emus
                    logger.debug(f"  Column {col_idx}: {col_width_pt:.1f}pt = {col_width_emus} EMUs")
            else:
                logger.warning(f"No column widths found in content, using default equal distribution")
            
            # CRITICAL: Use PDF row heights directly without scaling
            # User requirement: "不要缩放比例了"
            # Margins set to font_size/2 provide natural spacing
            # PDF row heights are trusted to be correct
            row_heights = content.get('row_heights', [])
            if row_heights and len(row_heights) == rows:
                logger.info(f"Using PDF row heights directly: {[f'{h:.1f}pt' for h in row_heights]}")
                
                for row_idx, pdf_row_height_pt in enumerate(row_heights):
                    # Use PDF height directly, no scaling
                    optimized_height_pt = pdf_row_height_pt
                    
                    # Apply safety bounds after optimization
                    # Min: 12pt (minimum readable row)
                    # Max: 200pt (prevent extreme outliers)
                    min_height_pt = 12.0
                    max_height_pt = 200.0
                    
                    final_row_height_pt = min(max(optimized_height_pt, min_height_pt), max_height_pt)
                    
                    # Convert to EMUs and apply
                    row_height_emus = int(final_row_height_pt * 12700)
                    table.rows[row_idx].height = row_height_emus
                    
                    logger.info(f"  Row {row_idx}: PDF={pdf_row_height_pt:.1f}pt → optimized={final_row_height_pt:.1f}pt ({row_height_emus} EMUs)")
            else:
                logger.info(f"No explicit row heights found, calculating from cell heights in grid")
                # Calculate row heights from grid data
                for row_idx, row_data in enumerate(grid):
                    # Find maximum cell height in this row
                    max_height_pt = 0
                    for cell_data in row_data:
                        cell_height = cell_data.get('height', 0)
                        if cell_height > max_height_pt:
                            max_height_pt = cell_height
                    
                    if max_height_pt > 0:
                        row_height_emus = int(max_height_pt * 12700)
                        table.rows[row_idx].height = row_height_emus
                        logger.debug(f"  Row {row_idx}: calculated {max_height_pt:.1f}pt = {row_height_emus} EMUs from cell data")
            
            # CRITICAL: Detect and apply cell merging for first column
            # In the reference table:
            # - Row 0: Header (no merge)
            # - Rows 1-6: "终端安全" should merge these 6 rows
            # - Rows 7-9: "态势感知" should merge these 3 rows
            # - Row 10: "防火墙" (single row, no merge)
            merge_map = self._detect_cell_merges(grid)
            
            # Populate table cells
            for row_idx, row_data in enumerate(grid):
                for col_idx, cell_data in enumerate(row_data):
                    if col_idx >= cols:
                        break
                        
                    cell = table.cell(row_idx, col_idx)
                    
                    # CRITICAL FIX: Smart line break detection in cell text
                    # Issue: Page 12 cells like "托管资产相关检查（托管资产是否被相应监测、响应组件完全覆盖）"
                    # should have line breaks between main text and parenthetical explanations
                    # Solution: Detect line breaks from text_elements Y positions
                    text_elements = cell_data.get('text_elements', [])
                    cell_text = cell_data.get('text', '')
                    
                    if text_elements and len(text_elements) > 1:
                        # Check if text elements have different Y positions (multiple lines)
                        # Group text elements by Y position with 3pt tolerance
                        y_tolerance = 3.0
                        y_groups = {}
                        
                        for text_elem in text_elements:
                            y = text_elem.get('y', 0)
                            y_key = round(y / y_tolerance) * y_tolerance
                            if y_key not in y_groups:
                                y_groups[y_key] = []
                            y_groups[y_key].append(text_elem)
                        
                        # If multiple Y groups, rebuild text with line breaks
                        if len(y_groups) > 1:
                            sorted_y_keys = sorted(y_groups.keys())
                            line_texts = []
                            
                            for y_key in sorted_y_keys:
                                # Sort text elements in this line by X position
                                line_elements = sorted(y_groups[y_key], key=lambda e: e.get('x', 0))
                                line_text = ''.join(e.get('content', e.get('text', '')) for e in line_elements)
                                line_texts.append(line_text)
                            
                            # Join lines with newline character
                            cell_text = '\n'.join(line_texts)
                            logger.debug(f"Detected {len(line_texts)} lines in cell, added line breaks: {line_texts}")
                    
                    # Set cell text
                    cell.text = cell_text
                    
                    # CRITICAL FIX: Use MINIMAL margins to allow row height matching PDF
                    # Issue: PowerPoint adds internal padding on top of our margin settings
                    # When we set margin = 3.76pt, PowerPoint renders ~5-6pt actual margin
                    # This prevents rows from shrinking to match PDF (21.5pt PDF vs 30pt+ PPT)
                    # 
                    # Solution: Set margins to 0.5pt (minimal)
                    # PowerPoint will add its own internal padding naturally
                    # This allows rows to shrink to PDF height while maintaining readability
                    # 
                    # Reference: PDF page 12 "托管服务器" cell analysis:
                    # - Cell height: 21.50pt (PDF)
                    # - With margin=3.76pt setting: PPT renders 30pt+ (TOO TALL)
                    # - With margin=0.5pt setting: PPT renders 21.5pt (MATCHES PDF)
                    from pptx.util import Pt as PtMargin
                    
                    # Use minimal margins from table_detector (default 0.5pt)
                    margin_top = cell_data.get('margin_top', 0.5)
                    margin_bottom = cell_data.get('margin_bottom', 0.5)
                    margin_left = cell_data.get('margin_left', 0.5)
                    margin_right = cell_data.get('margin_right', 0.5)
                    
                    # Apply minimal margins - DO NOT increase them
                    # PowerPoint's internal padding provides sufficient spacing
                    cell.margin_top = PtMargin(margin_top)
                    cell.margin_bottom = PtMargin(margin_bottom)
                    cell.margin_left = PtMargin(margin_left)
                    cell.margin_right = PtMargin(margin_right)
                    
                    logger.debug(f"Applied MINIMAL cell margins: T={margin_top:.1f}pt, B={margin_bottom:.1f}pt, "
                               f"L={margin_left:.1f}pt, R={margin_right:.1f}pt for cell '{cell_text[:15]}' "
                               f"(PowerPoint adds internal padding on top of these values)")
                    
                    # Apply cell background color
                    fill_color = cell_data.get('fill_color')
                    if fill_color:
                        from pptx.dml.color import RGBColor
                        from pptx.enum.dml import MSO_FILL
                        
                        cell.fill.solid()
                        if fill_color.startswith('#'):
                            rgb = tuple(int(fill_color[i:i+2], 16) for i in (1, 3, 5))
                            cell.fill.fore_color.rgb = RGBColor(*rgb)
                    
                    # Apply cell border (stroke) style using XML manipulation
                    # CRITICAL: Extract border properties from cell data
                    # For page 8 table: border should be RGB(222,227,237) #DEE3ED at 0.537pt
                    stroke_color = cell_data.get('stroke_color')
                    stroke_width = cell_data.get('stroke_width')
                    
                    # Set defaults if not provided
                    if stroke_color is None:
                        stroke_color = '#DEE3EC'  # Default gray border
                    if stroke_width is None:
                        stroke_width = 0.5  # Default to 0.5pt
                    
                    # Convert stroke width to EMUs (1 pt = 12700 EMUs)
                    # PDF page 8: 0.537pt = 6820 EMUs
                    border_width = int(stroke_width * 12700)
                    
                    from lxml import etree
                    tc = cell._tc
                    tcPr = tc.get_or_add_tcPr()
                    
                    # Parse stroke color
                    if stroke_color and stroke_color.startswith('#'):
                        rgb = tuple(int(stroke_color[i:i+2], 16) for i in (1, 3, 5))
                        rgb_hex = stroke_color[1:].upper()  # Remove # and uppercase
                        
                        # Create border elements for all four sides
                        # Using lnL (left), lnR (right), lnT (top), lnB (bottom)
                        for border_side in ['lnL', 'lnR', 'lnT', 'lnB']:
                            # Remove existing border if present
                            existing = tcPr.find(f'.//{{{tcPr.nsmap["a"]}}}{border_side}')
                            if existing is not None:
                                tcPr.remove(existing)
                            
                            # Create new border element
                            ln = etree.SubElement(tcPr, f'{{{tcPr.nsmap["a"]}}}{border_side}')
                            ln.set('w', str(border_width))
                            
                            # Add solid fill
                            solidFill = etree.SubElement(ln, f'{{{tcPr.nsmap["a"]}}}solidFill')
                            srgbClr = etree.SubElement(solidFill, f'{{{tcPr.nsmap["a"]}}}srgbClr')
                            srgbClr.set('val', rgb_hex)
                    
                    # Apply text formatting
                    # CRITICAL: Extract font properties from text_elements in cell
                    # For page 8 table:
                    # - Header row: FangSong, 7.5pt, RGB(112,122,137)
                    # - Data rows: FangSong, 7.5pt, RGB(20,22,26)
                    # 
                    # CRITICAL FIX: For multi-line cells, apply formatting to ALL paragraphs and runs
                    # Issue: Page 9 and 12 tables have cells with multiple lines (e.g., "托管资产相关检查\n（托管资产...）")
                    # When cell.text contains '\n', PowerPoint creates multiple paragraphs
                    # We must apply font formatting to EACH paragraph's runs, not just the first one
                    text_elements = cell_data.get('text_elements', [])
                    if text_elements and cell.text_frame.paragraphs:
                        from pptx.util import Pt as PtUtil
                        from pptx.dml.color import RGBColor as RGBColorUtil
                        
                        # Get first text element's style as representative for all text in cell
                        # (All text in a cell typically has the same font properties)
                        first_text = text_elements[0]
                        font_size = first_text.get('font_size')
                        font_color = first_text.get('color')
                        font_family = first_text.get('font_name')
                        is_bold_in_pdf = first_text.get('is_bold', False)
                        
                        # CRITICAL ENHANCEMENT: Apply text alignment detected from PDF
                        # Read alignment from cell_data (detected by table_detector._populate_table_cells)
                        text_alignment = cell_data.get('text_alignment', 'left')
                        
                        # Apply formatting to ALL paragraphs and runs in the cell
                        # This ensures multi-line text has consistent formatting
                        for para in cell.text_frame.paragraphs:
                            # Apply text alignment to paragraph
                            from pptx.enum.text import PP_ALIGN
                            if text_alignment == 'center':
                                para.alignment = PP_ALIGN.CENTER
                                logger.debug(f"Applied CENTER alignment to paragraph in cell '{cell.text[:20]}'")
                            elif text_alignment == 'right':
                                para.alignment = PP_ALIGN.RIGHT
                                logger.debug(f"Applied RIGHT alignment to paragraph in cell '{cell.text[:20]}'")
                            elif text_alignment == 'left':
                                para.alignment = PP_ALIGN.LEFT
                                logger.debug(f"Applied LEFT alignment to paragraph in cell '{cell.text[:20]}'")
                            
                            if not para.runs:
                                continue
                            
                            for run in para.runs:
                                # Apply font size
                                if font_size:
                                    run.font.size = PtUtil(font_size)

                                # Apply font color
                                if font_color and font_color.startswith('#') and len(font_color) == 7:
                                    rgb = tuple(int(font_color[i:i+2], 16) for i in (1, 3, 5))
                                    run.font.color.rgb = RGBColorUtil(*rgb)

                                # Apply font family
                                if font_family:
                                    run.font.name = font_family

                                # Apply bold property
                                run.font.bold = is_bold_in_pdf

                                # Set XML-level font properties for WPS compatibility
                                # This sets latin, ea (East Asian), and cs (Complex Script) font attributes
                                if font_family:
                                    is_cjk = is_cjk_font(font_family)
                                    set_run_font_xml(run, font_family, is_cjk=is_cjk)
                        
                        logger.debug(f"Applied formatting to {len(cell.text_frame.paragraphs)} paragraph(s) "
                                   f"in cell '{cell.text[:20]}': {font_family}, {font_size}pt, bold={is_bold_in_pdf}, "
                                   f"align={text_alignment}")
            
            # Apply cell merges after all cells are populated
            # Use the _detect_cell_merges method to identify merge regions from grid data
            from pptx.enum.text import MSO_ANCHOR
            
            # CRITICAL: Track which cells have been merged to avoid conflicts
            # Python-pptx can behave unexpectedly when trying to merge already-merged cells
            applied_merges = set()  # Track (row, col) of cells that are part of merges
            
            # CRITICAL: Sort merges by size (largest first) to avoid python-pptx conflicts
            # When merges overlap or are adjacent, applying larger merges first prevents issues
            def merge_size(m):
                rows = m['row_end'] - m['row_start'] + 1
                cols = m['col_end'] - m['col_start'] + 1
                return rows * cols
            
            sorted_merges = sorted(merge_map, key=merge_size, reverse=True)
            logger.info(f"Applying {len(sorted_merges)} merges (sorted by size, largest first)")
            
            for merge_info in sorted_merges:
                try:
                    row_start = merge_info['row_start']
                    row_end = merge_info['row_end']
                    col_start = merge_info['col_start']
                    col_end = merge_info['col_end']
                    
                    # CRITICAL: Check if any cell in this merge range is already part of another merge
                    # Python-pptx doesn't handle overlapping merges well
                    skip_merge = False
                    for r in range(row_start, row_end + 1):
                        for c in range(col_start, col_end + 1):
                            if (r, c) in applied_merges:
                                logger.warning(f"Skipping merge ({row_start},{col_start})-({row_end},{col_end}): "
                                             f"cell ({r},{c}) already merged")
                                skip_merge = True
                                break
                        if skip_merge:
                            break
                    
                    if skip_merge:
                        continue
                    
                    # Merge cells
                    merged_cell = table.cell(row_start, col_start)
                    end_cell = table.cell(row_end, col_end)
                    merged_cell.merge(end_cell)
                    
                    # Mark all cells in this range as merged
                    for r in range(row_start, row_end + 1):
                        for c in range(col_start, col_end + 1):
                            applied_merges.add((r, c))
                    
                    # Set vertical alignment to middle for merged cells
                    merged_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
                    logger.info(f"Merged cells: ({row_start},{col_start}) to ({row_end},{col_end}) and set vertical center")
                except Exception as e:
                    logger.error(f"Failed to merge cells: {e}")
            
            logger.info(f"Successfully rendered {rows}x{cols} table with {len(merge_map)} merges")
            return table_shape
            
        except Exception as e:
            logger.error(f"Failed to render table element: {e}", exc_info=True)
            return None
    
    
    def _detect_cell_merges(self, grid: List[List[Dict[str, Any]]]) -> List[Dict[str, Any]]:
        """
        Detect all cell merges in the table (both row and column merges).
        
        **INTELLIGENT MERGE DETECTION**:
        This method implements smart merge detection for complex table structures with parent-child relationships.
        
        Merge Logic:
        1. **Header row horizontal merges**: Merge consecutive empty cells in row 0
        2. **Parent-child vertical merges**: Detect parent items with child sub-items
           - If column N has text and column N+1 has text (sub-items), merge column N vertically
           - Count consecutive rows where column N+1 has text (these are sub-items)
           - Parent cell in column N merges downward to cover all sub-items
        3. **No-child horizontal merges**: If column N has text but column N+1 is empty (no sub-items),
           merge column N horizontally to column N+1
        4. **Simple vertical merges**: For remaining columns, merge consecutive empty cells below text cells
        
        Example (Page 12 table):
        - Row 1 col 1: "托管资产相关检查" has 2 sub-items in col 2 (rows 1-2) → merge col 1 rows 1-2
        - Row 3 col 1: "设备安全策略隐患" has NO sub-items → merge col 1-2 row 3
        - Row 6 col 1: "勒索风险" has 3 sub-items in col 2 (rows 6-8) → merge col 1 rows 6-8
        
        Args:
            grid: Table grid structure
            
        Returns:
            List of merge dictionaries with row_start, row_end, col_start, col_end
        """
        if not grid or len(grid) < 1:
            return []
        
        merge_map = []
        num_rows = len(grid)
        num_cols = len(grid[0]) if grid else 0
        

        
        # Track which cells have been merged to avoid duplicates
        merged_cells = set()
        
        # ========================================================================
        # STEP 1: Detect horizontal merges in header row (row 0)
        # ========================================================================
        row_idx = 0
        if row_idx < num_rows:
            col_idx = 0
            while col_idx < num_cols:
                cell_text = grid[row_idx][col_idx].get('text', '').strip()
                
                if cell_text and (row_idx, col_idx) not in merged_cells:
                    # Check for consecutive empty cells to the right
                    merge_start_col = col_idx
                    merge_end_col = col_idx
                    
                    next_col = col_idx + 1
                    while next_col < num_cols:
                        next_cell_text = grid[row_idx][next_col].get('text', '').strip()
                        if not next_cell_text:
                            merge_end_col = next_col
                            merged_cells.add((row_idx, next_col))
                            next_col += 1
                        else:
                            break
                    
                    # Create merge if span > 1 column
                    if merge_end_col > merge_start_col:
                        merge_map.append({
                            'row_start': row_idx,
                            'row_end': row_idx,
                            'col_start': merge_start_col,
                            'col_end': merge_end_col
                        })
                        logger.info(f"[HEADER] Horizontal merge: '{cell_text[:30]}' cols {merge_start_col}-{merge_end_col}")
                        merged_cells.add((row_idx, col_idx))
                
                col_idx += 1
        
        # ========================================================================
        # STEP 2: INTELLIGENT PARENT-CHILD DETECTION
        # For each column except the last, detect parent-child relationships
        # ========================================================================
        
        # Track which rows are part of a parent-child relationship
        # Key: row_idx, Value: True if this row is a child/sub-item row
        child_rows = set()
        
        for col_idx in range(num_cols - 1):  # Exclude last column (can't have children)
            row_idx = 1  # Start from row 1 (skip header)
            
            while row_idx < num_rows:
                if (row_idx, col_idx) in merged_cells:
                    row_idx += 1
                    continue
                
                cell_text = grid[row_idx][col_idx].get('text', '').strip()
                
                if cell_text:
                    # ============================================================
                    # CASE A: PARENT-CHILD RELATIONSHIP DETECTION
                    # Current column has text, check if next column has sub-items
                    # 
                    # CRITICAL: A sub-item is only valid if:
                    # 1. Next column (child column) has text
                    # 2. Current column (parent column) is EMPTY in subsequent rows (continuation)
                    # 
                    # If current column has text in subsequent rows, it's a NEW parent item,
                    # not a sub-item of the current parent.
                    # ============================================================
                    sub_item_count = 0
                    check_row = row_idx
                    next_col_idx = col_idx + 1
                    next_cell_text = grid[row_idx][next_col_idx].get('text', '').strip()
                    
                    if next_cell_text:
                        # Next column has text, check if there are multiple sub-items
                        while check_row < num_rows:
                            # Check if next column has text (is a sub-item)
                            sub_cell_text = grid[check_row][next_col_idx].get('text', '').strip()
                            current_col_text = grid[check_row][col_idx].get('text', '').strip()
                            
                            if sub_cell_text:
                                # Next column has text
                                if check_row == row_idx:
                                    # First row: this is the parent row itself, count it
                                    sub_item_count += 1
                                    check_row += 1
                                else:
                                    # Subsequent rows: only count if current column is EMPTY
                                    if not current_col_text:
                                        # Current column is empty, so next column text is a sub-item
                                        sub_item_count += 1
                                        check_row += 1
                                    else:
                                        # Current column has text, this is a NEW parent item
                                        # Stop counting sub-items here
                                        break
                            else:
                                # Next column is empty, no more sub-items
                                break
                    
                    # Only create vertical merge if there are MULTIPLE rows (sub_item_count > 1)
                    # A single row means no actual sub-items, just parent + child in same row
                    if sub_item_count > 1:
                        # Merge parent cell vertically across all sub-item rows
                        merge_end_row = row_idx + sub_item_count - 1
                        
                        # Mark all covered cells as merged
                        for r in range(row_idx + 1, merge_end_row + 1):
                            merged_cells.add((r, col_idx))
                            # CRITICAL: Mark these rows as child rows (contains sub-items)
                            # This prevents horizontal merging of sub-item cells in next column
                            child_rows.add(r)
                        
                        merge_map.append({
                            'row_start': row_idx,
                            'row_end': merge_end_row,
                            'col_start': col_idx,
                            'col_end': col_idx
                        })
                        logger.info(f"[PARENT-CHILD] Vertical merge col {col_idx}: '{cell_text[:30]}' rows {row_idx}-{merge_end_row} ({sub_item_count} sub-items)")
                        merged_cells.add((row_idx, col_idx))
                        
                        # Mark the parent row as well (for consistency)
                        child_rows.add(row_idx)
                        
                        # Also check if right-side columns should merge (e.g., "扣分说明" column)
                        # These columns typically span the same rows as the parent
                        for right_col in range(next_col_idx + 1, num_cols):
                            # Check if ALL rows in this range for right_col are empty EXCEPT the first
                            first_row_text = grid[row_idx][right_col].get('text', '').strip()
                            if first_row_text:
                                # Check if subsequent rows are empty
                                all_empty_below = all(
                                    not grid[r][right_col].get('text', '').strip()
                                    for r in range(row_idx + 1, merge_end_row + 1)
                                )
                                
                                if all_empty_below:
                                    # Merge this right column vertically too
                                    for r in range(row_idx + 1, merge_end_row + 1):
                                        merged_cells.add((r, right_col))
                                    
                                    merge_map.append({
                                        'row_start': row_idx,
                                        'row_end': merge_end_row,
                                        'col_start': right_col,
                                        'col_end': right_col
                                    })
                                    logger.info(f"[COMPANION] Vertical merge col {right_col}: '{first_row_text[:30]}' rows {row_idx}-{merge_end_row} (companion to parent)")
                                    merged_cells.add((row_idx, right_col))
                        
                        # Move to next parent item
                        row_idx = merge_end_row + 1
                        continue
                    
                    # ============================================================
                    # CASE B: NO CHILD RELATIONSHIP (no multi-row sub-items)
                    # This handles two scenarios:
                    # 1. Next column is empty → merge horizontally
                    # 2. Next column has text in same row only (single row) → no merge needed
                    # 
                    # CRITICAL FIX: For NO-CHILD case, we need to check if there are 
                    # subsequent rows with the SAME pattern (col_idx empty, next_col empty)
                    # If yes, we should merge ALL these rows together (vertical span)
                    # along with any companion columns (like "【扣分说明】【提分指南】")
                    # 
                    # CRITICAL CHECK: Skip horizontal merge if this is a sub-item row
                    # (i.e., part of a parent-child relationship from previous column)
                    # Sub-items like "托管服务器", "托管终端" should NOT merge horizontally
                    # ============================================================
                    if not next_cell_text and row_idx not in child_rows:
                        # Count how many consecutive rows have the same pattern
                        # (current col empty, next col empty)
                        consecutive_empty_rows = 1  # Start with current row
                        check_row = row_idx + 1
                        
                        while check_row < num_rows:
                            check_col_text = grid[check_row][col_idx].get('text', '').strip()
                            check_next_col_text = grid[check_row][next_col_idx].get('text', '').strip()
                            
                            # If both are empty, this row has same pattern
                            if not check_col_text and not check_next_col_text:
                                consecutive_empty_rows += 1
                                check_row += 1
                            else:
                                # Pattern changed, stop
                                break
                        
                        # If multiple consecutive rows, merge vertically
                        if consecutive_empty_rows > 1:
                            merge_end_row = row_idx + consecutive_empty_rows - 1
                            
                            # Merge current column (col_idx) vertically
                            for r in range(row_idx + 1, merge_end_row + 1):
                                merged_cells.add((r, col_idx))
                            
                            merge_map.append({
                                'row_start': row_idx,
                                'row_end': merge_end_row,
                                'col_start': col_idx,
                                'col_end': col_idx
                            })
                            logger.info(f"[NO-CHILD-MULTI] Vertical merge col {col_idx}: '{cell_text[:30]}' rows {row_idx}-{merge_end_row} (no sub-items, {consecutive_empty_rows} rows)")
                            merged_cells.add((row_idx, col_idx))
                            
                            # Also check companion columns (right-side columns)
                            # For NO-CHILD case, check if right columns should also merge
                            for right_col in range(next_col_idx + 1, num_cols):
                                first_row_text = grid[row_idx][right_col].get('text', '').strip()
                                if first_row_text:
                                    # Check if subsequent rows are empty
                                    all_empty_below = all(
                                        not grid[r][right_col].get('text', '').strip()
                                        for r in range(row_idx + 1, merge_end_row + 1)
                                    )
                                    
                                    if all_empty_below:
                                        # Merge this right column vertically too
                                        for r in range(row_idx + 1, merge_end_row + 1):
                                            merged_cells.add((r, right_col))
                                        
                                        merge_map.append({
                                            'row_start': row_idx,
                                            'row_end': merge_end_row,
                                            'col_start': right_col,
                                            'col_end': right_col
                                        })
                                        logger.info(f"[NO-CHILD-COMPANION] Vertical merge col {right_col}: '{first_row_text[:30]}' rows {row_idx}-{merge_end_row} (companion to no-child)")
                                        merged_cells.add((row_idx, right_col))
                            
                            # Move to next group
                            row_idx = merge_end_row + 1
                            continue
                        else:
                            # Single row, merge horizontally (original logic)
                            merge_map.append({
                                'row_start': row_idx,
                                'row_end': row_idx,
                                'col_start': col_idx,
                                'col_end': next_col_idx
                            })
                            logger.info(f"[NO-CHILD] Horizontal merge row {row_idx}: '{cell_text[:30]}' cols {col_idx}-{next_col_idx} (no sub-items)")
                            merged_cells.add((row_idx, col_idx))
                            merged_cells.add((row_idx, next_col_idx))
                
                row_idx += 1
        
        # ========================================================================
        # STEP 3: DETECT ZERO-WIDTH COLUMN MERGES (Page 9 specific case)
        # If a column has cells with zero width (missing) in multiple consecutive rows,
        # but has a tall cell in the first row, merge that tall cell vertically
        # Example: Page 9 col 0 row 0 has 154.8pt tall cell, rows 1-6 have zero-width cells
        # ========================================================================
        for col_idx in range(num_cols):
            row_idx = 1  # Start from row 1 (skip header)
            
            while row_idx < num_rows:
                if (row_idx, col_idx) in merged_cells:
                    row_idx += 1
                    continue
                
                cell_data = grid[row_idx][col_idx]
                cell_text = cell_data.get('text', '').strip()
                cell_width = cell_data.get('width', 0)
                cell_height = cell_data.get('height', 0)
                
                # CRITICAL: Check if this cell is effectively "missing" (width = 0)
                # AND if subsequent rows in the same column also have missing cells
                # This indicates a physical merge in the PDF that we need to replicate
                if cell_width == 0 and not cell_text and (row_idx, col_idx) not in merged_cells:
                    # Count consecutive zero-width cells in this column
                    consecutive_missing = 1
                    check_row = row_idx + 1
                    
                    while check_row < num_rows:
                        check_cell = grid[check_row][col_idx]
                        check_width = check_cell.get('width', 0)
                        check_text = check_cell.get('text', '').strip()
                        
                        if check_width == 0 and not check_text and (check_row, col_idx) not in merged_cells:
                            consecutive_missing += 1
                            check_row += 1
                        else:
                            break
                    
                    # If we have multiple consecutive missing cells (>=2),
                    # this is likely a physical merge in the PDF
                    # Find the cell above (or below) that should span these rows
                    if consecutive_missing >= 2:
                        # Check if previous row has a tall cell that should span these rows
                        prev_row = row_idx - 1
                        if prev_row >= 0:
                            prev_cell = grid[prev_row][col_idx]
                            prev_height = prev_cell.get('height', 0)
                            prev_text = prev_cell.get('text', '').strip()
                            
                            # If previous row has a tall cell (>50pt), merge it with missing cells
                            # Calculate expected height for this merged region
                            # If prev cell height roughly matches the span, it's a physical merge
                            expected_height_min = consecutive_missing * 15  # Conservative estimate (min 15pt per row)
                            expected_height_max = consecutive_missing * 30  # Liberal estimate (max 30pt per row)
                            
                            if prev_height >= expected_height_min and prev_height <= expected_height_max * 1.5:
                                # This is a physical merge! Merge prev_row with current rows
                                merge_end_row = row_idx + consecutive_missing - 1
                                
                                # Mark all cells as merged
                                for r in range(row_idx, merge_end_row + 1):
                                    merged_cells.add((r, col_idx))
                                
                                merge_map.append({
                                    'row_start': prev_row,
                                    'row_end': merge_end_row,
                                    'col_start': col_idx,
                                    'col_end': col_idx
                                })
                                logger.info(f"[PHYSICAL-MERGE] Detected physical merge col {col_idx}: '{prev_text[:30]}' rows {prev_row}-{merge_end_row} (prev_height={prev_height:.1f}pt, missing_rows={consecutive_missing})")
                                merged_cells.add((prev_row, col_idx))
                                
                                # Move past merged region
                                row_idx = merge_end_row + 1
                                continue
                    
                    # If not a physical merge, just skip this missing cell
                    row_idx += 1
                    continue
                
                # If cell has text, check for empty cells below (simple case)
                if cell_text:
                    merge_start_row = row_idx
                    merge_end_row = row_idx
                    
                    # Look ahead for empty cells in the same column
                    next_row = row_idx + 1
                    consecutive_empty_rows = 0
                    
                    while next_row < num_rows:
                        next_cell_text = grid[next_row][col_idx].get('text', '').strip()
                        
                        # If next cell is empty, count it
                        if not next_cell_text and (next_row, col_idx) not in merged_cells:
                            consecutive_empty_rows += 1
                            merge_end_row = next_row
                            next_row += 1
                        else:
                            # Next cell has text or is already merged, stop
                            break
                    
                    # CRITICAL FIX: Don't merge if empty rows are likely "reserved data rows"
                    #判断标准:
                    # 1. 当前单元格在表头行或第2行 (row_idx <= 1)
                    # 2. 下方有大量连续空行 (>= 3行)
                    # 3. 这些空行占据剩余行数的大部分 (>= 50%)
                    # 4. 所有列在这些行中都是空的(表明整行都是空的)
                    # 
                    # 这种情况下，空行很可能是预留的数据行，应该保留，不要合并
                    # 例如: season_report_del.pdf 第16页，表头下方有5个空行用于填写数据
                    should_skip_merge = False
                    
                    if row_idx <= 1 and consecutive_empty_rows >= 3:
                        remaining_rows = num_rows - row_idx - 1
                        empty_row_ratio = consecutive_empty_rows / remaining_rows if remaining_rows > 0 else 0
                        
                        if empty_row_ratio >= 0.5:
                            # 检查这些空行在所有列中是否都是空的
                            all_columns_empty = True
                            for check_row in range(row_idx + 1, merge_end_row + 1):
                                for check_col in range(num_cols):
                                    check_text = grid[check_row][check_col].get('text', '').strip()
                                    if check_text:
                                        all_columns_empty = False
                                        break
                                if not all_columns_empty:
                                    break
                            
                            if all_columns_empty:
                                should_skip_merge = True
                                logger.info(f"[SKIP-MERGE] Preserving {consecutive_empty_rows} empty data rows "
                                          f"after header row {row_idx} col {col_idx} ('{cell_text[:30]}') - "
                                          f"likely reserved for data entry ({empty_row_ratio*100:.0f}% of remaining rows)")
                    
                    # Create merge only if not skipped and span > 1 row
                    if not should_skip_merge and merge_end_row > merge_start_row:
                        # Mark cells as merged
                        for r in range(row_idx + 1, merge_end_row + 1):
                            merged_cells.add((r, col_idx))
                        
                        merge_map.append({
                            'row_start': merge_start_row,
                            'row_end': merge_end_row,
                            'col_start': col_idx,
                            'col_end': col_idx
                        })
                        logger.info(f"[SIMPLE] Vertical merge col {col_idx}: '{cell_text[:30]}' rows {merge_start_row}-{merge_end_row}")
                        merged_cells.add((row_idx, col_idx))
                    elif should_skip_merge:
                        # Skip merge, just move to next row
                        row_idx += 1
                        continue
                    
                    # Move to next group
                    row_idx = merge_end_row + 1
                else:
                    # Empty cell, skip
                    row_idx += 1
        
        # ========================================================================
        # STEP 4: SIMPLE VERTICAL MERGES for remaining empty cells
        # Process any remaining unmerged cells
        # ========================================================================
        
        logger.info(f"Total detected merges: {len(merge_map)}")
        return merge_map

    def render_element(self, slide, element: SlideElement) -> Any:
        """
        Render any element type to the slide.
        
        Args:
            slide: PowerPoint slide object
            element: SlideElement to render
            
        Returns:
            Created shape/object
        """
        if element.type == 'text':
            return self.render_text(slide, element)
        elif element.type == 'image':
            return self.render_image(slide, element)
        elif element.type == 'shape':
            return self.render_shape(slide, element)
        elif element.type == 'table':
            return self.render_table(slide, element)
        else:
            logger.warning(f"Unknown element type: {element.type}")
            return None

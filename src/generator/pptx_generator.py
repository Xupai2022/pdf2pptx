"""
PPTX Generator - Main generator for creating PowerPoint presentations
"""

import logging
from typing import List, Dict, Any, Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from ..rebuilder.slide_model import SlideModel
from ..mapper.style_mapper import StyleMapper
from .element_renderer import ElementRenderer

logger = logging.getLogger(__name__)


class PPTXGenerator:
    """
    Generates PowerPoint presentations from slide models.
    """
    
    def __init__(self, config: Dict[str, Any]):
        """
        Initialize PPTX Generator.
        
        Args:
            config: Configuration dictionary
        """
        self.config = config
        # Support both full config and generator-only config
        generator_config = config.get('generator', config)
        self.template_path = generator_config.get('template')
        self.preserve_layout = generator_config.get('preserve_layout', True)
        self.compress_images = generator_config.get('compress_images', False)
        
        # Initialize presentation
        if self.template_path and Path(self.template_path).exists():
            self.prs = Presentation(self.template_path)
        else:
            self.prs = Presentation()
        
        # CRITICAL: Slide dimensions will be set dynamically based on each PDF's actual size
        # This happens in add_slide_from_model() when the first slide is added
        # For now, store config defaults but don't apply them yet
        rebuilder_config = config.get('rebuilder', {})
        self.default_slide_width = rebuilder_config.get('slide_width', 10.0)
        self.default_slide_height = rebuilder_config.get('slide_height', 7.5)
        self.slide_dimensions_set = False  # Track if dimensions have been set
        logger.info(f"PPTXGenerator initialized with default dimensions {self.default_slide_width}\" x {self.default_slide_height}\"")
        
        # Initialize style mapper and element renderer
        self.style_mapper = StyleMapper(config)
        self.element_renderer = ElementRenderer(self.style_mapper)
    
    def add_slide_from_model(self, slide_model: SlideModel) -> Any:
        """
        Add a slide to the presentation from a slide model.

        Args:
            slide_model: SlideModel to render

        Returns:
            Created slide object
        """
        # CRITICAL: Set slide dimensions from first slide model (all slides have same PDF dimensions)
        if not self.slide_dimensions_set:
            self.prs.slide_width = Inches(slide_model.width)
            self.prs.slide_height = Inches(slide_model.height)
            self.slide_dimensions_set = True
            logger.info(f"Set PPT dimensions to {slide_model.width:.2f}\" x {slide_model.height:.2f}\" (matching PDF size)")
        
        # Use blank slide layout
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_layout)

        logger.info(f"Creating slide {slide_model.slide_number + 1} with {len(slide_model.elements)} elements")

        # Update style mapper with slide-specific scale factor if available
        if hasattr(slide_model, 'scale_factor') and slide_model.scale_factor is not None:
            self.style_mapper.update_font_scale(slide_model.scale_factor)
            logger.info(f"Applied slide-specific scale factor: {slide_model.scale_factor:.3f}")

        # Set background if specified
        if slide_model.background_color:
            self._set_slide_background_color(slide, slide_model.background_color)

        # Render elements in z-index order
        for element in slide_model.elements:
            try:
                self.element_renderer.render_element(slide, element)
            except Exception as e:
                logger.error(f"Failed to render element: {e}")

        return slide
    
    def generate_from_models(self, slide_models: List[SlideModel]) -> bool:
        """
        Generate presentation from list of slide models.
        
        Args:
            slide_models: List of SlideModel objects
            
        Returns:
            True if successful
        """
        try:
            for slide_model in slide_models:
                self.add_slide_from_model(slide_model)
            
            logger.info(f"Generated presentation with {len(slide_models)} slides")
            return True
            
        except Exception as e:
            logger.error(f"Failed to generate presentation: {e}")
            return False
    
    def save(self, output_path: str) -> bool:
        """
        Save the presentation to file.
        
        Args:
            output_path: Path to save the PPTX file
            
        Returns:
            True if successful
        """
        try:
            # Ensure output directory exists
            output_file = Path(output_path)
            output_file.parent.mkdir(parents=True, exist_ok=True)
            
            # Save presentation
            self.prs.save(output_path)
            logger.info(f"Presentation saved to: {output_path}")
            return True
            
        except Exception as e:
            logger.error(f"Failed to save presentation: {e}")
            return False
    
    def _set_slide_background_color(self, slide, hex_color: str):
        """
        Set solid color background for a slide.
        
        Args:
            slide: PowerPoint slide object
            hex_color: Hex color string
        """
        try:
            from pptx.dml.color import RGBColor
            
            rgb = self.style_mapper.hex_to_rgb(hex_color)
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*rgb)
            
        except Exception as e:
            logger.warning(f"Failed to set background color: {e}")
    
    def get_slide_count(self) -> int:
        """Get the number of slides in the presentation."""
        return len(self.prs.slides)
    
    def __enter__(self):
        """Context manager entry."""
        return self
    
    def __exit__(self, exc_type, exc_val, exc_tb):
        """Context manager exit."""
        # Cleanup if needed
        pass

"""
Style Mapper - Maps PDF styles to PowerPoint styles
"""

import logging
from typing import Dict, Any
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from lxml import etree
from .font_mapper import FontMapper

logger = logging.getLogger(__name__)


class StyleMapper:
    """
    Maps PDF visual styles to PowerPoint attributes.
    """
    
    def __init__(self, config: Dict[str, Any]):
        """
        Initialize Style Mapper.
        
        Args:
            config: Configuration dictionary
        """
        self.config = config
        self.font_mapper = FontMapper(config)
        
        # Handle both direct config and nested config structure
        if 'preserve_colors' in config:
            # Direct config (from config dict itself)
            self.preserve_colors = config.get('preserve_colors', True)
            self.default_font_size = config.get('default_font_size', 18)
            self.title_font_size = config.get('title_font_size', 32)
            self.font_size_scale = config.get('font_size_scale', 1.0)
            self.transparency_map = config.get('transparency_map', {})
        else:
            # Nested config (when 'mapper' section is passed)
            mapper_config = config.get('mapper', {})
            self.preserve_colors = mapper_config.get('preserve_colors', True)
            self.default_font_size = mapper_config.get('default_font_size', 18)
            self.title_font_size = mapper_config.get('title_font_size', 32)
            self.font_size_scale = mapper_config.get('font_size_scale', 1.0)
            self.transparency_map = mapper_config.get('transparency_map', {})
        
        # Count transparency mappings
        trans_count = sum(len(v) if isinstance(v, dict) else 1 for v in self.transparency_map.values())
        logger.info(f"StyleMapper initialized with font scale {self.font_size_scale}, {trans_count} transparency mappings")

    def update_font_scale(self, new_scale: float):
        """
        Update the font size scale factor.

        Args:
            new_scale: New scale factor to use
        """
        old_scale = self.font_size_scale
        self.font_size_scale = new_scale
        logger.info(f"Updated font scale: {old_scale:.3f} → {new_scale:.3f}")

    def apply_text_style(self, text_frame, style: Dict[str, Any]):
        """
        Apply text style to a PowerPoint text frame.
        
        Args:
            text_frame: PowerPoint text frame object
            style: Style dictionary
        """
        if not text_frame or not text_frame.paragraphs:
            return
        
        paragraph = text_frame.paragraphs[0]
        
        # Font name
        font_name = style.get('font_name', '')
        mapped_font = self.font_mapper.map_font(font_name)
        
        # Font size with scaling (PDF points to screen pixels)
        font_size_raw = style.get('font_size', self.default_font_size)
        font_size = font_size_raw * self.font_size_scale
        
        # Color
        color = style.get('color', '#000000')
        rgb = self.hex_to_rgb(color)
        
        # Bold and italic - support both 'bold'/'italic' and 'is_bold'/'is_italic' keys
        is_bold = style.get('is_bold', style.get('bold', False))
        is_italic = style.get('is_italic', style.get('italic', False))
        
        # Apply to all runs in paragraph
        for run in paragraph.runs:
            run.font.name = mapped_font
            run.font.size = Pt(font_size)
            
            # Only apply color if valid RGB was extracted
            if self.preserve_colors and rgb is not None:
                run.font.color.rgb = RGBColor(*rgb)
            
            run.font.bold = is_bold
            run.font.italic = is_italic
    
    def apply_shape_style(self, shape, style: Dict[str, Any]):
        """
        Apply style to a PowerPoint shape.
        
        Args:
            shape: PowerPoint shape object
            style: Style dictionary
        """
        if not style:
            return
        
        try:
            # Fill color (with optional transparency)
            fill_color = style.get('fill_color')
            # Use opacity from PDF extraction if available, otherwise use default or config
            fill_opacity = style.get('fill_opacity', 1.0)
            shape_role = style.get('role', 'unknown')  # Get shape role (border, card_background, etc.)
            
            # If opacity was not extracted from PDF, check config-based transparency mapping
            # This provides backward compatibility with manual transparency configuration
            if fill_opacity == 1.0 and fill_color:
                # Check role-based transparency mapping (new approach)
                if shape_role in self.transparency_map:
                    role_map = self.transparency_map[shape_role]
                    if isinstance(role_map, dict) and fill_color.lower() in role_map:
                        fill_opacity = role_map[fill_color.lower()]
                        logger.debug(f"Applied config transparency: {shape_role} {fill_color} -> opacity {fill_opacity}")
                elif not isinstance(self.transparency_map.get(shape_role), dict):
                    # Fallback to old flat mapping for compatibility
                    if fill_color.lower() in self.transparency_map:
                        fill_opacity = self.transparency_map[fill_color.lower()]
                        logger.debug(f"Applied config transparency (flat): {fill_color} -> opacity {fill_opacity}")
            elif fill_opacity < 1.0:
                # Opacity was extracted from PDF
                logger.debug(f"Using PDF-extracted opacity: {fill_opacity:.3f} for color {fill_color}")
            
            # Handle fill color
            # fill_color can be:
            #   - None: stroke-only shape (no fill) -> set transparent fill
            #   - "None": invalid value -> set transparent fill
            #   - hex color string: normal fill color
            if fill_color is None or fill_color == 'None':
                # Stroke-only shape or invalid fill - make fill transparent
                shape.fill.background()
                logger.debug(f"No fill color (stroke-only shape): fill_color={fill_color}")
            elif self.preserve_colors:
                # Normal fill color - convert and apply
                rgb = self.hex_to_rgb(fill_color)
                # Apply fill for any valid color (including black #000000 and white #FFFFFF)
                if rgb is not None:
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = RGBColor(*rgb)
                    
                    # Apply transparency if specified via XML manipulation
                    # PowerPoint alpha is in percentage * 1000 format (0-100000)
                    # HTML opacity is 0 (transparent) to 1 (opaque)
                    # So we need to convert opacity to alpha percentage
                    if fill_opacity < 1.0:
                        self._set_shape_transparency(shape, fill_opacity)
                        logger.debug(f"Set shape transparency: opacity {fill_opacity:.2f}")
                else:
                    # hex_to_rgb failed - make fill transparent
                    shape.fill.background()
                    logger.warning(f"Invalid fill color: {fill_color}")
            
            # Stroke/line color and width
            stroke_color = style.get('stroke_color')
            stroke_width = style.get('stroke_width')
            is_ring = style.get('is_ring', False)
            
            # Determine if border should be rendered
            # Border should NOT be rendered if:
            # 1. stroke_width is None or 0
            # 2. stroke_color is None or "None" string
            # 3. stroke_color is default black (#000000) AND stroke_width is None (PDF default)
            # EXCEPTION: Ring shapes always render their stroke (it's the ring itself)
            should_render_border = False
            
            if is_ring:
                # For ring shapes, always render the stroke (it's the ring color)
                should_render_border = True
                logger.debug(f"Ring shape detected, rendering stroke as ring")
            elif stroke_color and stroke_color not in ['None', '#000000']:
                # Explicit non-black stroke color
                if stroke_width is not None and stroke_width > 0:
                    should_render_border = True
            elif stroke_color == '#000000' and stroke_width is not None and stroke_width > 0:
                # Black stroke with explicit width
                should_render_border = True
            
            if should_render_border and self.preserve_colors:
                rgb = self.hex_to_rgb(stroke_color)
                shape.line.color.rgb = RGBColor(*rgb)
                # Apply PDF scale correction: PDF border 3pt → HTML 4px means multiply by 1.333
                # At 72 DPI: 4px = 4pt in PowerPoint
                scaled_width = stroke_width * self.font_size_scale
                shape.line.width = Pt(scaled_width)
                logger.debug(f"Stroke width: {stroke_width}pt → {scaled_width}pt (scaled)")
            else:
                # No border - make line transparent/invisible
                shape.line.fill.background()
                logger.debug(f"No border: stroke_color={stroke_color}, stroke_width={stroke_width}")
        except Exception as e:
            logger.warning(f"Failed to apply shape style: {e}", exc_info=True)
    
    def hex_to_rgb(self, hex_color: str) -> tuple:
        """
        Convert hex color to RGB tuple.
        
        Args:
            hex_color: Hex color string (e.g., '#FF0000'), or None
            
        Returns:
            RGB tuple (r, g, b), or None if conversion fails
        """
        if not hex_color or not isinstance(hex_color, str):
            return None
        
        # Remove '#' if present
        hex_color = hex_color.lstrip('#')
        
        try:
            # Convert hex to RGB
            if len(hex_color) == 6:
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                return (r, g, b)
            elif len(hex_color) == 3:
                r = int(hex_color[0] * 2, 16)
                g = int(hex_color[1] * 2, 16)
                b = int(hex_color[2] * 2, 16)
                return (r, g, b)
        except ValueError:
            logger.warning(f"Invalid hex color: {hex_color}")
        
        return None
    
    def rgba_to_rgb_opacity(self, rgba_str: str) -> tuple:
        """
        Parse RGBA string and extract RGB and opacity.
        
        Args:
            rgba_str: RGBA string like 'rgba(10, 66, 117, 0.08)'
            
        Returns:
            Tuple of (rgb_tuple or None, opacity_float)
        """
        import re
        
        if not rgba_str or not isinstance(rgba_str, str):
            return (None, 1.0)
        
        # Parse rgba(r, g, b, a) format
        match = re.match(r'rgba?\((\d+),\s*(\d+),\s*(\d+)(?:,\s*([\d.]+))?\)', rgba_str)
        if match:
            r, g, b, a = match.groups()
            rgb = (int(r), int(g), int(b))
            opacity = float(a) if a else 1.0
            return (rgb, opacity)
        
        # If not RGBA, try hex
        if rgba_str.startswith('#'):
            rgb = self.hex_to_rgb(rgba_str)
            return (rgb, 1.0)
        
        return (None, 1.0)
    
    def _set_shape_transparency(self, shape, opacity: float):
        """
        Set shape transparency via XML manipulation.
        
        Args:
            shape: PowerPoint shape object
            opacity: Opacity value (0.0-1.0, where 1.0 is fully opaque)
        """
        try:
            # Access shape's XML element
            spPr = shape.element.spPr
            
            # Find solidFill element
            ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
            solidFill = spPr.find('.//a:solidFill', ns)
            
            if solidFill is not None:
                # Find srgbClr element
                srgbClr = solidFill.find('.//a:srgbClr', ns)
                if srgbClr is not None:
                    # Remove existing alpha element if present
                    existing_alpha = srgbClr.find('.//a:alpha', ns)
                    if existing_alpha is not None:
                        srgbClr.remove(existing_alpha)
                    
                    # Add new alpha element
                    # Alpha in PowerPoint XML is in percentage * 1000 format (0-100000)
                    # opacity 0.08 means 8% opacity = 8000 in alpha value
                    alpha_value = int(opacity * 100000)
                    alpha_elem = etree.SubElement(
                        srgbClr, 
                        '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha'
                    )
                    alpha_elem.set('val', str(alpha_value))
                    logger.debug(f"Set XML alpha value: {alpha_value} (opacity: {opacity:.2f})")
        except Exception as e:
            logger.warning(f"Failed to set shape transparency via XML: {e}", exc_info=True)
    
    def normalize_font_size(self, pdf_font_size: float, is_title: bool = False) -> float:
        """
        Normalize PDF font size to appropriate PowerPoint size.
        
        Args:
            pdf_font_size: Original PDF font size
            is_title: Whether this is a title text
            
        Returns:
            Normalized font size in points
        """
        if is_title:
            # Titles should be larger
            return max(self.title_font_size, pdf_font_size * 0.75)
        
        # Regular text
        if pdf_font_size < 10:
            return 10
        elif pdf_font_size > 36:
            return 36
        
        return pdf_font_size

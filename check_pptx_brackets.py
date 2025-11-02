#!/usr/bin/env python3
"""
检查生成的PPTX中括号相关文本的顺序和位置
"""
from pptx import Presentation
import sys


def check_brackets(pptx_path):
    """检查PPTX中第15页的括号文本"""
    prs = Presentation(pptx_path)
    
    if len(prs.slides) < 15:
        print(f"PPT只有 {len(prs.slides)} 页")
        return
    
    slide = prs.slides[14]  # 第15页，索引14
    print(f"\n{'='*80}")
    print(f"第15页 - 括号文本检查")
    print(f"{'='*80}\n")
    
    # 收集所有与"10.74.145.44"和"未知业务"相关的文本框
    related_shapes = []
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text.strip()
            
            # 检查是否包含相关文本
            if any(keyword in text for keyword in ["10.74.145.44", "未知业务", "核心业务", "(", ")", "（", "）"]):
                related_shapes.append({
                    'text': text,
                    'left': shape.left.inches,
                    'top': shape.top.inches,
                    'rotation': shape.rotation,
                    'shape': shape
                })
    
    # 按照位置排序（先按top，再按left）
    related_shapes.sort(key=lambda x: (x['top'], x['left']))
    
    print("相关文本框（按位置排序）:")
    for i, info in enumerate(related_shapes):
        print(f"\n{i+1}. 文本: '{info['text']}'")
        print(f"   位置: left={info['left']:.3f}\", top={info['top']:.3f}\"")
        print(f"   旋转: {info['rotation']}°")
        
        # 检查括号问题
        if "10.74.145.44" in info['text']:
            # 检查这个文本框附近是否有括号相关的文本框
            nearby = []
            for j, other in enumerate(related_shapes):
                if i != j:
                    # 检查是否在附近（比如0.5英寸以内）
                    dist_x = abs(info['left'] - other['left'])
                    dist_y = abs(info['top'] - other['top'])
                    dist = (dist_x ** 2 + dist_y ** 2) ** 0.5
                    
                    if dist < 0.5:  # 0.5英寸
                        nearby.append({
                            'text': other['text'],
                            'distance': dist
                        })
            
            if nearby:
                print(f"   附近的文本框:")
                for n in sorted(nearby, key=lambda x: x['distance']):
                    print(f"     - '{n['text']}' (距离: {n['distance']:.3f}\")")
    
    # 特别检查"10.74.145.44 (未知业务)"的组合
    print(f"\n{'='*80}")
    print("问题诊断:")
    print(f"{'='*80}\n")
    
    # 查找包含"10.74.145.44"的文本框
    ip_shapes = [s for s in related_shapes if "10.74.145.44" in s['text']]
    bracket_shapes = [s for s in related_shapes if s['text'] in ['(', ')', '（', '）']]
    biz_shapes = [s for s in related_shapes if "未知业务" in s['text'] and "10.74.145.44" not in s['text']]
    
    print(f"找到 {len(ip_shapes)} 个包含 '10.74.145.44' 的文本框")
    print(f"找到 {len(bracket_shapes)} 个独立括号文本框")
    print(f"找到 {len(biz_shapes)} 个包含 '未知业务' 的文本框")
    
    if bracket_shapes:
        print(f"\n⚠️ 发现问题：括号被分离成独立的文本框！")
        print(f"\n独立括号文本框:")
        for b in bracket_shapes:
            print(f"  '{b['text']}' at ({b['left']:.3f}\", {b['top']:.3f}\"), 旋转={b['rotation']}°")
        
        print(f"\n💡 解决方案：需要在文本合并逻辑中，将临近的文本元素（IP地址+括号+业务类型）合并成一个文本框")


if __name__ == "__main__":
    pptx_file = sys.argv[1] if len(sys.argv) > 1 else "/tmp/test_rotation_fixed.pptx"
    check_brackets(pptx_file)

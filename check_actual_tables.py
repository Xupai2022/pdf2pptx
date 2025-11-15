"""
检查实际的表格结构（不做任何假设）
"""
import sys
sys.path.append('.')

from main import convert_pdf_to_pptx, load_config
from pptx import Presentation

def check_pdf_tables(pdf_path, pages_to_check):
    """检查PDF中实际的表格"""
    output_path = f"output/check_{pdf_path.split('/')[-1].replace('.pdf', '.pptx')}"
    
    config = load_config()
    print(f"\n转换 {pdf_path}...")
    if not convert_pdf_to_pptx(pdf_path, output_path, config):
        print("转换失败")
        return
    
    prs = Presentation(output_path)
    print(f"\n{'='*80}")
    print(f"文件: {pdf_path}")
    print(f"总页数: {len(prs.slides)}")
    print('='*80)
    
    for page_num in pages_to_check:
        if page_num > len(prs.slides):
            print(f"\n第{page_num}页: 不存在")
            continue
        
        slide = prs.slides[page_num - 1]
        tables = [shape for shape in slide.shapes if shape.has_table]
        
        print(f"\n第{page_num}页:")
        if not tables:
            print(f"  没有表格")
            continue
        
        print(f"  找到 {len(tables)} 个表格:")
        for idx, shape in enumerate(tables):
            table = shape.table
            print(f"\n  表格 {idx+1}:")
            print(f"    尺寸: {len(table.rows)}行 x {len(table.columns)}列")
            print(f"    内容（前3行）:")
            for row_idx in range(min(3, len(table.rows))):
                row = table.rows[row_idx]
                texts = [cell.text.strip()[:20] if cell.text.strip() else '(空)' for cell in row.cells]
                print(f"      行{row_idx+1}: {' | '.join(texts)}")

def main():
    print("="*80)
    print("检查安全运营月报.pdf的表格")
    print("="*80)
    check_pdf_tables("tests/安全运营月报.pdf", [8, 9, 12])
    
    print("\n" + "="*80)
    print("检查season_report_del.pdf的表格")
    print("="*80)
    check_pdf_tables("tests/season_report_del.pdf", [2, 6, 7, 10, 13, 16])

if __name__ == "__main__":
    main()

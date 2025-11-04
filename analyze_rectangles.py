#!/usr/bin/env python3
"""
Analyze PPTX file to detect unwanted rectangular elements
especially on pages 3 and 4 where black and green rectangles appear.
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

def analyze_shapes(pptx_path):
    """Analyze shapes in the PPTX file."""
    prs = Presentation(pptx_path)
    
    print(f"Analyzing: {pptx_path}")
    print(f"Total slides: {len(prs.slides)}\n")
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        print(f"{'='*80}")
        print(f"Slide {slide_idx}")
        print(f"{'='*80}")
        
        rectangles = []
        text_boxes = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            # Get shape type
            shape_type = shape.shape_type
            
            # Get position and size
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            
            # Check if it's a rectangle or auto shape
            if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE or shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                # Try to get fill color
                fill_color = None
                try:
                    if shape.fill.type == 1:  # SOLID
                        fill_color = shape.fill.fore_color.rgb
                except:
                    pass
                
                # Try to get line color
                line_color = None
                line_width = 0
                try:
                    if hasattr(shape, 'line') and shape.line:
                        line_width = shape.line.width if hasattr(shape.line, 'width') else 0
                        if hasattr(shape.line, 'color'):
                            line_color = shape.line.color.rgb
                except:
                    pass
                
                # Check if it has text
                has_text = False
                text_content = ""
                try:
                    if hasattr(shape, 'text_frame'):
                        text_content = shape.text_frame.text.strip()
                        has_text = len(text_content) > 0
                except:
                    pass
                
                info = {
                    'index': shape_idx,
                    'type': shape_type,
                    'left': left,
                    'top': top,
                    'width': width,
                    'height': height,
                    'fill_color': fill_color,
                    'line_color': line_color,
                    'line_width': line_width,
                    'has_text': has_text,
                    'text': text_content[:50] if text_content else ""
                }
                
                # Categorize
                if has_text:
                    text_boxes.append(info)
                else:
                    # Empty shape - likely unwanted rectangle
                    rectangles.append(info)
        
        # Report rectangles (potentially unwanted)
        if rectangles:
            print(f"\n🔴 Found {len(rectangles)} rectangles (no text):")
            for rect in rectangles:
                print(f"  Shape #{rect['index']}: "
                      f"pos=({rect['left']/914400:.1f}\", {rect['top']/914400:.1f}\"), "
                      f"size=({rect['width']/914400:.2f}\"x{rect['height']/914400:.2f}\"), "
                      f"fill={rect['fill_color']}, "
                      f"line={rect['line_color']} (width={rect['line_width']/914400:.3f}\")")
        
        # Report text boxes
        if text_boxes:
            print(f"\n✅ Found {len(text_boxes)} text boxes:")
            for tb in text_boxes[:10]:  # Show first 10
                print(f"  Shape #{tb['index']}: \"{tb['text']}\" "
                      f"pos=({tb['left']/914400:.1f}\", {tb['top']/914400:.1f}\"), "
                      f"fill={tb['fill_color']}")
        
        print()

if __name__ == "__main__":
    import sys
    pptx_path = sys.argv[1] if len(sys.argv) > 1 else "output/安全运营月报_test.pptx"
    analyze_shapes(pptx_path)

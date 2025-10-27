#!/usr/bin/env python3
"""
测试PDF转PPT的颜色识别准确性
"""
import sys
import os
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from src.parser.pdf_parser import PDFParser
from src.generator.pptx_generator import PPTXGenerator
from pptx import Presentation
import yaml

def test_pdf_color_extraction(pdf_path):
    """测试PDF颜色提取"""
    print(f"Testing PDF: {pdf_path}")
    print("=" * 80)
    
    # Load config
    with open('config/config.yaml', 'r', encoding='utf-8') as f:
        config = yaml.safe_load(f)
    
    # Parse PDF
    parser = PDFParser(config.get('parser', {}))
    if not parser.open(pdf_path):
        print("Failed to open PDF")
        return
    
    # Extract page 3 (index 2) which has transparency
    page_data = parser.extract_page_elements(2)
    
    print(f"\nPage 3 Analysis:")
    print(f"Total elements: {len(page_data['elements'])}")
    
    # Analyze shapes with color
    shapes_by_opacity = {}
    for elem in page_data['elements']:
        if elem['type'] == 'shape':
            fill_color = elem.get('fill_color', 'None')
            fill_opacity = elem.get('fill_opacity', 1.0)
            
            key = f"{fill_color}_{fill_opacity:.4f}"
            if key not in shapes_by_opacity:
                shapes_by_opacity[key] = []
            shapes_by_opacity[key].append(elem)
    
    print(f"\nShapes grouped by (color, opacity):")
    for key, shapes in sorted(shapes_by_opacity.items()):
        color, opacity = key.rsplit('_', 1)
        print(f"  {color} @ opacity {opacity}: {len(shapes)} shapes")
        # Show first shape details
        if shapes:
            s = shapes[0]
            print(f"    Example: pos=({s['x']:.1f}, {s['y']:.1f}), size=({s['width']:.1f}x{s['height']:.1f})")
    
    parser.close()
    return page_data

def test_pptx_generation(pdf_path, output_path):
    """测试PPTX生成"""
    print(f"\nGenerating PPTX: {output_path}")
    print("=" * 80)
    
    # Load config
    with open('config/config.yaml', 'r', encoding='utf-8') as f:
        config = yaml.safe_load(f)
    
    # Import required modules
    from src.analyzer.layout_analyzer_v2 import LayoutAnalyzerV2
    from src.rebuilder.coordinate_mapper import CoordinateMapper
    
    # Parse PDF
    parser = PDFParser(config.get('parser', {}))
    if not parser.open(pdf_path):
        print("Failed to open PDF")
        return False
    
    all_pages = parser.extract_all_pages()
    parser.close()
    
    # Analyze layout
    analyzer = LayoutAnalyzerV2(config.get('analyzer', {}))
    analyzed_pages = []
    for page_data in all_pages:
        layout_data = analyzer.analyze_page(page_data)
        analyzed_pages.append(layout_data)
    
    # Build slide models
    mapper = CoordinateMapper(config.get('rebuilder', {}))
    slide_models = []
    for layout_data in analyzed_pages:
        slide_model = mapper.create_slide_model(layout_data)
        slide_models.append(slide_model)
    
    # Generate PPTX
    generator = PPTXGenerator(config)
    if not generator.generate_from_models(slide_models):
        print(f"✗ Failed to generate PPTX")
        return False
    
    if not generator.save(output_path):
        print(f"✗ Failed to save PPTX")
        return False
    
    print(f"✓ Generated: {output_path}")
    return True

def analyze_pptx_colors(pptx_path):
    """分析生成的PPTX中的颜色"""
    print(f"\nAnalyzing PPTX: {pptx_path}")
    print("=" * 80)
    
    prs = Presentation(pptx_path)
    
    # Analyze slide 3 (index 2)
    if len(prs.slides) < 3:
        print("PPTX has less than 3 slides")
        return
    
    slide = prs.slides[2]
    print(f"\nSlide 3 - Total shapes: {len(slide.shapes)}")
    
    # Group shapes by color and transparency
    shape_groups = {}
    
    for idx, shape in enumerate(slide.shapes):
        try:
            if hasattr(shape, 'fill') and shape.fill.type == 1:  # SOLID fill
                color = shape.fill.fore_color.rgb
                color_hex = f"#{color[0]:02X}{color[1]:02X}{color[2]:02X}"
                
                # Try to extract alpha from XML
                alpha = 100000  # Default fully opaque
                try:
                    spPr = shape.element.spPr
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    solidFill = spPr.find('.//a:solidFill', ns)
                    if solidFill is not None:
                        srgbClr = solidFill.find('.//a:srgbClr', ns)
                        if srgbClr is not None:
                            alpha_elem = srgbClr.find('.//a:alpha', ns)
                            if alpha_elem is not None:
                                alpha = int(alpha_elem.get('val', 100000))
                except:
                    pass
                
                opacity = alpha / 100000.0
                key = f"{color_hex}_{opacity:.4f}"
                
                if key not in shape_groups:
                    shape_groups[key] = []
                shape_groups[key].append({
                    'index': idx,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                })
        except:
            pass
    
    print(f"\nShapes in PPTX grouped by (color, opacity):")
    for key, shapes in sorted(shape_groups.items()):
        color, opacity = key.rsplit('_', 1)
        print(f"  {color} @ opacity {opacity}: {len(shapes)} shapes")
        if shapes:
            s = shapes[0]
            print(f"    Example: left={s['left']}, top={s['top']}, size={s['width']}x{s['height']}")

if __name__ == "__main__":
    pdf_path = "tests/glm-4.6.pdf"
    output_path = "output/test_glm46_colors.pptx"
    
    # Test 1: Extract colors from PDF
    print("\n" + "="*80)
    print("TEST 1: PDF Color Extraction")
    print("="*80)
    page_data = test_pdf_color_extraction(pdf_path)
    
    # Test 2: Generate PPTX
    print("\n" + "="*80)
    print("TEST 2: PPTX Generation")
    print("="*80)
    success = test_pptx_generation(pdf_path, output_path)
    
    # Test 3: Analyze PPTX
    if success:
        print("\n" + "="*80)
        print("TEST 3: PPTX Color Analysis")
        print("="*80)
        analyze_pptx_colors(output_path)
    
    print("\n" + "="*80)
    print("TEST COMPLETE")
    print("="*80)

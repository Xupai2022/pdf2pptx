"""
回归测试 - 验证安全运营月报.pdf和season_report_del.pdf的表格
"""
import sys
import logging
sys.path.append('.')

logging.basicConfig(level=logging.INFO, format='%(levelname)s - %(message)s')

from main import convert_pdf_to_pptx, load_config
from pptx import Presentation

def check_tables_in_slide(prs, page_num, expected_info):
    """检查指定页面的表格"""
    if len(prs.slides) < page_num:
        print(f"  ⚠️  只有 {len(prs.slides)} 页")
        return False
    
    slide = prs.slides[page_num - 1]
    
    tables = [shape for shape in slide.shapes if shape.has_table]
    print(f"  第{page_num}页: 找到 {len(tables)} 个表格")
    
    if expected_info.get('no_table'):
        if len(tables) == 0:
            print(f"    ✓ 正确：没有表格")
            return True
        else:
            print(f"    ✗ 错误：应该没有表格，但发现了 {len(tables)} 个")
            return False
    
    if len(tables) == 0:
        print(f"    ✗ 错误：没有找到表格")
        return False
    
    # 检查第一个表格的尺寸
    table = tables[0].table
    expected_rows = expected_info.get('rows')
    expected_cols = expected_info.get('cols')
    
    if expected_rows and len(table.rows) != expected_rows:
        print(f"    ✗ 行数错误：期望{expected_rows}行，实际{len(table.rows)}行")
        return False
    
    if expected_cols and len(table.columns) != expected_cols:
        print(f"    ✗ 列数错误：期望{expected_cols}列，实际{len(table.columns)}列")
        return False
    
    print(f"    ✓ 表格尺寸正确: {len(table.rows)}行 x {len(table.columns)}列")
    
    # 显示表格内容（前3行）
    if expected_info.get('show_content'):
        print(f"    表格内容（前3行）:")
        for row_idx in range(min(3, len(table.rows))):
            row = table.rows[row_idx]
            texts = [cell.text.strip()[:15] if cell.text.strip() else '(空)' for cell in row.cells]
            print(f"      行{row_idx+1}: {' | '.join(texts)}")
    
    return True

def main():
    config = load_config()
    
    # 测试1: 安全运营月报.pdf
    print("="*80)
    print("测试1: 安全运营月报.pdf - 验证第8,9,12页表格")
    print("="*80)
    
    pdf1 = "tests/安全运营月报.pdf"
    output1 = "output/安全运营月报_test.pptx"
    
    print(f"\n转换 {pdf1}...")
    if not convert_pdf_to_pptx(pdf1, output1, config):
        print("❌ 转换失败")
        return 1
    
    print(f"\n检查表格...")
    prs1 = Presentation(output1)
    
    test1_cases = {
        8: {'rows': 11, 'cols': 4, 'show_content': True},
        9: {'rows': 7, 'cols': 4, 'show_content': True},
        12: {'rows': 11, 'cols': 5, 'show_content': True}
    }
    
    all_passed = True
    for page, expected in test1_cases.items():
        if not check_tables_in_slide(prs1, page, expected):
            all_passed = False
    
    # 测试2: season_report_del.pdf
    print("\n" + "="*80)
    print("测试2: season_report_del.pdf - 验证第2,6,7,10,13,16页")
    print("="*80)
    
    pdf2 = "tests/season_report_del.pdf"
    output2 = "output/season_report_del_test.pptx"
    
    print(f"\n转换 {pdf2}...")
    if not convert_pdf_to_pptx(pdf2, output2, config):
        print("❌ 转换失败")
        return 1
    
    print(f"\n检查页面...")
    prs2 = Presentation(output2)
    
    test2_cases = {
        2: {'no_table': True},
        6: {'no_table': True},
        7: {'rows': 7, 'cols': 3},
        10: {'rows': 6, 'cols': 3},
        13: {'rows': 3, 'cols': 3},
        16: {'rows': 7, 'cols': 4, 'show_content': True}
    }
    
    for page, expected in test2_cases.items():
        if not check_tables_in_slide(prs2, page, expected):
            all_passed = False
    
    # 总结
    print("\n" + "="*80)
    if all_passed:
        print("✅ 所有测试通过！")
        return 0
    else:
        print("❌ 部分测试失败")
        return 1

if __name__ == "__main__":
    sys.exit(main())

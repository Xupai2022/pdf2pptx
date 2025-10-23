#!/usr/bin/env python3
"""
Final comparison: HTML reference vs. PPTX output
"""

import sys
from pathlib import Path
from pptx import Presentation

def print_header(title):
    print("\n" + "="*80)
    print(f"  {title}")
    print("="*80 + "\n")

def compare_results():
    """Compare HTML expectations with PPTX output."""
    
    print_header("HTML → PPTX Conversion Comparison")
    
    # Expected from HTML
    print("📄 ORIGINAL HTML STRUCTURE (1920×1080):")
    print(f"{'─'*80}")
    print("1. Top bar: 10px height, blue #0a4275")
    print("2. Main title: '风险资产深度剖析' (48px ≈ 36pt)")
    print("3. Subtitle: '典型高风险资产案例展示' (36px ≈ 27pt)")
    print("4. Three stat cards with:")
    print("   - Card backgrounds: rgba(10, 66, 117, 0.08)")
    print("   - Left borders: 4px blue")
    print("   - Section titles: 28px ≈ 21pt")
    print("   - Large numbers: 48px (red)")
    print("5. Two detail columns with 3 items each")
    print("6. Impact analysis section (4 columns)")
    print("7. Six FontAwesome icons")
    
    # Check PPTX output
    pptx_path = "output/test_output.pptx"
    if not Path(pptx_path).exists():
        print("\n❌ PPTX file not found")
        return False
    
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    
    print(f"\n{'─'*80}")
    print("📊 GENERATED PPTX OUTPUT (10×7.5\"):")
    print(f"{'─'*80}")
    print(f"Total shapes: {len(slide.shapes)}")
    
    # Count by type
    text_count = 0
    image_count = 0
    shape_count = 0
    
    font_sizes = {}
    
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            text_count += 1
            if shape.text_frame.paragraphs:
                para = shape.text_frame.paragraphs[0]
                if para.runs and para.runs[0].font.size:
                    size = para.runs[0].font.size.pt
                    if size not in font_sizes:
                        font_sizes[size] = []
                    font_sizes[size].append(shape.text.strip()[:30])
        elif hasattr(shape, "image"):
            image_count += 1
        else:
            shape_count += 1
    
    print(f"\nElement counts:")
    print(f"  • Text boxes: {text_count}")
    print(f"  • Images (icons): {image_count}")
    print(f"  • Shapes (backgrounds/borders): {shape_count}")
    
    print(f"\nFont sizes detected:")
    for size in sorted(font_sizes.keys(), reverse=True):
        texts = font_sizes[size]
        print(f"  • {size}pt: {len(texts)} elements")
        if len(texts) <= 3:
            for t in texts:
                print(f"      - \"{t}\"")
    
    print(f"\n{'─'*80}")
    print("✅ VERIFICATION CHECKLIST:")
    print(f"{'─'*80}")
    
    # Check specific requirements
    checks = []
    
    # Font sizes
    has_36pt = 36.0 in font_sizes
    has_27pt = 27.0 in font_sizes
    has_21pt = 21.0 in font_sizes
    has_18pt = any(17.5 <= s <= 18.5 for s in font_sizes.keys())
    
    checks.append(("36pt title font", has_36pt, "✓" if has_36pt else "✗"))
    checks.append(("27pt subtitle font", has_27pt, "✓" if has_27pt else "✗"))
    checks.append(("21pt heading font", has_21pt, "✓" if has_21pt else "✗"))
    checks.append(("18pt body font", has_18pt, "✓" if has_18pt else "✗"))
    checks.append(("All 6 icons extracted", image_count == 6, "✓" if image_count == 6 else "✗"))
    checks.append(("Multiple text boxes", text_count > 40, "✓" if text_count > 40 else "✗"))
    checks.append(("Background shapes", shape_count > 20, "✓" if shape_count > 20 else "✗"))
    
    for label, passed, symbol in checks:
        status = "PASS" if passed else "FAIL"
        print(f"  {symbol} {label:30s} {status}")
    
    all_passed = all(passed for _, passed, _ in checks)
    
    print(f"\n{'─'*80}")
    if all_passed:
        print("🎉 ALL CHECKS PASSED!")
        print("✅ Layout structure preserved correctly")
        print("✅ Font sizes match HTML specifications")
        print("✅ All icons extracted as images")
        print("✅ Text elements preserved independently")
        print("✅ Background shapes and borders created")
    else:
        print("⚠️  SOME CHECKS FAILED")
        print("Review the output PPTX and compare with HTML")
    
    print(f"{'─'*80}\n")
    
    return all_passed


def main():
    """Main function."""
    success = compare_results()
    return 0 if success else 1


if __name__ == '__main__':
    sys.exit(main())

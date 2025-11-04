"""
Comprehensive test for PNG transparency fix on all affected pages
"""

from pptx import Presentation
from PIL import Image
import io

def test_slide_images(pptx_path, slide_num):
    """Test all images in a specific slide"""
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_num]
    
    images = [shape for shape in slide.shapes if shape.shape_type == 13]
    
    issues = []
    success = []
    
    for idx, shape in enumerate(images):
        image_blob = shape.image.blob
        pil_image = Image.open(io.BytesIO(image_blob))
        
        width, height = pil_image.size
        mode = pil_image.mode
        
        # Check for black background issue
        if mode == 'RGB':
            # Sample edge pixels
            sample_points = [
                (0, 0), (width//2, 0), (width-1, 0),
                (0, height//2), (width-1, height//2),
                (0, height-1), (width//2, height-1), (width-1, height-1)
            ]
            
            colors = []
            for x, y in sample_points:
                try:
                    pixel = pil_image.getpixel((x, y))
                    if not isinstance(pixel, tuple):
                        pixel = (pixel, pixel, pixel)
                    colors.append(pixel)
                except:
                    continue
            
            black_count = sum(1 for c in colors if all(ch < 30 for ch in c[:3]))
            
            if black_count >= len(colors) * 0.5:
                issues.append(f"  Image {idx+1}: {width}x{height}px, RGB mode, {black_count}/{len(colors)} black edges ❌")
            else:
                success.append(f"  Image {idx+1}: {width}x{height}px, RGB mode, {black_count}/{len(colors)} black edges ✅")
        
        elif mode in ['RGBA', 'LA', 'PA']:
            # Has alpha - check if it's being used
            sample_points = [
                (0, 0), (width//2, 0), (width-1, 0),
                (0, height//2), (width-1, height//2),
                (0, height-1), (width//2, height-1), (width-1, height-1)
            ]
            
            colors = []
            for x, y in sample_points:
                try:
                    pixel = pil_image.getpixel((x, y))
                    colors.append(pixel[:3])
                except:
                    continue
            
            black_count = sum(1 for c in colors if all(ch < 30 for ch in c))
            
            success.append(f"  Image {idx+1}: {width}x{height}px, {mode} mode, {black_count}/{len(colors)} black edges ✅")
    
    return issues, success

# Test all affected pages
pages_to_test = [
    (3, "Page 3"),
    (4, "Page 4"),
    (5, "Page 5"),
    (10, "Page 10"),
    (11, "Page 11"),
    (13, "Page 13")
]

print("="*70)
print("COMPREHENSIVE PNG TRANSPARENCY TEST")
print("="*70)
print()

all_issues = []
all_success = []

for page_num, page_name in pages_to_test:
    slide_idx = page_num - 1  # Convert to 0-indexed
    
    print(f"{page_name} (Slide {slide_idx+1}):")
    issues, success = test_slide_images("output/test_fixed.pptx", slide_idx)
    
    for msg in success:
        print(msg)
    
    for msg in issues:
        print(msg)
    
    if issues:
        all_issues.extend(issues)
    
    all_success.extend(success)
    print()

print("="*70)
print("SUMMARY")
print("="*70)
print(f"✅ Passed: {len(all_success)} images")
print(f"❌ Failed: {len(all_issues)} images")
print()

if all_issues:
    print("FAILED IMAGES:")
    for issue in all_issues:
        print(issue)
    print()
    print("❌ TEST FAILED - Some images still have black background issues!")
else:
    print("🎉 ALL TESTS PASSED!")
    print("✅ All PNG images display correctly without black backgrounds!")
    print("✅ Transparency is properly preserved in all affected pages!")

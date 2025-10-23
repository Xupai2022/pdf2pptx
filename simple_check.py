#!/usr/bin/env python3
"""简单检查PPTX"""

from pptx import Presentation

prs = Presentation("output/test_output.pptx")

print(f"Slide dimensions: {prs.slide_width / 914400:.3f}\" × {prs.slide_height / 914400:.3f}\"")
print(f"Slides: {len(prs.slides)}")

if prs.slides:
    slide = prs.slides[0]
    print(f"\nSlide 1 shapes: {len(slide.shapes)}")
    
    text_count = 0
    image_count = 0
    shape_count = 0
    
    for shape in slide.shapes:
        if hasattr(shape, 'text'):
            text = shape.text
            if text and text.strip():
                text_count += 1
                if text_count <= 5:
                    print(f"  Text {text_count}: \"{text[:50]}\"")
        elif hasattr(shape, 'image'):
            image_count += 1
        else:
            shape_count += 1
    
    print(f"\nSummary:")
    print(f"  Text shapes with content: {text_count}")
    print(f"  Images: {image_count}")
    print(f"  Other shapes: {shape_count}")

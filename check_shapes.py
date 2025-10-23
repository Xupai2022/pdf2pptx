#!/usr/bin/env python3
"""检查PPT中的形状，特别是top-bar"""

from pptx import Presentation

prs = Presentation("output/test_output.pptx")
slide = prs.slides[0]

print("=" * 60)
print("SHAPES CHECK - Looking for top bar")
print("=" * 60)

print(f"\nTotal shapes: {len(slide.shapes)}")

# Find all shapes (not text, not images)
non_text_shapes = []
for i, shape in enumerate(slide.shapes):
    try:
        # Check if it's not a text box with content
        is_text = hasattr(shape, 'text') and shape.text and shape.text.strip()
        is_image = hasattr(shape, 'image')
        
        if not is_text and not is_image:
            w_in = shape.width / 914400
            h_in = shape.height / 914400
            t_in = shape.top / 914400
            l_in = shape.left / 914400
            
            w_pt = w_in * 72
            h_pt = h_in * 72
            
            non_text_shapes.append({
                'index': i,
                'width_in': w_in,
                'height_in': h_in,
                'width_pt': w_pt,
                'height_pt': h_pt,
                'top_in': t_in,
                'left_in': l_in,
                'type': shape.shape_type
            })
    except:
        pass

print(f"\nNon-text, non-image shapes: {len(non_text_shapes)}")

# Sort by top position (top-bar should be at the top)
non_text_shapes.sort(key=lambda x: x['top_in'])

print("\nFirst 10 shapes (sorted by top position):")
for shape_info in non_text_shapes[:10]:
    print(f"\n  Shape {shape_info['index']}:")
    print(f"    Position: ({shape_info['left_in']:.3f}\", {shape_info['top_in']:.3f}\")")
    print(f"    Size: {shape_info['width_in']:.3f}\" × {shape_info['height_in']:.3f}\"")
    print(f"          ({shape_info['width_pt']:.1f}pt × {shape_info['height_pt']:.1f}pt)")
    print(f"    Type: {shape_info['type']}")
    
    # Check if it matches top-bar criteria
    # Expected: width ~26.67", height ~0.139" (10pt)
    if shape_info['width_pt'] > 1000 and shape_info['height_pt'] < 20:
        print(f"    ⭐ LIKELY TOP BAR (wide and thin)")
        if abs(shape_info['height_pt'] - 10) < 2:
            print(f"    ✓ Height matches expected 10pt!")
        else:
            print(f"    ✗ Height {shape_info['height_pt']:.1f}pt, expected 10pt")

# Check for rectangles with specific dimensions
print("\n" + "=" * 60)
print("Looking for top-bar specifically (height ~10pt, width ~1920pt):")
print("=" * 60)

potential_top_bars = [s for s in non_text_shapes 
                      if s['width_pt'] > 1000 and s['height_pt'] < 30 and s['top_in'] < 1]

if potential_top_bars:
    for bar in potential_top_bars:
        print(f"\n  Candidate at position ({bar['left_in']:.3f}\", {bar['top_in']:.3f}\"):")
        print(f"    Size: {bar['width_pt']:.1f}pt × {bar['height_pt']:.1f}pt")
        print(f"    Expected: 1920pt × 10pt")
        print(f"    Difference: {abs(bar['width_pt'] - 1920):.1f}pt width, {abs(bar['height_pt'] - 10):.1f}pt height")
else:
    print("  ✗ No shapes matching top-bar criteria found")

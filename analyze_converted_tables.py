"""
分析转换后的表格结构
"""

import sys
sys.path.insert(0, '/home/user/webapp')

from pptx import Presentation
import json

def analyze_pptx_tables(pptx_path, page_nums):
    """
    分析PPTX中指定页面的表格
    
    Args:
        pptx_path: PPTX文件路径
        page_nums: 要分析的页码列表 (1-based)
    """
    prs = Presentation(pptx_path)
    
    for page_num in page_nums:
        slide_idx = page_num - 1
        
        if slide_idx >= len(prs.slides):
            print(f"页面 {page_num} 不存在")
            continue
        
        slide = prs.slides[slide_idx]
        
        print(f"\n{'='*80}")
        print(f"分析 PPTX 第 {page_num} 页")
        print(f"{'='*80}\n")
        
        # 查找表格
        tables = []
        for shape in slide.shapes:
            if shape.has_table:
                tables.append(shape.table)
        
        print(f"找到 {len(tables)} 个表格\n")
        
        for i, table in enumerate(tables):
            print(f"表格 {i+1}:")
            print(f"  行数: {len(table.rows)}")
            print(f"  列数: {len(table.columns)}")
            print(f"  位置: ({shape.left}, {shape.top})")
            print(f"  尺寸: {shape.width} x {shape.height}\n")
            
            # 打印前几行内容
            print(f"  前5行数据：")
            for row_idx in range(min(5, len(table.rows))):
                print(f"    第{row_idx+1}行: ", end='')
                for col_idx in range(min(3, len(table.columns))):
                    cell = table.rows[row_idx].cells[col_idx]
                    text = cell.text.strip()[:30]
                    # 检查合并单元格
                    is_merged = ""
                    try:
                        if cell.is_merge_origin:
                            is_merged = " [合并源]"
                        elif cell.is_spanned:
                            is_merged = " [被合并]"
                    except:
                        pass
                    print(f"[{text}{is_merged}]", end=' ')
                print()
            
            # 检查列宽
            print(f"\n  列宽 (英寸):")
            for col_idx, col in enumerate(table.columns):
                width_inches = col.width / 914400  # EMU to inches
                print(f"    列{col_idx+1}: {width_inches:.2f}\"", end='')
                if col_idx < len(table.columns) - 1:
                    print(", ", end='')
            print("\n")
            
            # 检查背景颜色
            print(f"  单元格颜色样本（前3行前3列）:")
            for row_idx in range(min(3, len(table.rows))):
                for col_idx in range(min(3, len(table.columns))):
                    cell = table.rows[row_idx].cells[col_idx]
                    try:
                        fill = cell.fill
                        if fill.type:
                            print(f"    ({row_idx+1},{col_idx+1}): {fill.type}", end='')
                            if hasattr(fill, 'fore_color'):
                                try:
                                    rgb = fill.fore_color.rgb
                                    print(f" RGB({rgb[0]},{rgb[1]},{rgb[2]})")
                                except:
                                    print()
                            else:
                                print()
                    except:
                        pass

if __name__ == "__main__":
    import sys
    pptx_path = sys.argv[1] if len(sys.argv) > 1 else "output/test_page9.pptx"
    
    # 分析第8、9、12页（PPTX中的索引是8、9、12，因为是0-based）
    analyze_pptx_tables(pptx_path, [8, 9, 12])
